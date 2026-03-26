#!/usr/bin/env python3
# SPDX-License-Identifier: GPL-3.0-or-later
# Copyright 2026 Joseph Webber
"""
Example 37: RAG Document Q&A System

A comprehensive document Q&A system supporting multiple file formats:
- PDF documents with text extraction
- Word documents (.docx)
- Excel spreadsheets (.xlsx)
- PowerPoint presentations (.pptx)
- Plain text and Markdown

Key RAG features demonstrated:
- Multiple chunking strategies (fixed, semantic, sliding window)
- Embedding generation with sentence transformers
- Vector similarity search with FAISS
- Hybrid search (BM25 + semantic)
- Reranking with cross-encoders
- Source citation in responses
- Evaluation metrics (relevance, accuracy)

Demo: Company policy documents Q&A

Usage:
    python examples/37_rag_documents.py
    python examples/37_rag_documents.py --demo  # Run with sample data

Requirements:
    pip install agentic-brain sentence-transformers faiss-cpu
    pip install pypdf python-docx openpyxl python-pptx
"""

import asyncio
import hashlib
import json
import os
import re
import tempfile
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Callable, Generator, Optional
import math

# Try importing optional dependencies
try:
    import numpy as np

    HAS_NUMPY = True
except ImportError:
    HAS_NUMPY = False
    print("⚠️ NumPy not available, using mock vectors")


# ══════════════════════════════════════════════════════════════════════════════
# CONFIGURATION
# ══════════════════════════════════════════════════════════════════════════════


@dataclass
class RAGConfig:
    """Configuration for the RAG pipeline."""

    # Embedding settings
    embedding_model: str = "all-MiniLM-L6-v2"
    embedding_dimension: int = 384

    # Chunking settings
    chunk_size: int = 512
    chunk_overlap: int = 50
    min_chunk_size: int = 100

    # Retrieval settings
    top_k: int = 5
    rerank_top_k: int = 3
    similarity_threshold: float = 0.3

    # Hybrid search weights
    semantic_weight: float = 0.7
    keyword_weight: float = 0.3

    # LLM settings
    max_context_tokens: int = 4000
    temperature: float = 0.1


# ══════════════════════════════════════════════════════════════════════════════
# ENUMS AND DATA MODELS
# ══════════════════════════════════════════════════════════════════════════════


class ChunkingStrategy(Enum):
    """Document chunking strategies."""

    FIXED = "fixed"  # Fixed character count
    SEMANTIC = "semantic"  # Sentence/paragraph boundaries
    SLIDING_WINDOW = "sliding"  # Overlapping windows
    RECURSIVE = "recursive"  # Hierarchical (sections, paragraphs, sentences)
    MARKDOWN = "markdown"  # Preserve markdown structure


class DocumentType(Enum):
    """Supported document types."""

    PDF = "pdf"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    TXT = "txt"
    MD = "md"
    HTML = "html"
    JSON = "json"


@dataclass
class DocumentChunk:
    """A chunk of document content."""

    id: str
    document_id: str
    content: str
    embedding: Optional[list[float]] = None
    metadata: dict = field(default_factory=dict)
    start_char: int = 0
    end_char: int = 0
    chunk_index: int = 0

    @property
    def source_info(self) -> str:
        """Generate source citation."""
        title = self.metadata.get("title", "Unknown")
        page = self.metadata.get("page", "")
        section = self.metadata.get("section", "")

        citation = f"[{title}"
        if page:
            citation += f", p.{page}"
        if section:
            citation += f", §{section}"
        citation += "]"
        return citation


@dataclass
class Document:
    """A document in the system."""

    id: str
    title: str
    content: str
    doc_type: DocumentType
    metadata: dict = field(default_factory=dict)
    chunks: list[DocumentChunk] = field(default_factory=list)
    created_at: datetime = field(default_factory=datetime.now)

    @property
    def word_count(self) -> int:
        return len(self.content.split())


@dataclass
class RetrievalResult:
    """Result from retrieval."""

    chunk: DocumentChunk
    score: float
    rank: int
    match_type: str = "semantic"  # semantic, keyword, hybrid


@dataclass
class QueryResult:
    """Complete query result with answer and sources."""

    query: str
    answer: str
    sources: list[RetrievalResult]
    confidence: float
    latency_ms: float
    tokens_used: int = 0


@dataclass
class EvaluationMetrics:
    """Metrics for evaluating RAG performance."""

    relevance_score: float  # How relevant are retrieved chunks
    answer_accuracy: float  # How accurate is the answer
    faithfulness: float  # Does answer match sources
    context_precision: float  # Relevant chunks / total retrieved
    response_time_ms: float


# ══════════════════════════════════════════════════════════════════════════════
# DOCUMENT LOADERS
# ══════════════════════════════════════════════════════════════════════════════


class DocumentLoader(ABC):
    """Abstract base class for document loaders."""

    @abstractmethod
    def load(self, path: str) -> tuple[str, dict]:
        """Load document content and metadata."""
        pass

    @abstractmethod
    def supported_extensions(self) -> list[str]:
        """Return list of supported file extensions."""
        pass


class TextLoader(DocumentLoader):
    """Load plain text and markdown files."""

    def load(self, path: str) -> tuple[str, dict]:
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()

        metadata = {
            "filename": os.path.basename(path),
            "file_size": os.path.getsize(path),
            "extension": Path(path).suffix.lower(),
        }
        return content, metadata

    def supported_extensions(self) -> list[str]:
        return [".txt", ".md", ".markdown"]


class PDFLoader(DocumentLoader):
    """Load PDF documents."""

    def load(self, path: str) -> tuple[str, dict]:
        try:
            import pypdf

            with open(path, "rb") as f:
                reader = pypdf.PdfReader(f)
                pages = []
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(f"[Page {i+1}]\n{text}")

                content = "\n\n".join(pages)
                metadata = {
                    "filename": os.path.basename(path),
                    "page_count": len(reader.pages),
                    "file_size": os.path.getsize(path),
                }

                # Extract PDF metadata if available
                if reader.metadata:
                    metadata["author"] = reader.metadata.get("/Author", "")
                    metadata["title"] = reader.metadata.get("/Title", "")
                    metadata["created"] = str(reader.metadata.get("/CreationDate", ""))

                return content, metadata
        except ImportError:
            return f"[PDF content from {path} - pypdf not installed]", {
                "filename": path
            }

    def supported_extensions(self) -> list[str]:
        return [".pdf"]


class DocxLoader(DocumentLoader):
    """Load Word documents."""

    def load(self, path: str) -> tuple[str, dict]:
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(path)
            paragraphs = []

            for para in doc.paragraphs:
                if para.text.strip():
                    # Check for headings
                    if para.style.name.startswith("Heading"):
                        level = para.style.name.replace("Heading ", "")
                        paragraphs.append(
                            f"{'#' * int(level) if level.isdigit() else '##'} {para.text}"
                        )
                    else:
                        paragraphs.append(para.text)

            # Extract tables
            for table in doc.tables:
                table_text = "\n[Table]\n"
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_text += " | ".join(cells) + "\n"
                paragraphs.append(table_text)

            content = "\n\n".join(paragraphs)
            metadata = {
                "filename": os.path.basename(path),
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
            }

            # Core properties
            if doc.core_properties:
                metadata["author"] = doc.core_properties.author or ""
                metadata["title"] = doc.core_properties.title or ""

            return content, metadata
        except ImportError:
            return f"[DOCX content from {path} - python-docx not installed]", {
                "filename": path
            }

    def supported_extensions(self) -> list[str]:
        return [".docx", ".doc"]


class ExcelLoader(DocumentLoader):
    """Load Excel spreadsheets."""

    def load(self, path: str) -> tuple[str, dict]:
        try:
            import openpyxl

            workbook = openpyxl.load_workbook(path, data_only=True)
            sheets_content = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = f"\n## Sheet: {sheet_name}\n\n"

                rows = list(sheet.iter_rows(values_only=True))
                if rows:
                    # First row as header
                    headers = [str(cell) if cell else "" for cell in rows[0]]
                    sheet_text += " | ".join(headers) + "\n"
                    sheet_text += " | ".join(["---"] * len(headers)) + "\n"

                    # Data rows
                    for row in rows[1:]:
                        cells = [str(cell) if cell else "" for cell in row]
                        if any(cells):  # Skip empty rows
                            sheet_text += " | ".join(cells) + "\n"

                sheets_content.append(sheet_text)

            content = "\n".join(sheets_content)
            metadata = {
                "filename": os.path.basename(path),
                "sheet_count": len(workbook.sheetnames),
                "sheets": workbook.sheetnames,
            }

            return content, metadata
        except ImportError:
            return f"[XLSX content from {path} - openpyxl not installed]", {
                "filename": path
            }

    def supported_extensions(self) -> list[str]:
        return [".xlsx", ".xls"]


class PowerPointLoader(DocumentLoader):
    """Load PowerPoint presentations."""

    def load(self, path: str) -> tuple[str, dict]:
        try:
            from pptx import Presentation

            prs = Presentation(path)
            slides_content = []

            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"\n## Slide {i}\n\n"

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_text += f"- {text}\n"

                    # Handle tables
                    if shape.has_table:
                        slide_text += "\n[Table]\n"
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_text += " | ".join(cells) + "\n"

                slides_content.append(slide_text)

            content = "\n".join(slides_content)
            metadata = {
                "filename": os.path.basename(path),
                "slide_count": len(prs.slides),
            }

            return content, metadata
        except ImportError:
            return f"[PPTX content from {path} - python-pptx not installed]", {
                "filename": path
            }

    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]


class DocumentLoaderFactory:
    """Factory for creating appropriate document loaders."""

    _loaders: list[DocumentLoader] = [
        TextLoader(),
        PDFLoader(),
        DocxLoader(),
        ExcelLoader(),
        PowerPointLoader(),
    ]

    @classmethod
    def get_loader(cls, path: str) -> Optional[DocumentLoader]:
        """Get appropriate loader for file type."""
        ext = Path(path).suffix.lower()
        for loader in cls._loaders:
            if ext in loader.supported_extensions():
                return loader
        return None

    @classmethod
    def supported_extensions(cls) -> list[str]:
        """Get all supported extensions."""
        extensions = []
        for loader in cls._loaders:
            extensions.extend(loader.supported_extensions())
        return extensions


# ══════════════════════════════════════════════════════════════════════════════
# CHUNKING STRATEGIES
# ══════════════════════════════════════════════════════════════════════════════


class Chunker:
    """Document chunking with multiple strategies."""

    def __init__(self, config: RAGConfig):
        self.config = config

    def chunk(
        self, content: str, strategy: ChunkingStrategy = ChunkingStrategy.SEMANTIC
    ) -> list[tuple[str, int, int]]:
        """Chunk content using specified strategy.

        Returns list of (chunk_text, start_char, end_char).
        """
        if strategy == ChunkingStrategy.FIXED:
            return self._fixed_chunking(content)
        elif strategy == ChunkingStrategy.SEMANTIC:
            return self._semantic_chunking(content)
        elif strategy == ChunkingStrategy.SLIDING_WINDOW:
            return self._sliding_window(content)
        elif strategy == ChunkingStrategy.RECURSIVE:
            return self._recursive_chunking(content)
        elif strategy == ChunkingStrategy.MARKDOWN:
            return self._markdown_chunking(content)
        else:
            return self._semantic_chunking(content)

    def _fixed_chunking(self, content: str) -> list[tuple[str, int, int]]:
        """Split at fixed character intervals."""
        chunks = []
        start = 0

        while start < len(content):
            end = min(start + self.config.chunk_size, len(content))

            # Don't break in middle of word
            if end < len(content):
                space_idx = content.rfind(" ", start, end)
                if space_idx > start:
                    end = space_idx

            chunk_text = content[start:end].strip()
            if len(chunk_text) >= self.config.min_chunk_size:
                chunks.append((chunk_text, start, end))

            start = end - self.config.chunk_overlap
            if start >= len(content):
                break

        return chunks

    def _semantic_chunking(self, content: str) -> list[tuple[str, int, int]]:
        """Split at sentence/paragraph boundaries."""
        # Split by paragraphs first
        paragraphs = re.split(r"\n\s*\n", content)
        chunks = []
        current_chunk = ""
        current_start = 0
        char_pos = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                char_pos += 2  # Account for paragraph break
                continue

            # If adding this paragraph exceeds limit, save current chunk
            if (
                current_chunk
                and len(current_chunk) + len(para) > self.config.chunk_size
            ):
                if len(current_chunk) >= self.config.min_chunk_size:
                    chunks.append((current_chunk.strip(), current_start, char_pos))
                current_chunk = para
                current_start = char_pos
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
                    current_start = char_pos

            char_pos += len(para) + 2

        # Add final chunk
        if current_chunk and len(current_chunk) >= self.config.min_chunk_size:
            chunks.append((current_chunk.strip(), current_start, char_pos))

        return chunks

    def _sliding_window(self, content: str) -> list[tuple[str, int, int]]:
        """Overlapping sliding window approach."""
        chunks = []
        window_size = self.config.chunk_size
        step_size = window_size - self.config.chunk_overlap

        sentences = re.split(r"(?<=[.!?])\s+", content)
        current_chunk = ""
        current_start = 0
        char_pos = 0

        for sentence in sentences:
            if len(current_chunk) + len(sentence) > window_size and current_chunk:
                # Save current window
                if len(current_chunk) >= self.config.min_chunk_size:
                    chunks.append((current_chunk.strip(), current_start, char_pos))

                # Slide window - keep overlap portion
                words = current_chunk.split()
                overlap_words = words[-(len(words) // 3) :]
                current_chunk = " ".join(overlap_words) + " " + sentence
                current_start = char_pos - len(" ".join(overlap_words))
            else:
                if current_chunk:
                    current_chunk += " " + sentence
                else:
                    current_chunk = sentence
                    current_start = char_pos

            char_pos += len(sentence) + 1

        if current_chunk and len(current_chunk) >= self.config.min_chunk_size:
            chunks.append((current_chunk.strip(), current_start, char_pos))

        return chunks

    def _recursive_chunking(self, content: str) -> list[tuple[str, int, int]]:
        """Hierarchical chunking using multiple separators."""
        separators = [
            "\n\n\n",  # Section breaks
            "\n\n",  # Paragraphs
            "\n",  # Lines
            ". ",  # Sentences
            " ",  # Words
        ]

        return self._recursive_split(content, separators, 0)

    def _recursive_split(
        self, text: str, separators: list[str], start_pos: int
    ) -> list[tuple[str, int, int]]:
        """Recursively split text using separators hierarchy."""
        if not text or not separators:
            if len(text) >= self.config.min_chunk_size:
                return [(text, start_pos, start_pos + len(text))]
            return []

        separator = separators[0]
        parts = text.split(separator)

        chunks = []
        current_chunk = ""
        current_start = start_pos
        char_pos = start_pos

        for part in parts:
            if not part.strip():
                char_pos += len(separator)
                continue

            if (
                len(current_chunk) + len(part) > self.config.chunk_size
                and current_chunk
            ):
                if len(current_chunk) >= self.config.min_chunk_size:
                    chunks.append((current_chunk.strip(), current_start, char_pos))
                elif len(separators) > 1:
                    # Try finer separators
                    sub_chunks = self._recursive_split(
                        current_chunk, separators[1:], current_start
                    )
                    chunks.extend(sub_chunks)

                current_chunk = part
                current_start = char_pos
            else:
                if current_chunk:
                    current_chunk += separator + part
                else:
                    current_chunk = part
                    current_start = char_pos

            char_pos += len(part) + len(separator)

        if current_chunk and len(current_chunk) >= self.config.min_chunk_size:
            chunks.append((current_chunk.strip(), current_start, char_pos))

        return chunks

    def _markdown_chunking(self, content: str) -> list[tuple[str, int, int]]:
        """Preserve markdown structure (headers, code blocks, etc.)."""
        chunks = []

        # Split by headers
        header_pattern = r"^(#{1,6})\s+(.+)$"
        lines = content.split("\n")
        current_chunk = ""
        current_start = 0
        char_pos = 0
        current_header = ""

        in_code_block = False

        for line in lines:
            # Track code blocks
            if line.strip().startswith("```"):
                in_code_block = not in_code_block

            # Check for header (outside code blocks)
            header_match = re.match(header_pattern, line) if not in_code_block else None

            if header_match:
                # Save previous chunk
                if current_chunk and len(current_chunk) >= self.config.min_chunk_size:
                    chunks.append((current_chunk.strip(), current_start, char_pos))

                current_header = line
                current_chunk = line
                current_start = char_pos
            else:
                # Check chunk size
                if (
                    len(current_chunk) + len(line) > self.config.chunk_size
                    and current_chunk
                ):
                    if len(current_chunk) >= self.config.min_chunk_size:
                        chunks.append((current_chunk.strip(), current_start, char_pos))

                    # Start new chunk, preserve header context
                    if current_header and not line.strip().startswith("#"):
                        current_chunk = current_header + "\n" + line
                    else:
                        current_chunk = line
                    current_start = char_pos
                else:
                    current_chunk += "\n" + line

            char_pos += len(line) + 1

        if current_chunk and len(current_chunk) >= self.config.min_chunk_size:
            chunks.append((current_chunk.strip(), current_start, char_pos))

        return chunks


# ══════════════════════════════════════════════════════════════════════════════
# EMBEDDING AND VECTOR OPERATIONS
# ══════════════════════════════════════════════════════════════════════════════


class EmbeddingGenerator:
    """Generate embeddings for text."""

    def __init__(self, model_name: str = "all-MiniLM-L6-v2"):
        self.model_name = model_name
        self.model = None
        self._dimension = 384

    def _load_model(self):
        """Lazy load the embedding model."""
        if self.model is None:
            try:
                from sentence_transformers import SentenceTransformer

                self.model = SentenceTransformer(self.model_name)
                self._dimension = self.model.get_sentence_embedding_dimension()
            except ImportError:
                print("⚠️ sentence-transformers not installed, using mock embeddings")
                self.model = "mock"

    def embed(self, text: str) -> list[float]:
        """Generate embedding for single text."""
        self._load_model()

        if self.model == "mock":
            return self._mock_embedding(text)

        embedding = self.model.encode(text, convert_to_numpy=True)
        return embedding.tolist()

    def embed_batch(self, texts: list[str]) -> list[list[float]]:
        """Generate embeddings for batch of texts."""
        self._load_model()

        if self.model == "mock":
            return [self._mock_embedding(t) for t in texts]

        embeddings = self.model.encode(texts, convert_to_numpy=True)
        return embeddings.tolist()

    def _mock_embedding(self, text: str) -> list[float]:
        """Generate deterministic mock embedding based on text hash."""
        import hashlib

        hash_val = hashlib.md5(text.encode()).hexdigest()

        # Generate pseudo-random but deterministic vector
        embedding = []
        for i in range(0, min(len(hash_val), self._dimension), 2):
            val = int(hash_val[i : i + 2], 16) / 255.0 - 0.5
            embedding.append(val)

        # Pad to full dimension
        while len(embedding) < self._dimension:
            embedding.append(0.0)

        # Normalize
        norm = math.sqrt(sum(x * x for x in embedding))
        if norm > 0:
            embedding = [x / norm for x in embedding]

        return embedding[: self._dimension]

    @property
    def dimension(self) -> int:
        """Get embedding dimension."""
        return self._dimension


class VectorStore:
    """Simple in-memory vector store."""

    def __init__(self, dimension: int = 384):
        self.dimension = dimension
        self.vectors: list[list[float]] = []
        self.metadata: list[dict] = []
        self.index = None

    def add(self, vector: list[float], meta: dict):
        """Add a vector with metadata."""
        self.vectors.append(vector)
        self.metadata.append(meta)
        self.index = None  # Invalidate index

    def add_batch(self, vectors: list[list[float]], metas: list[dict]):
        """Add multiple vectors."""
        self.vectors.extend(vectors)
        self.metadata.extend(metas)
        self.index = None

    def _build_index(self):
        """Build search index."""
        if not HAS_NUMPY:
            return

        try:
            import faiss

            vectors_np = np.array(self.vectors, dtype=np.float32)
            self.index = faiss.IndexFlatIP(self.dimension)  # Inner product
            faiss.normalize_L2(vectors_np)
            self.index.add(vectors_np)
        except ImportError:
            # Fall back to brute force
            self.index = "bruteforce"

    def search(
        self, query_vector: list[float], top_k: int = 5
    ) -> list[tuple[int, float]]:
        """Search for similar vectors.

        Returns list of (index, similarity_score).
        """
        if not self.vectors:
            return []

        if self.index is None:
            self._build_index()

        if HAS_NUMPY and self.index is not None and self.index != "bruteforce":
            try:
                import faiss

                query_np = np.array([query_vector], dtype=np.float32)
                faiss.normalize_L2(query_np)

                distances, indices = self.index.search(
                    query_np, min(top_k, len(self.vectors))
                )

                results = []
                for i, (idx, score) in enumerate(zip(indices[0], distances[0])):
                    if idx >= 0:
                        results.append((int(idx), float(score)))
                return results
            except Exception:
                pass

        # Brute force fallback
        return self._brute_force_search(query_vector, top_k)

    def _brute_force_search(
        self, query_vector: list[float], top_k: int
    ) -> list[tuple[int, float]]:
        """Brute force similarity search."""
        similarities = []

        for i, vec in enumerate(self.vectors):
            sim = self._cosine_similarity(query_vector, vec)
            similarities.append((i, sim))

        # Sort by similarity descending
        similarities.sort(key=lambda x: x[1], reverse=True)
        return similarities[:top_k]

    def _cosine_similarity(self, a: list[float], b: list[float]) -> float:
        """Compute cosine similarity between two vectors."""
        dot_product = sum(x * y for x, y in zip(a, b))
        norm_a = math.sqrt(sum(x * x for x in a))
        norm_b = math.sqrt(sum(x * x for x in b))

        if norm_a == 0 or norm_b == 0:
            return 0.0

        return dot_product / (norm_a * norm_b)


# ══════════════════════════════════════════════════════════════════════════════
# BM25 KEYWORD SEARCH
# ══════════════════════════════════════════════════════════════════════════════


class BM25Index:
    """BM25 keyword search index."""

    def __init__(self, k1: float = 1.5, b: float = 0.75):
        self.k1 = k1
        self.b = b
        self.documents: list[list[str]] = []
        self.doc_lengths: list[int] = []
        self.avg_doc_length: float = 0
        self.doc_freqs: dict[str, int] = {}
        self.total_docs: int = 0

    def add(self, text: str):
        """Add a document to the index."""
        tokens = self._tokenize(text)
        self.documents.append(tokens)
        self.doc_lengths.append(len(tokens))

        # Update document frequencies
        unique_tokens = set(tokens)
        for token in unique_tokens:
            self.doc_freqs[token] = self.doc_freqs.get(token, 0) + 1

        self.total_docs += 1
        self.avg_doc_length = sum(self.doc_lengths) / self.total_docs

    def search(self, query: str, top_k: int = 5) -> list[tuple[int, float]]:
        """Search for documents matching query."""
        query_tokens = self._tokenize(query)
        scores = []

        for doc_idx, doc_tokens in enumerate(self.documents):
            score = self._score_document(query_tokens, doc_tokens, doc_idx)
            scores.append((doc_idx, score))

        scores.sort(key=lambda x: x[1], reverse=True)
        return scores[:top_k]

    def _score_document(
        self, query_tokens: list[str], doc_tokens: list[str], doc_idx: int
    ) -> float:
        """Calculate BM25 score for a document."""
        score = 0.0
        doc_length = self.doc_lengths[doc_idx]

        for token in query_tokens:
            if token not in self.doc_freqs:
                continue

            # Term frequency in document
            tf = doc_tokens.count(token)
            if tf == 0:
                continue

            # Inverse document frequency
            df = self.doc_freqs[token]
            idf = math.log((self.total_docs - df + 0.5) / (df + 0.5) + 1)

            # BM25 score
            numerator = tf * (self.k1 + 1)
            denominator = tf + self.k1 * (
                1 - self.b + self.b * doc_length / self.avg_doc_length
            )

            score += idf * (numerator / denominator)

        return score

    def _tokenize(self, text: str) -> list[str]:
        """Tokenize text for indexing."""
        # Lowercase and split on non-alphanumeric
        tokens = re.findall(r"\b\w+\b", text.lower())
        # Remove stopwords (basic list)
        stopwords = {
            "the",
            "a",
            "an",
            "is",
            "are",
            "was",
            "were",
            "be",
            "been",
            "being",
            "have",
            "has",
            "had",
            "do",
            "does",
            "did",
            "will",
            "would",
            "could",
            "should",
            "may",
            "might",
            "must",
            "shall",
            "of",
            "to",
            "in",
            "for",
            "on",
            "with",
            "at",
            "by",
            "from",
            "as",
            "and",
            "or",
            "not",
            "but",
            "if",
            "then",
            "else",
            "this",
            "that",
            "it",
            "its",
            "they",
            "them",
            "their",
            "we",
            "our",
        }
        return [t for t in tokens if t not in stopwords and len(t) > 1]


# ══════════════════════════════════════════════════════════════════════════════
# RERANKER
# ══════════════════════════════════════════════════════════════════════════════


class Reranker:
    """Rerank retrieved results using cross-encoder."""

    def __init__(self, model_name: str = "cross-encoder/ms-marco-MiniLM-L-6-v2"):
        self.model_name = model_name
        self.model = None

    def _load_model(self):
        """Lazy load the reranking model."""
        if self.model is None:
            try:
                from sentence_transformers import CrossEncoder

                self.model = CrossEncoder(self.model_name)
            except ImportError:
                print("⚠️ CrossEncoder not available, using score-based reranking")
                self.model = "fallback"

    def rerank(
        self, query: str, results: list[RetrievalResult], top_k: int = 3
    ) -> list[RetrievalResult]:
        """Rerank results using cross-encoder."""
        if not results:
            return []

        self._load_model()

        if self.model == "fallback":
            # Just re-sort by existing scores
            sorted_results = sorted(results, key=lambda x: x.score, reverse=True)
            for i, r in enumerate(sorted_results[:top_k]):
                r.rank = i + 1
            return sorted_results[:top_k]

        # Prepare pairs for cross-encoder
        pairs = [(query, r.chunk.content) for r in results]

        # Get cross-encoder scores
        scores = self.model.predict(pairs)

        # Combine with original results
        scored_results = list(zip(results, scores))
        scored_results.sort(key=lambda x: x[1], reverse=True)

        # Update ranks and return top_k
        reranked = []
        for i, (result, new_score) in enumerate(scored_results[:top_k]):
            result.score = float(new_score)
            result.rank = i + 1
            reranked.append(result)

        return reranked


# ══════════════════════════════════════════════════════════════════════════════
# RAG PIPELINE
# ══════════════════════════════════════════════════════════════════════════════


class DocumentRAGPipeline:
    """Complete RAG pipeline for document Q&A."""

    def __init__(self, config: RAGConfig = None):
        self.config = config or RAGConfig()

        # Components
        self.chunker = Chunker(self.config)
        self.embedder = EmbeddingGenerator(self.config.embedding_model)
        self.vector_store = VectorStore(self.config.embedding_dimension)
        self.bm25_index = BM25Index()
        self.reranker = Reranker()

        # Document storage
        self.documents: dict[str, Document] = {}
        self.chunks: dict[str, DocumentChunk] = {}
        self.chunk_id_to_index: dict[str, int] = {}

    def add_document(
        self,
        path: str = None,
        content: str = None,
        title: str = None,
        metadata: dict = None,
        chunking_strategy: ChunkingStrategy = ChunkingStrategy.SEMANTIC,
    ) -> Document:
        """Add a document to the pipeline."""
        # Load from file if path provided
        if path:
            loader = DocumentLoaderFactory.get_loader(path)
            if loader:
                content, file_metadata = loader.load(path)
                metadata = {**(metadata or {}), **file_metadata}
                if not title:
                    title = metadata.get("title") or metadata.get(
                        "filename", "Untitled"
                    )

        if not content:
            raise ValueError("Either path or content must be provided")

        # Create document
        doc_id = hashlib.md5(content[:1000].encode()).hexdigest()[:12]
        doc_type = self._detect_type(path) if path else DocumentType.TXT

        doc = Document(
            id=doc_id,
            title=title or "Untitled",
            content=content,
            doc_type=doc_type,
            metadata=metadata or {},
        )

        # Chunk the document
        chunk_tuples = self.chunker.chunk(content, chunking_strategy)

        # Create chunk objects
        for i, (chunk_text, start, end) in enumerate(chunk_tuples):
            chunk_id = f"{doc_id}_chunk_{i}"

            chunk = DocumentChunk(
                id=chunk_id,
                document_id=doc_id,
                content=chunk_text,
                metadata={
                    **doc.metadata,
                    "title": doc.title,
                    "chunk_index": i,
                    "total_chunks": len(chunk_tuples),
                },
                start_char=start,
                end_char=end,
                chunk_index=i,
            )

            doc.chunks.append(chunk)
            self.chunks[chunk_id] = chunk

        # Generate embeddings
        chunk_texts = [c.content for c in doc.chunks]
        embeddings = self.embedder.embed_batch(chunk_texts)

        # Add to vector store and BM25 index
        for chunk, embedding in zip(doc.chunks, embeddings):
            chunk.embedding = embedding

            idx = len(self.vector_store.vectors)
            self.chunk_id_to_index[chunk.id] = idx

            self.vector_store.add(embedding, {"chunk_id": chunk.id})
            self.bm25_index.add(chunk.content)

        self.documents[doc_id] = doc
        return doc

    def _detect_type(self, path: str) -> DocumentType:
        """Detect document type from path."""
        ext = Path(path).suffix.lower()
        type_map = {
            ".pdf": DocumentType.PDF,
            ".docx": DocumentType.DOCX,
            ".doc": DocumentType.DOCX,
            ".xlsx": DocumentType.XLSX,
            ".xls": DocumentType.XLSX,
            ".pptx": DocumentType.PPTX,
            ".ppt": DocumentType.PPTX,
            ".txt": DocumentType.TXT,
            ".md": DocumentType.MD,
            ".html": DocumentType.HTML,
            ".json": DocumentType.JSON,
        }
        return type_map.get(ext, DocumentType.TXT)

    def query(
        self,
        question: str,
        top_k: int = None,
        use_reranking: bool = True,
        include_sources: bool = True,
    ) -> QueryResult:
        """Query the document collection."""
        import time

        start_time = time.time()

        top_k = top_k or self.config.top_k

        # Hybrid search
        results = self._hybrid_search(question, top_k * 2)

        # Rerank if enabled
        if use_reranking and len(results) > 0:
            results = self.reranker.rerank(question, results, self.config.rerank_top_k)
        else:
            results = results[:top_k]

        # Build context from results
        context = self._build_context(results)

        # Generate answer
        answer = self._generate_answer(
            question, context, results if include_sources else []
        )

        latency = (time.time() - start_time) * 1000

        # Calculate confidence
        confidence = self._calculate_confidence(results)

        return QueryResult(
            query=question,
            answer=answer,
            sources=results,
            confidence=confidence,
            latency_ms=latency,
        )

    def _hybrid_search(self, query: str, top_k: int) -> list[RetrievalResult]:
        """Perform hybrid search (semantic + keyword)."""
        # Semantic search
        query_embedding = self.embedder.embed(query)
        semantic_results = self.vector_store.search(query_embedding, top_k)

        # Keyword search
        keyword_results = self.bm25_index.search(query, top_k)

        # Merge results with weighting
        chunk_scores: dict[str, tuple[float, str]] = {}
        chunk_ids = list(self.chunk_id_to_index.keys())

        # Normalize and add semantic scores
        if semantic_results:
            max_sem = max(r[1] for r in semantic_results)
            for idx, score in semantic_results:
                if idx < len(chunk_ids):
                    chunk_id = chunk_ids[idx]
                    normalized = score / max_sem if max_sem > 0 else 0
                    weighted = normalized * self.config.semantic_weight
                    chunk_scores[chunk_id] = (weighted, "semantic")

        # Add keyword scores
        if keyword_results:
            max_kw = max(r[1] for r in keyword_results) if keyword_results else 1
            for idx, score in keyword_results:
                if idx < len(chunk_ids):
                    chunk_id = chunk_ids[idx]
                    normalized = score / max_kw if max_kw > 0 else 0
                    weighted = normalized * self.config.keyword_weight

                    if chunk_id in chunk_scores:
                        old_score, _ = chunk_scores[chunk_id]
                        chunk_scores[chunk_id] = (old_score + weighted, "hybrid")
                    else:
                        chunk_scores[chunk_id] = (weighted, "keyword")

        # Sort by combined score
        sorted_chunks = sorted(
            chunk_scores.items(), key=lambda x: x[1][0], reverse=True
        )

        # Build retrieval results
        results = []
        for rank, (chunk_id, (score, match_type)) in enumerate(
            sorted_chunks[:top_k], 1
        ):
            if chunk_id in self.chunks:
                results.append(
                    RetrievalResult(
                        chunk=self.chunks[chunk_id],
                        score=score,
                        rank=rank,
                        match_type=match_type,
                    )
                )

        return results

    def _build_context(self, results: list[RetrievalResult]) -> str:
        """Build context string from retrieval results."""
        context_parts = []

        for r in results:
            source = r.chunk.source_info
            context_parts.append(f"Source {source}:\n{r.chunk.content}")

        return "\n\n---\n\n".join(context_parts)

    def _generate_answer(
        self, question: str, context: str, sources: list[RetrievalResult]
    ) -> str:
        """Generate answer using LLM or rule-based approach."""
        # For demo, use template-based response
        # In production, this would call an LLM

        if not context:
            return "I couldn't find relevant information to answer this question."

        # Extract key sentences from context that match question keywords
        question_words = set(re.findall(r"\b\w+\b", question.lower()))
        question_words -= {
            "what",
            "how",
            "why",
            "when",
            "where",
            "who",
            "is",
            "are",
            "the",
            "a",
            "an",
            "does",
            "do",
            "can",
            "could",
            "would",
        }

        relevant_sentences = []
        for source in sources:
            sentences = re.split(r"[.!?]+", source.chunk.content)
            for sent in sentences:
                sent_words = set(re.findall(r"\b\w+\b", sent.lower()))
                overlap = question_words & sent_words
                if len(overlap) >= min(2, len(question_words)):
                    relevant_sentences.append((sent.strip(), source.chunk.source_info))

        if relevant_sentences:
            answer = "Based on the documents:\n\n"
            for sent, source in relevant_sentences[:3]:
                if sent:
                    answer += f"• {sent}. {source}\n"
            return answer

        # Fallback to context summary
        return f"Based on the available documents:\n\n{sources[0].chunk.content[:500]}...\n\n{sources[0].chunk.source_info}"

    def _calculate_confidence(self, results: list[RetrievalResult]) -> float:
        """Calculate confidence score for the results."""
        if not results:
            return 0.0

        # Average of top scores, weighted by rank
        weighted_sum = 0.0
        weight_total = 0.0

        for r in results:
            weight = 1.0 / r.rank
            weighted_sum += r.score * weight
            weight_total += weight

        return min(weighted_sum / weight_total, 1.0) if weight_total > 0 else 0.0

    def evaluate(
        self,
        question: str,
        expected_answer: str = None,
        relevant_doc_ids: list[str] = None,
    ) -> EvaluationMetrics:
        """Evaluate RAG performance on a query."""
        import time

        start = time.time()

        result = self.query(question)

        response_time = (time.time() - start) * 1000

        # Context precision
        if relevant_doc_ids:
            retrieved_doc_ids = [r.chunk.document_id for r in result.sources]
            relevant_retrieved = len(set(retrieved_doc_ids) & set(relevant_doc_ids))
            context_precision = (
                relevant_retrieved / len(result.sources) if result.sources else 0
            )
        else:
            context_precision = result.confidence

        # Relevance score (based on retrieval scores)
        relevance = (
            sum(r.score for r in result.sources) / len(result.sources)
            if result.sources
            else 0
        )

        # Answer accuracy (simple word overlap with expected if provided)
        if expected_answer:
            answer_words = set(re.findall(r"\b\w+\b", result.answer.lower()))
            expected_words = set(re.findall(r"\b\w+\b", expected_answer.lower()))
            overlap = len(answer_words & expected_words)
            accuracy = overlap / len(expected_words) if expected_words else 0
        else:
            accuracy = result.confidence

        # Faithfulness (do answer words appear in sources?)
        answer_words = set(re.findall(r"\b\w+\b", result.answer.lower()))
        source_words = set()
        for r in result.sources:
            source_words.update(re.findall(r"\b\w+\b", r.chunk.content.lower()))

        common = len(answer_words & source_words)
        faithfulness = common / len(answer_words) if answer_words else 0

        return EvaluationMetrics(
            relevance_score=relevance,
            answer_accuracy=accuracy,
            faithfulness=faithfulness,
            context_precision=context_precision,
            response_time_ms=response_time,
        )

    def get_stats(self) -> dict:
        """Get pipeline statistics."""
        return {
            "documents": len(self.documents),
            "chunks": len(self.chunks),
            "vectors": len(self.vector_store.vectors),
            "embedding_model": self.config.embedding_model,
            "chunk_size": self.config.chunk_size,
        }


# ══════════════════════════════════════════════════════════════════════════════
# SAMPLE DATA FOR DEMO
# ══════════════════════════════════════════════════════════════════════════════

SAMPLE_POLICY_DOCUMENTS = [
    {
        "title": "Employee Handbook - Leave Policy",
        "content": """
# Leave Policy

## Annual Leave
All full-time employees are entitled to 20 days of annual leave per calendar year.
Leave accrues monthly at a rate of 1.67 days per month. Leave can be carried over
to the next year, up to a maximum of 10 days.

### Requesting Leave
1. Submit leave requests at least 2 weeks in advance
2. Use the HR portal or email hr@company.com
3. Await approval from your direct manager
4. Emergency leave may be granted with shorter notice

## Sick Leave
Employees are entitled to 10 days of sick leave per year. Medical certificates
are required for absences of 3 or more consecutive days.

### Sick Leave Notification
- Notify your manager before your shift starts
- Call or text if email is not possible
- Provide medical certificate if required

## Parental Leave
New parents are entitled to 12 weeks of paid parental leave. This can be taken
flexibly within the first year of the child's birth or adoption.
""",
    },
    {
        "title": "IT Security Policy",
        "content": """
# Information Technology Security Policy

## Password Requirements
All passwords must meet the following criteria:
- Minimum 12 characters
- At least one uppercase letter
- At least one lowercase letter
- At least one number
- At least one special character
- Changed every 90 days
- Cannot reuse last 5 passwords

## Device Security
Company devices must:
1. Have full disk encryption enabled
2. Run approved antivirus software
3. Have automatic screen lock after 5 minutes
4. Be reported immediately if lost or stolen

### Remote Work Security
When working remotely:
- Always use the company VPN
- Do not use public WiFi without VPN
- Keep work data on company devices only
- Lock your screen when away from device

## Data Classification
- PUBLIC: Can be shared freely
- INTERNAL: Company employees only
- CONFIDENTIAL: Need-to-know basis
- RESTRICTED: Special authorization required
""",
    },
    {
        "title": "Expense Reimbursement Policy",
        "content": """
# Expense Reimbursement Policy

## Eligible Expenses
The company will reimburse the following business expenses:
- Travel (flights, trains, taxis)
- Accommodation (up to $200/night)
- Meals (up to $50/day)
- Client entertainment (pre-approved)
- Office supplies for remote work

## Submission Process
1. Keep all original receipts
2. Submit expenses within 30 days
3. Use the expense management system
4. Include project code if applicable
5. Add business justification

### Approval Limits
- Under $100: Self-approval
- $100-$500: Manager approval
- $500-$2000: Director approval
- Over $2000: VP approval required

## Travel Booking
All travel must be booked through the approved travel portal.
Economy class is standard for flights under 6 hours.
Business class may be approved for flights over 6 hours.

### Per Diem Rates
- Domestic travel: $75/day
- International travel: Varies by location (see rate table)
- Day trips: $30 meal allowance
""",
    },
    {
        "title": "Code of Conduct",
        "content": """
# Code of Conduct

## Our Values
We are committed to maintaining a workplace that is:
- Respectful and inclusive
- Free from harassment and discrimination
- Safe and healthy
- Ethical and honest

## Expected Behavior
All employees must:
1. Treat colleagues with respect and dignity
2. Communicate professionally
3. Maintain confidentiality
4. Report concerns through proper channels
5. Comply with all laws and regulations

## Prohibited Conduct
The following behaviors are strictly prohibited:
- Harassment of any kind
- Discrimination based on protected characteristics
- Theft or fraud
- Violence or threats
- Substance abuse at work
- Conflicts of interest

## Reporting Violations
Report violations to:
- Your manager
- HR department
- Ethics hotline (anonymous): 1-800-ETHICS
- compliance@company.com

All reports will be investigated confidentially.
""",
    },
]


# ══════════════════════════════════════════════════════════════════════════════
# MAIN DEMO
# ══════════════════════════════════════════════════════════════════════════════


def run_demo():
    """Run interactive document Q&A demo."""
    print("=" * 70)
    print("🧠 Document Q&A System - RAG Demo")
    print("=" * 70)

    # Create pipeline
    config = RAGConfig(chunk_size=400, chunk_overlap=50, top_k=5, rerank_top_k=3)
    pipeline = DocumentRAGPipeline(config)

    # Load sample documents
    print("\n📄 Loading company policy documents...")
    for doc_data in SAMPLE_POLICY_DOCUMENTS:
        doc = pipeline.add_document(
            content=doc_data["content"],
            title=doc_data["title"],
            chunking_strategy=ChunkingStrategy.SEMANTIC,
        )
        print(f"  ✅ {doc.title} ({len(doc.chunks)} chunks)")

    # Show stats
    stats = pipeline.get_stats()
    print(f"\n📊 Pipeline Stats:")
    print(f"   Documents: {stats['documents']}")
    print(f"   Chunks: {stats['chunks']}")
    print(f"   Embedding Model: {stats['embedding_model']}")

    # Sample queries
    sample_queries = [
        "How many days of annual leave do I get?",
        "What are the password requirements?",
        "How do I submit an expense report?",
        "What is the policy on harassment?",
        "Can I work remotely?",
    ]

    print("\n" + "=" * 70)
    print("💡 Sample questions you can ask:")
    for q in sample_queries:
        print(f"   • {q}")

    # Interactive loop
    print("\n" + "=" * 70)
    print("💬 Ask questions about the policies (type 'quit' to exit)")
    print("   Commands: 'eval' - run evaluation, 'stats' - show stats")
    print("=" * 70)

    while True:
        try:
            query = input("\n❓ Your question: ").strip()
        except (EOFError, KeyboardInterrupt):
            break

        if not query:
            continue

        if query.lower() == "quit":
            break

        if query.lower() == "stats":
            stats = pipeline.get_stats()
            print(f"\n📊 Pipeline Statistics:")
            for key, value in stats.items():
                print(f"   {key}: {value}")
            continue

        if query.lower() == "eval":
            print("\n📈 Running evaluation on sample queries...")
            total_metrics = {
                "relevance": 0,
                "accuracy": 0,
                "faithfulness": 0,
                "precision": 0,
                "time": 0,
            }

            for q in sample_queries:
                metrics = pipeline.evaluate(q)
                total_metrics["relevance"] += metrics.relevance_score
                total_metrics["accuracy"] += metrics.answer_accuracy
                total_metrics["faithfulness"] += metrics.faithfulness
                total_metrics["precision"] += metrics.context_precision
                total_metrics["time"] += metrics.response_time_ms

            n = len(sample_queries)
            print(f"\n   Average Metrics ({n} queries):")
            print(f"   Relevance:   {total_metrics['relevance']/n:.2f}")
            print(f"   Accuracy:    {total_metrics['accuracy']/n:.2f}")
            print(f"   Faithfulness: {total_metrics['faithfulness']/n:.2f}")
            print(f"   Precision:   {total_metrics['precision']/n:.2f}")
            print(f"   Avg Time:    {total_metrics['time']/n:.0f}ms")
            continue

        # Query the pipeline
        result = pipeline.query(query)

        print(
            f"\n🤖 Answer (confidence: {result.confidence:.2f}, {result.latency_ms:.0f}ms):\n"
        )
        print(result.answer)

        if result.sources:
            print(f"\n📚 Sources ({len(result.sources)}):")
            for r in result.sources:
                print(
                    f"   • {r.chunk.source_info} (score: {r.score:.2f}, {r.match_type})"
                )

    print("\n👋 Goodbye!")


def main():
    """Main entry point."""
    import sys

    if "--help" in sys.argv or "-h" in sys.argv:
        print(__doc__)
        return

    run_demo()


if __name__ == "__main__":
    main()
