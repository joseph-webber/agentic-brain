# SPDX-License-Identifier: Apache-2.0

"""Unified table extraction across office document formats.

This module provides a single table extraction service that works across common
office formats and exposes a normalized table model with helpers for export and
downstream analysis.
"""

from __future__ import annotations

import csv
import html
import json
import logging
import re
import zipfile
from dataclasses import dataclass, field
from datetime import datetime, timedelta
from html.parser import HTMLParser
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

try:
    import pandas as pd

    PANDAS_AVAILABLE = True
except ImportError:  # pragma: no cover - optional dependency
    pd = None
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None
    PPTX_AVAILABLE = False

from .exceptions import UnsupportedOfficeFormatError
from .models import OfficeFormat

logger = logging.getLogger(__name__)

W_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
X_NS = {
    "x": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "pr": "http://schemas.openxmlformats.org/package/2006/relationships",
}
A_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
ODF_NS = {
    "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
    "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
}

EMPTY_RE = re.compile(r"^\s*$")
INT_RE = re.compile(r"^[+-]?\d+$")
FLOAT_RE = re.compile(r"^[+-]?(?:\d{1,3}(?:,\d{3})+|\d+)(?:\.\d+)?$")
PERCENT_RE = re.compile(r"^[+-]?(?:\d{1,3}(?:,\d{3})+|\d+)(?:\.\d+)?%$")
CURRENCY_RE = re.compile(r"^[\$£€¥]\s*[+-]?(?:\d{1,3}(?:,\d{3})+|\d+)(?:\.\d+)?$")
BOOL_TRUE = {"true", "yes", "y", "1"}
BOOL_FALSE = {"false", "no", "n", "0"}
DATE_HINT_RE = re.compile(r"[dmyhs]", re.IGNORECASE)
BUILTIN_DATE_FORMATS = {14, 15, 16, 17, 22, 27, 30, 36, 45, 46, 47, 50, 57}


@dataclass(slots=True)
class Cell:
    """Normalized cell representation."""

    value: Any = None
    text: str = ""
    data_type: str = "empty"
    rowspan: int = 1
    colspan: int = 1
    is_header: bool = False
    is_placeholder: bool = False
    formula: str | None = None
    style: dict[str, Any] = field(default_factory=dict)
    source: dict[str, Any] = field(default_factory=dict)

    def __post_init__(self) -> None:
        if not self.text and self.value is not None:
            self.text = self._stringify(self.value)
        if self.value is None and self.text:
            inferred_value, inferred_type = infer_scalar(self.text)
            self.value = inferred_value
            self.data_type = inferred_type
        elif self.value is not None and self.data_type == "empty":
            _, inferred_type = infer_scalar(self.text or self._stringify(self.value))
            self.data_type = inferred_type

    @staticmethod
    def _stringify(value: Any) -> str:
        if isinstance(value, datetime):
            return value.isoformat(sep=" ", timespec="seconds")
        if value is None:
            return ""
        return str(value)

    def to_json(self) -> dict[str, Any]:
        """Convert the cell to a JSON-serializable dictionary."""
        value = self.value.isoformat() if isinstance(self.value, datetime) else self.value
        return {
            "value": value,
            "text": self.text,
            "data_type": self.data_type,
            "rowspan": self.rowspan,
            "colspan": self.colspan,
            "is_header": self.is_header,
            "is_placeholder": self.is_placeholder,
            "formula": self.formula,
            "style": self.style,
            "source": self.source,
        }


@dataclass(slots=True)
class Table:
    """Normalized table with convenience export helpers."""

    rows: list[list[Cell]]
    headers: list[str] = field(default_factory=list)
    header_row_count: int = 0
    style: dict[str, Any] = field(default_factory=dict)
    source: dict[str, Any] = field(default_factory=dict)

    @property
    def row_count(self) -> int:
        """Return total row count including header rows."""
        return len(self.rows)

    @property
    def col_count(self) -> int:
        """Return the maximum number of columns in the normalized grid."""
        return max((len(row) for row in self.rows), default=0)

    @property
    def has_merged_cells(self) -> bool:
        """Return whether the table contains merged cells."""
        return any(
            not cell.is_placeholder and (cell.rowspan > 1 or cell.colspan > 1)
            for row in self.rows
            for cell in row
        )

    @property
    def body_rows(self) -> list[list[Cell]]:
        """Return rows excluding detected header rows."""
        return self.rows[self.header_row_count :]

    def get_cell(self, row: int, col: int) -> Cell:
        """Get a cell by zero-based row and column index."""
        if row < 0 or row >= self.row_count:
            raise IndexError(f"Row index {row} out of range")
        if col < 0 or col >= len(self.rows[row]):
            raise IndexError(f"Column index {col} out of range")
        return self.rows[row][col]

    def to_dataframe(self) -> pd.DataFrame:
        """Convert the table to a pandas DataFrame."""
        if not PANDAS_AVAILABLE:
            raise ImportError(
                "pandas is required for DataFrame export. Install with: pip install pandas"
            )
        return pd.DataFrame(self._body_value_matrix(), columns=self._resolved_headers())

    def to_csv(self, path: str | Path) -> Path:
        """Write the table to a CSV file."""
        output_path = Path(path).expanduser()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with output_path.open("w", encoding="utf-8", newline="") as handle:
            writer = csv.writer(handle)
            writer.writerow(self._resolved_headers())
            writer.writerows(self._body_text_matrix())
        return output_path

    def to_markdown(self) -> str:
        """Render the table as GitHub-flavored Markdown."""
        headers = self._resolved_headers()
        if not headers:
            return ""
        lines = [
            "| " + " | ".join(_escape_markdown(header) for header in headers) + " |",
            "| " + " | ".join("---" for _ in headers) + " |",
        ]
        for row in self._body_text_matrix():
            padded = row + [""] * max(0, len(headers) - len(row))
            lines.append(
                "| "
                + " | ".join(_escape_markdown(value) for value in padded[: len(headers)])
                + " |"
            )
        return "\n".join(lines)

    def to_html(self) -> str:
        """Render the table as semantic HTML."""
        headers = self._resolved_headers()
        html_lines = ['<table border="1">']
        if headers:
            html_lines.append("  <thead>")
            html_lines.append(
                "    <tr>"
                + "".join(f"<th>{html.escape(header)}</th>" for header in headers)
                + "</tr>"
            )
            html_lines.append("  </thead>")

        html_lines.append("  <tbody>")
        for row in self.body_rows:
            html_lines.append("    <tr>")
            for cell in row:
                if cell.is_placeholder:
                    continue
                attrs: list[str] = []
                if cell.rowspan > 1:
                    attrs.append(f'rowspan="{cell.rowspan}"')
                if cell.colspan > 1:
                    attrs.append(f'colspan="{cell.colspan}"')
                css = _style_to_css(cell.style)
                if css:
                    attrs.append(f'style="{css}"')
                attr_text = f" {' '.join(attrs)}" if attrs else ""
                html_lines.append(f"      <td{attr_text}>{html.escape(cell.text)}</td>")
            html_lines.append("    </tr>")
        html_lines.append("  </tbody>")
        html_lines.append("</table>")
        return "\n".join(html_lines)

    def to_json(self) -> dict[str, Any]:
        """Convert the table to a JSON-serializable dictionary."""
        headers = self._resolved_headers()
        return {
            "headers": headers,
            "header_row_count": self.header_row_count,
            "row_count": self.row_count,
            "col_count": self.col_count,
            "has_merged_cells": self.has_merged_cells,
            "rows": [[cell.to_json() for cell in row] for row in self.rows],
            "records": [
                dict(zip(headers, row, strict=False)) for row in self._body_value_matrix()
            ],
            "style": self.style,
            "source": self.source,
        }

    def _resolved_headers(self) -> list[str]:
        if self.headers:
            return self.headers
        if self.col_count == 0:
            return []
        return [f"Column {index}" for index in range(1, self.col_count + 1)]

    def _body_text_matrix(self) -> list[list[str]]:
        return [
            [cell.text for cell in row[: self.col_count]]
            + [""] * max(0, self.col_count - len(row))
            for row in self.body_rows
        ]

    def _body_value_matrix(self) -> list[list[Any]]:
        return [
            [cell.value for cell in row[: self.col_count]]
            + [None] * max(0, self.col_count - len(row))
            for row in self.body_rows
        ]


class TableExtractor:
    """Extract tables from office documents with a unified data model."""

    def extract_tables(self, path: str | Path) -> list[Table]:
        """Auto-detect the office format and extract all tables."""
        document_path = self._resolve_path(path)
        format_name = self._detect_format(document_path)
        dispatch = {
            OfficeFormat.DOCX: self.extract_tables_from_docx,
            OfficeFormat.XLSX: self.extract_tables_from_xlsx,
            OfficeFormat.PPTX: self.extract_tables_from_pptx,
            OfficeFormat.PAGES: self.extract_tables_from_pages,
            OfficeFormat.NUMBERS: self.extract_tables_from_numbers,
            OfficeFormat.ODT: self.extract_tables_from_odt,
        }
        extractor = dispatch.get(format_name)
        if extractor is None:
            raise UnsupportedOfficeFormatError(format_name)
        return extractor(document_path)

    def extract_all_tables(self, paths: list[str | Path]) -> dict[str, list[Table]]:
        """Extract tables from multiple documents."""
        return {
            str(self._resolve_path(path)): self.extract_tables(path)
            for path in paths
        }

    def extract_tables_from_docx(self, path: str | Path) -> list[Table]:
        """Extract tables from a DOCX document."""
        document_path = self._require_suffix(path, ".docx")
        tables: list[Table] = []

        with zipfile.ZipFile(document_path) as archive:
            document_root = ET.fromstring(archive.read("word/document.xml"))
            style_names = self._read_docx_style_names(archive)

            for index, table_element in enumerate(document_root.findall(".//w:tbl", W_NS), start=1):
                table_style_id = self._find_attr(
                    table_element.find("w:tblPr/w:tblStyle", W_NS),
                    "val",
                    namespace=W_NS["w"],
                )
                table_style = {
                    "style_id": table_style_id,
                    "style_name": style_names.get(table_style_id or "", table_style_id),
                    "format": "docx",
                }
                grid = self._parse_docx_table(table_element)
                finalized = self._finalize_table(
                    grid,
                    source={"path": str(document_path), "format": "docx", "table_index": index},
                    style=table_style,
                )
                if finalized:
                    tables.append(finalized)

        return tables

    def extract_tables_from_xlsx(self, path: str | Path) -> list[Table]:
        """Extract tables from an XLSX workbook."""
        workbook_path = self._require_suffix(path, ".xlsx")
        tables: list[Table] = []

        with zipfile.ZipFile(workbook_path) as archive:
            workbook_root = ET.fromstring(archive.read("xl/workbook.xml"))
            workbook_rels = self._read_relationships(archive, "xl/_rels/workbook.xml.rels")
            shared_strings = self._read_shared_strings(archive)
            style_map = self._read_xlsx_styles(archive)

            for sheet_index, sheet in enumerate(
                workbook_root.findall("x:sheets/x:sheet", X_NS),
                start=1,
            ):
                name = sheet.get("name", f"Sheet{sheet_index}")
                relation_id = self._find_attr(
                    sheet, "id", namespace=X_NS["r"]
                )
                target = workbook_rels.get(relation_id or "")
                if not target:
                    continue
                target_path = "xl/" + target.lstrip("/")
                sheet_root = ET.fromstring(archive.read(target_path))
                grid = self._parse_xlsx_sheet(sheet_root, shared_strings, style_map)
                for table_number, table_grid in enumerate(self._split_grid_into_tables(grid), start=1):
                    finalized = self._finalize_table(
                        table_grid,
                        source={
                            "path": str(workbook_path),
                            "format": "xlsx",
                            "sheet": name,
                            "sheet_index": sheet_index,
                            "table_index": table_number,
                        },
                        style={"format": "xlsx", "sheet": name},
                    )
                    if finalized:
                        tables.append(finalized)

        return tables

    def extract_tables_from_pptx(self, path: str | Path) -> list[Table]:
        """Extract tables from a PPTX presentation."""
        presentation_path = self._require_suffix(path, ".pptx")
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PPTX table extraction. "
                "Install with: pip install python-pptx"
            )

        presentation = Presentation(str(presentation_path))
        tables: list[Table] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if not getattr(shape, "has_table", False):
                    continue
                grid: list[list[Cell]] = []
                table = shape.table
                for row in table.rows:
                    normalized_row: list[Cell] = []
                    for cell in row.cells:
                        if getattr(cell, "is_spanned", False):
                            normalized_row.append(
                                Cell(
                                    text=cell.text_frame.text.strip()
                                    if cell.text_frame is not None
                                    else "",
                                    is_placeholder=True,
                                    source={"format": "pptx", "merged": True},
                                )
                            )
                            continue

                        cell_style = self._extract_pptx_style(cell)
                        normalized = Cell(
                            text=self._clean_text(cell.text_frame.text if cell.text_frame else ""),
                            rowspan=getattr(cell, "span_height", 1) or 1,
                            colspan=getattr(cell, "span_width", 1) or 1,
                            style=cell_style,
                            source={"format": "pptx"},
                        )
                        normalized_row.append(normalized)
                        for _ in range(1, normalized.colspan):
                            normalized_row.append(
                                Cell(
                                    value=normalized.value,
                                    text=normalized.text,
                                    data_type=normalized.data_type,
                                    is_placeholder=True,
                                    style=dict(normalized.style),
                                    source={"format": "pptx", "merged": True},
                                )
                            )
                    grid.append(normalized_row)

                finalized = self._finalize_table(
                    grid,
                    source={
                        "path": str(presentation_path),
                        "format": "pptx",
                        "slide_index": slide_index,
                        "shape_index": shape_index,
                    },
                    style={"format": "pptx", "shape_name": getattr(shape, "name", None)},
                )
                if finalized:
                    tables.append(finalized)

        return tables

    def extract_tables_from_pages(self, path: str | Path) -> list[Table]:
        """Extract tables from a Pages package using package artifact heuristics."""
        return self._extract_tables_from_iwork_package(path, OfficeFormat.PAGES)

    def extract_tables_from_numbers(self, path: str | Path) -> list[Table]:
        """Extract tables from a Numbers package using package artifact heuristics."""
        return self._extract_tables_from_iwork_package(path, OfficeFormat.NUMBERS)

    def extract_tables_from_odt(self, path: str | Path) -> list[Table]:
        """Extract tables from an ODT document."""
        document_path = self._require_suffix(path, ".odt")
        tables: list[Table] = []

        with zipfile.ZipFile(document_path) as archive:
            root = ET.fromstring(archive.read("content.xml"))
            for index, table_element in enumerate(root.findall(".//table:table", ODF_NS), start=1):
                logical_rows: list[list[dict[str, Any]]] = []
                for row_element in table_element.findall("table:table-row", ODF_NS):
                    repeat_rows = int(
                        row_element.get(f"{{{ODF_NS['table']}}}number-rows-repeated", "1")
                    )
                    row_cells: list[dict[str, Any]] = []
                    for cell_element in list(row_element):
                        tag_name = _local_name(cell_element.tag)
                        if tag_name not in {"table-cell", "covered-table-cell"}:
                            continue

                        repeated = int(
                            cell_element.get(
                                f"{{{ODF_NS['table']}}}number-columns-repeated",
                                "1",
                            )
                        )
                        rowspan = int(
                            cell_element.get(
                                f"{{{ODF_NS['table']}}}number-rows-spanned",
                                "1",
                            )
                        )
                        colspan = int(
                            cell_element.get(
                                f"{{{ODF_NS['table']}}}number-columns-spanned",
                                "1",
                            )
                        )
                        style_name = cell_element.get(f"{{{ODF_NS['table']}}}style-name")
                        text_value = self._clean_text(" ".join(cell_element.itertext()))

                        if tag_name == "covered-table-cell":
                            for _ in range(repeated):
                                row_cells.append(
                                    {
                                        "text": text_value,
                                        "rowspan": 1,
                                        "colspan": 1,
                                        "placeholder": True,
                                        "style": {"style_name": style_name},
                                    }
                                )
                            continue

                        for _ in range(repeated):
                            row_cells.append(
                                {
                                    "text": text_value,
                                    "rowspan": rowspan,
                                    "colspan": colspan,
                                    "style": {"style_name": style_name},
                                }
                            )

                    for _ in range(repeat_rows):
                        logical_rows.append(list(row_cells))

                grid = self._logical_rows_to_grid(logical_rows, source_format="odt")
                finalized = self._finalize_table(
                    grid,
                    source={"path": str(document_path), "format": "odt", "table_index": index},
                    style={
                        "format": "odt",
                        "style_name": table_element.get(f"{{{ODF_NS['table']}}}style-name"),
                        "table_name": table_element.get(f"{{{ODF_NS['table']}}}name"),
                    },
                )
                if finalized:
                    tables.append(finalized)

        return tables

    def _extract_tables_from_iwork_package(
        self, path: str | Path, format_name: OfficeFormat
    ) -> list[Table]:
        document_path = self._require_suffix(path, f".{format_name.value}")
        tables: list[Table] = []
        package_entries = self._iter_package_entries(document_path)

        for entry_name, payload in package_entries:
            suffix = Path(entry_name).suffix.lower()
            try:
                if suffix == ".csv":
                    tables.extend(
                        self._tables_from_delimited_bytes(
                            payload,
                            delimiter=",",
                            source={
                                "path": str(document_path),
                                "format": format_name.value,
                                "entry": entry_name,
                            },
                        )
                    )
                elif suffix == ".tsv":
                    tables.extend(
                        self._tables_from_delimited_bytes(
                            payload,
                            delimiter="\t",
                            source={
                                "path": str(document_path),
                                "format": format_name.value,
                                "entry": entry_name,
                            },
                        )
                    )
                elif suffix in {".html", ".htm"}:
                    tables.extend(
                        self._tables_from_html(
                            payload.decode("utf-8", errors="ignore"),
                            source={
                                "path": str(document_path),
                                "format": format_name.value,
                                "entry": entry_name,
                            },
                        )
                    )
                elif suffix == ".xml":
                    tables.extend(
                        self._tables_from_generic_xml(
                            payload,
                            source={
                                "path": str(document_path),
                                "format": format_name.value,
                                "entry": entry_name,
                            },
                        )
                    )
            except Exception as exc:  # pragma: no cover - best effort parser
                logger.debug("Skipping package entry %s: %s", entry_name, exc)

        return tables

    def _parse_docx_table(self, table_element: ET.Element) -> list[list[Cell]]:
        grid: list[list[Cell]] = []
        active_vertical: dict[int, Cell] = {}

        for row_element in table_element.findall("w:tr", W_NS):
            row: list[Cell] = []
            col_index = 0

            for cell_element in row_element.findall("w:tc", W_NS):
                colspan = int(
                    self._find_attr(
                        cell_element.find("w:tcPr/w:gridSpan", W_NS),
                        "val",
                        namespace=W_NS["w"],
                    )
                    or 1
                )
                vmerge_element = cell_element.find("w:tcPr/w:vMerge", W_NS)
                vmerge_value = (
                    self._find_attr(vmerge_element, "val", namespace=W_NS["w"])
                    if vmerge_element is not None
                    else None
                )
                text_value = self._clean_text("".join(text.text or "" for text in cell_element.findall(".//w:t", W_NS)))
                cell_style = self._extract_docx_cell_style(cell_element)

                if vmerge_element is not None and vmerge_value in {None, "continue"}:
                    anchor = active_vertical.get(col_index)
                    if anchor is None:
                        anchor = Cell(text=text_value, style=cell_style, source={"format": "docx"})
                    anchor.rowspan += 1
                    row.append(
                        Cell(
                            value=anchor.value,
                            text=anchor.text,
                            data_type=anchor.data_type,
                            is_placeholder=True,
                            style=dict(anchor.style),
                            source={"format": "docx", "merged": True},
                        )
                    )
                    for _offset in range(1, colspan):
                        row.append(
                            Cell(
                                value=anchor.value,
                                text=anchor.text,
                                data_type=anchor.data_type,
                                is_placeholder=True,
                                style=dict(anchor.style),
                                source={"format": "docx", "merged": True},
                            )
                        )
                    for offset in range(colspan):
                        active_vertical[col_index + offset] = anchor
                    col_index += colspan
                    continue

                cell = Cell(
                    text=text_value,
                    rowspan=1,
                    colspan=colspan,
                    style=cell_style,
                    source={"format": "docx"},
                )
                row.append(cell)
                for _ in range(1, colspan):
                    row.append(
                        Cell(
                            value=cell.value,
                            text=cell.text,
                            data_type=cell.data_type,
                            is_placeholder=True,
                            style=dict(cell.style),
                            source={"format": "docx", "merged": True},
                        )
                    )

                if vmerge_element is not None and vmerge_value == "restart":
                    for offset in range(colspan):
                        active_vertical[col_index + offset] = cell
                else:
                    for offset in range(colspan):
                        active_vertical.pop(col_index + offset, None)
                col_index += colspan

            grid.append(row)

        return grid

    def _parse_xlsx_sheet(
        self,
        sheet_root: ET.Element,
        shared_strings: list[str],
        style_map: dict[int, dict[str, Any]],
    ) -> list[list[Cell]]:
        raw_cells: dict[tuple[int, int], Cell] = {}
        max_row = 0
        max_col = 0

        for row_element in sheet_root.findall("x:sheetData/x:row", X_NS):
            for cell_element in row_element.findall("x:c", X_NS):
                reference = cell_element.get("r")
                if not reference:
                    continue
                row_number, col_number = self._cell_reference_to_indexes(reference)
                max_row = max(max_row, row_number + 1)
                max_col = max(max_col, col_number + 1)
                cell_type = cell_element.get("t")
                style_id = int(cell_element.get("s", "0"))
                formula = (
                    self._clean_text(cell_element.findtext("x:f", default="", namespaces=X_NS))
                    or None
                )
                value = self._parse_xlsx_value(
                    cell_element,
                    cell_type,
                    shared_strings,
                    style_map.get(style_id, {}),
                )
                text_value = value.isoformat(sep=" ", timespec="seconds") if isinstance(
                    value, datetime
                ) else ("" if value is None else str(value))
                raw_cells[(row_number, col_number)] = Cell(
                    value=value,
                    text=text_value,
                    data_type=self._xlsx_data_type(value, cell_type, style_map.get(style_id, {})),
                    formula=formula,
                    style=style_map.get(style_id, {}),
                    source={"format": "xlsx", "reference": reference, "style_id": style_id},
                )

        for merge in sheet_root.findall("x:mergeCells/x:mergeCell", X_NS):
            ref = merge.get("ref")
            if not ref or ":" not in ref:
                continue
            start_ref, end_ref = ref.split(":", 1)
            start_row, start_col = self._cell_reference_to_indexes(start_ref)
            end_row, end_col = self._cell_reference_to_indexes(end_ref)
            max_row = max(max_row, end_row + 1)
            max_col = max(max_col, end_col + 1)

            anchor = raw_cells.get((start_row, start_col))
            if anchor is None:
                anchor = Cell(source={"format": "xlsx", "reference": start_ref})
                raw_cells[(start_row, start_col)] = anchor

            anchor.rowspan = max(anchor.rowspan, end_row - start_row + 1)
            anchor.colspan = max(anchor.colspan, end_col - start_col + 1)
            anchor.source["merge_range"] = ref

            for row_number in range(start_row, end_row + 1):
                for col_number in range(start_col, end_col + 1):
                    if (row_number, col_number) == (start_row, start_col):
                        continue
                    raw_cells[(row_number, col_number)] = Cell(
                        value=anchor.value,
                        text=anchor.text,
                        data_type=anchor.data_type,
                        is_placeholder=True,
                        style=dict(anchor.style),
                        source={"format": "xlsx", "merge_range": ref},
                    )

        if max_row == 0 or max_col == 0:
            return []

        grid: list[list[Cell]] = []
        for row_number in range(max_row):
            row: list[Cell] = []
            for col_number in range(max_col):
                row.append(raw_cells.get((row_number, col_number), Cell()))
            grid.append(row)
        return grid

    def _logical_rows_to_grid(
        self, logical_rows: list[list[dict[str, Any]]], source_format: str
    ) -> list[list[Cell]]:
        grid: list[list[Cell]] = []
        active_rowspans: dict[int, tuple[Cell, int]] = {}

        for logical_row in logical_rows:
            row: list[Cell] = []
            col_index = 0

            def inject_active_spans(target_row: list[Cell]) -> None:
                nonlocal col_index
                while col_index in active_rowspans:
                    anchor, remaining = active_rowspans[col_index]
                    target_row.append(
                        Cell(
                            value=anchor.value,
                            text=anchor.text,
                            data_type=anchor.data_type,
                            is_placeholder=True,
                            style=dict(anchor.style),
                            source={"format": source_format, "merged": True},
                        )
                    )
                    if remaining <= 1:
                        del active_rowspans[col_index]
                    else:
                        active_rowspans[col_index] = (anchor, remaining - 1)
                    col_index += 1

            inject_active_spans(row)
            for spec in logical_row:
                inject_active_spans(row)
                if spec.get("placeholder"):
                    row.append(
                        Cell(
                            text=self._clean_text(spec.get("text", "")),
                            is_placeholder=True,
                            style=dict(spec.get("style", {})),
                            source={"format": source_format, "merged": True},
                        )
                    )
                    col_index += 1
                    continue

                cell = Cell(
                    text=self._clean_text(spec.get("text", "")),
                    rowspan=max(1, int(spec.get("rowspan", 1))),
                    colspan=max(1, int(spec.get("colspan", 1))),
                    is_header=bool(spec.get("is_header")),
                    style=dict(spec.get("style", {})),
                    source={"format": source_format},
                )
                row.append(cell)
                for _ in range(1, cell.colspan):
                    row.append(
                        Cell(
                            value=cell.value,
                            text=cell.text,
                            data_type=cell.data_type,
                            is_placeholder=True,
                            style=dict(cell.style),
                            source={"format": source_format, "merged": True},
                        )
                    )
                if cell.rowspan > 1:
                    for offset in range(cell.colspan):
                        active_rowspans[col_index + offset] = (cell, cell.rowspan - 1)
                col_index += cell.colspan
            grid.append(row)

        return grid

    def _tables_from_delimited_bytes(
        self, payload: bytes, delimiter: str, source: dict[str, Any]
    ) -> list[Table]:
        text_payload = payload.decode("utf-8", errors="ignore")
        parsed_rows = list(csv.reader(text_payload.splitlines(), delimiter=delimiter))
        grid = [[Cell(text=self._clean_text(value), source=source) for value in row] for row in parsed_rows]
        finalized = self._finalize_table(grid, source=source, style={"format": source["format"]})
        return [finalized] if finalized else []

    def _tables_from_html(self, html_text: str, source: dict[str, Any]) -> list[Table]:
        parser = _HTMLTableParser()
        parser.feed(html_text)
        tables: list[Table] = []
        for index, logical_rows in enumerate(parser.tables, start=1):
            grid = self._logical_rows_to_grid(logical_rows, source_format=source["format"])
            finalized = self._finalize_table(
                grid,
                source={**source, "table_index": index},
                style={"format": source["format"], "source_type": "html"},
            )
            if finalized:
                tables.append(finalized)
        return tables

    def _tables_from_generic_xml(
        self, payload: bytes, source: dict[str, Any]
    ) -> list[Table]:
        root = ET.fromstring(payload)
        tables: list[Table] = []

        for index, table_element in enumerate(
            [
                element
                for element in root.iter()
                if _local_name(element.tag) in {"table", "tbl"}
            ],
            start=1,
        ):
            logical_rows: list[list[dict[str, Any]]] = []
            for row_element in list(table_element):
                if _local_name(row_element.tag) not in {"tr", "table-row", "row"}:
                    continue
                row_specs: list[dict[str, Any]] = []
                for cell_element in list(row_element):
                    if _local_name(cell_element.tag) not in {
                        "td",
                        "th",
                        "tc",
                        "table-cell",
                        "cell",
                    }:
                        continue
                    row_specs.append(
                        {
                            "text": self._clean_text(" ".join(cell_element.itertext())),
                            "rowspan": int(cell_element.get("rowspan", "1")),
                            "colspan": int(cell_element.get("colspan", "1")),
                            "is_header": _local_name(cell_element.tag) == "th",
                            "style": {
                                key: value
                                for key, value in cell_element.attrib.items()
                                if "style" in key.lower() or "class" in key.lower()
                            },
                        }
                    )
                if row_specs:
                    logical_rows.append(row_specs)

            if not logical_rows:
                continue

            grid = self._logical_rows_to_grid(logical_rows, source_format=source["format"])
            finalized = self._finalize_table(
                grid,
                source={**source, "table_index": index},
                style={"format": source["format"], "source_type": "xml"},
            )
            if finalized:
                tables.append(finalized)

        return tables

    def _finalize_table(
        self,
        grid: list[list[Cell]],
        source: dict[str, Any],
        style: dict[str, Any] | None = None,
    ) -> Table | None:
        trimmed = self._trim_grid(grid)
        if not trimmed:
            return None
        normalized = self._pad_grid(trimmed)
        header_row_count = self._detect_header_row_count(normalized)
        headers = self._build_headers(normalized, header_row_count)

        for row_index in range(min(header_row_count, len(normalized))):
            for cell in normalized[row_index]:
                if not cell.is_placeholder:
                    cell.is_header = True

        return Table(
            rows=normalized,
            headers=headers,
            header_row_count=header_row_count,
            style=style or {},
            source=source,
        )

    def _trim_grid(self, grid: list[list[Cell]]) -> list[list[Cell]]:
        if not grid:
            return []

        active_rows = [
            index for index, row in enumerate(grid) if any(self._cell_occupies_space(cell) for cell in row)
        ]
        if not active_rows:
            return []

        trimmed_rows = [grid[index] for index in range(active_rows[0], active_rows[-1] + 1)]
        max_cols = max((len(row) for row in trimmed_rows), default=0)
        active_cols = [
            col
            for col in range(max_cols)
            if any(
                col < len(row) and self._cell_occupies_space(row[col])
                for row in trimmed_rows
            )
        ]
        if not active_cols:
            return []

        start_col = active_cols[0]
        end_col = active_cols[-1] + 1
        return [row[start_col:end_col] for row in trimmed_rows]

    def _pad_grid(self, grid: list[list[Cell]]) -> list[list[Cell]]:
        width = max((len(row) for row in grid), default=0)
        return [row + [Cell() for _ in range(width - len(row))] for row in grid]

    def _split_grid_into_tables(self, grid: list[list[Cell]]) -> list[list[list[Cell]]]:
        if not grid:
            return []

        occupied = {
            (row_index, col_index)
            for row_index, row in enumerate(grid)
            for col_index, cell in enumerate(row)
            if self._cell_occupies_space(cell)
        }
        components: list[set[tuple[int, int]]] = []

        while occupied:
            start = occupied.pop()
            queue = [start]
            component = {start}
            while queue:
                row_index, col_index = queue.pop()
                for neighbor in (
                    (row_index - 1, col_index),
                    (row_index + 1, col_index),
                    (row_index, col_index - 1),
                    (row_index, col_index + 1),
                ):
                    if neighbor in occupied:
                        occupied.remove(neighbor)
                        component.add(neighbor)
                        queue.append(neighbor)
            components.append(component)

        tables: list[list[list[Cell]]] = []
        for component in sorted(components, key=lambda items: (min(r for r, _ in items), min(c for _, c in items))):
            min_row = min(row for row, _ in component)
            max_row = max(row for row, _ in component)
            min_col = min(col for _, col in component)
            max_col = max(col for _, col in component)
            subgrid: list[list[Cell]] = []
            for row_number in range(min_row, max_row + 1):
                subgrid.append(
                    [
                        grid[row_number][col_number]
                        if col_number < len(grid[row_number])
                        else Cell()
                        for col_number in range(min_col, max_col + 1)
                    ]
                )
            tables.append(subgrid)
        return tables

    def _detect_header_row_count(self, grid: list[list[Cell]]) -> int:
        if not grid:
            return 0

        max_candidates = min(3, len(grid))
        scores = [self._header_score(grid[index], grid[index + 1] if index + 1 < len(grid) else None) for index in range(max_candidates)]
        if not scores or max(scores) < 0.45:
            return 0

        header_rows = 0
        for index, score in enumerate(scores):
            if score < 0.45:
                break
            header_rows += 1
            if index >= 1 and score < 0.6:
                break
        return max(1, min(header_rows, 2))

    def _header_score(self, row: list[Cell], next_row: list[Cell] | None) -> float:
        visible_cells = [cell for cell in row if not cell.is_placeholder and cell.text]
        if not visible_cells:
            return 0.0

        text_ratio = sum(cell.data_type == "text" for cell in visible_cells) / len(visible_cells)
        unique_ratio = len({cell.text.casefold() for cell in visible_cells}) / len(visible_cells)
        compact_ratio = sum(len(cell.text.split()) <= 5 for cell in visible_cells) / len(visible_cells)
        numeric_ratio = sum(cell.data_type in {"int", "float"} for cell in visible_cells) / len(visible_cells)
        score = (text_ratio * 0.45) + (unique_ratio * 0.2) + (compact_ratio * 0.2) - (
            numeric_ratio * 0.35
        )

        if next_row:
            next_cells = [cell for cell in next_row if not cell.is_placeholder and cell.text]
            if next_cells:
                next_numeric = sum(cell.data_type in {"int", "float"} for cell in next_cells) / len(next_cells)
                score += next_numeric * 0.15
        return score

    def _build_headers(self, grid: list[list[Cell]], header_row_count: int) -> list[str]:
        if not grid:
            return []
        if header_row_count <= 0:
            return [f"Column {index}" for index in range(1, len(grid[0]) + 1)]

        header_rows = grid[:header_row_count]
        headers: list[str] = []
        seen: dict[str, int] = {}
        for column_index in range(len(grid[0])):
            parts = []
            for row in header_rows:
                if column_index >= len(row):
                    continue
                cell = row[column_index]
                if cell.is_placeholder or not cell.text:
                    continue
                parts.append(cell.text)
            header = " / ".join(dict.fromkeys(parts)) or f"Column {column_index + 1}"
            count = seen.get(header.casefold(), 0)
            seen[header.casefold()] = count + 1
            if count:
                header = f"{header} ({count + 1})"
            headers.append(header)
        return headers

    def _read_docx_style_names(self, archive: zipfile.ZipFile) -> dict[str, str]:
        style_names: dict[str, str] = {}
        try:
            styles_root = ET.fromstring(archive.read("word/styles.xml"))
        except KeyError:
            return style_names

        for style in styles_root.findall("w:style", W_NS):
            style_id = self._find_attr(style, "styleId", namespace=W_NS["w"])
            style_name = self._find_attr(style.find("w:name", W_NS), "val", namespace=W_NS["w"])
            if style_id and style_name:
                style_names[style_id] = style_name
        return style_names

    def _extract_docx_cell_style(self, cell_element: ET.Element) -> dict[str, Any]:
        tc_pr = cell_element.find("w:tcPr", W_NS)
        paragraph_alignment = self._find_attr(
            cell_element.find(".//w:pPr/w:jc", W_NS),
            "val",
            namespace=W_NS["w"],
        )
        if tc_pr is None:
            return {"alignment": paragraph_alignment}

        return {
            "alignment": paragraph_alignment,
            "vertical_align": self._find_attr(
                tc_pr.find("w:vAlign", W_NS), "val", namespace=W_NS["w"]
            ),
            "background_color": self._find_attr(
                tc_pr.find("w:shd", W_NS), "fill", namespace=W_NS["w"]
            ),
            "width": self._find_attr(tc_pr.find("w:tcW", W_NS), "w", namespace=W_NS["w"]),
        }

    def _read_relationships(
        self, archive: zipfile.ZipFile, rels_path: str
    ) -> dict[str, str]:
        try:
            root = ET.fromstring(archive.read(rels_path))
        except KeyError:
            return {}
        relationships: dict[str, str] = {}
        for relationship in root.findall(
            "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"
        ):
            relation_id = relationship.get("Id")
            target = relationship.get("Target")
            if relation_id and target:
                relationships[relation_id] = target
        return relationships

    def _read_shared_strings(self, archive: zipfile.ZipFile) -> list[str]:
        try:
            root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
        except KeyError:
            return []
        values: list[str] = []
        for item in root.findall("x:si", X_NS):
            values.append(self._clean_text("".join(item.itertext())))
        return values

    def _read_xlsx_styles(self, archive: zipfile.ZipFile) -> dict[int, dict[str, Any]]:
        try:
            root = ET.fromstring(archive.read("xl/styles.xml"))
        except KeyError:
            return {}

        custom_formats = {
            int(num_fmt.get("numFmtId", "0")): num_fmt.get("formatCode", "")
            for num_fmt in root.findall("x:numFmts/x:numFmt", X_NS)
        }
        style_map: dict[int, dict[str, Any]] = {}
        for index, xf in enumerate(root.findall("x:cellXfs/x:xf", X_NS)):
            num_fmt_id = int(xf.get("numFmtId", "0"))
            format_code = custom_formats.get(num_fmt_id, "")
            alignment = xf.find("x:alignment", X_NS)
            style_map[index] = {
                "num_fmt_id": num_fmt_id,
                "format_code": format_code,
                "is_date": num_fmt_id in BUILTIN_DATE_FORMATS
                or bool(DATE_HINT_RE.search(format_code)),
                "horizontal": alignment.get("horizontal") if alignment is not None else None,
                "vertical": alignment.get("vertical") if alignment is not None else None,
            }
        return style_map

    def _parse_xlsx_value(
        self,
        cell_element: ET.Element,
        cell_type: str | None,
        shared_strings: list[str],
        style: dict[str, Any],
    ) -> Any:
        raw_value = cell_element.findtext("x:v", default="", namespaces=X_NS)
        if raw_value == "":
            inline = cell_element.find("x:is", X_NS)
            return self._clean_text("".join(inline.itertext())) if inline is not None else None
        if cell_type == "s":
            index = int(raw_value)
            return shared_strings[index] if index < len(shared_strings) else raw_value
        if cell_type == "b":
            return raw_value == "1"
        if cell_type == "d":
            try:
                return datetime.fromisoformat(raw_value)
            except ValueError:
                return raw_value
        if style.get("is_date"):
            try:
                serial = float(raw_value)
                return datetime(1899, 12, 30) + timedelta(days=serial)
            except ValueError:
                return raw_value
        if INT_RE.match(raw_value):
            return int(raw_value)
        if FLOAT_RE.match(raw_value):
            return float(raw_value.replace(",", ""))
        return raw_value

    def _xlsx_data_type(
        self, value: Any, cell_type: str | None, style: dict[str, Any]
    ) -> str:
        if value is None or value == "":
            return "empty"
        if isinstance(value, bool):
            return "bool"
        if isinstance(value, int):
            return "int"
        if isinstance(value, float):
            return "date" if style.get("is_date") else "float"
        if isinstance(value, datetime):
            return "date"
        if cell_type == "str":
            return "text"
        return infer_scalar(str(value))[1]

    def _extract_pptx_style(self, cell: Any) -> dict[str, Any]:
        style: dict[str, Any] = {}
        try:
            fill = cell.fill
            if getattr(fill, "fore_color", None) is not None:
                rgb = getattr(fill.fore_color, "rgb", None)
                if rgb is not None:
                    style["background_color"] = f"#{rgb}"
        except Exception:  # pragma: no cover - defensive
            pass

        try:
            style["margin_left"] = cell.margin_left
            style["margin_right"] = cell.margin_right
            style["margin_top"] = cell.margin_top
            style["margin_bottom"] = cell.margin_bottom
        except Exception:  # pragma: no cover - defensive
            pass
        return style

    def _iter_package_entries(self, path: Path) -> list[tuple[str, bytes]]:
        if path.is_dir():
            return [
                (str(entry.relative_to(path)), entry.read_bytes())
                for entry in path.rglob("*")
                if entry.is_file()
            ]
        if zipfile.is_zipfile(path):
            with zipfile.ZipFile(path) as archive:
                return [(name, archive.read(name)) for name in archive.namelist() if not name.endswith("/")]
        raise ValueError(f"Unsupported package container for {path}")

    def _resolve_path(self, path: str | Path) -> Path:
        resolved_path = Path(path).expanduser().resolve()
        if not resolved_path.exists():
            raise FileNotFoundError(f"File not found: {resolved_path}")
        return resolved_path

    def _require_suffix(self, path: str | Path, suffix: str) -> Path:
        resolved_path = self._resolve_path(path)
        if resolved_path.suffix.lower() != suffix:
            raise ValueError(f"Expected {suffix} file, got {resolved_path.name}")
        return resolved_path

    def _detect_format(self, path: Path) -> OfficeFormat:
        mapping = {
            ".docx": OfficeFormat.DOCX,
            ".xlsx": OfficeFormat.XLSX,
            ".pptx": OfficeFormat.PPTX,
            ".pages": OfficeFormat.PAGES,
            ".numbers": OfficeFormat.NUMBERS,
            ".odt": OfficeFormat.ODT,
        }
        try:
            return mapping[path.suffix.lower()]
        except KeyError as exc:
            raise UnsupportedOfficeFormatError(path.suffix.lower().lstrip(".")) from exc

    def _cell_reference_to_indexes(self, reference: str) -> tuple[int, int]:
        letters = "".join(character for character in reference if character.isalpha())
        numbers = "".join(character for character in reference if character.isdigit())
        column = 0
        for letter in letters.upper():
            column = (column * 26) + (ord(letter) - 64)
        return int(numbers) - 1, column - 1

    def _find_attr(
        self, element: ET.Element | None, name: str, namespace: str | None = None
    ) -> str | None:
        if element is None:
            return None
        if namespace:
            return element.get(f"{{{namespace}}}{name}")
        return element.get(name)

    def _clean_text(self, value: str | None) -> str:
        if value is None:
            return ""
        return " ".join(value.replace("\xa0", " ").split())

    def _cell_occupies_space(self, cell: Cell) -> bool:
        return cell.is_placeholder or bool(cell.text) or cell.rowspan > 1 or cell.colspan > 1


class _HTMLTableParser(HTMLParser):
    """Minimal HTML table parser preserving spans and header cells."""

    def __init__(self) -> None:
        super().__init__()
        self.tables: list[list[list[dict[str, Any]]]] = []
        self._table_stack: list[list[list[dict[str, Any]]]] = []
        self._current_row: list[dict[str, Any]] | None = None
        self._current_cell: dict[str, Any] | None = None

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attributes = dict(attrs)
        if tag == "table":
            table: list[list[dict[str, Any]]] = []
            self._table_stack.append(table)
        elif tag == "tr" and self._table_stack:
            self._current_row = []
        elif tag in {"td", "th"} and self._current_row is not None:
            self._current_cell = {
                "text": "",
                "rowspan": int(attributes.get("rowspan", "1") or "1"),
                "colspan": int(attributes.get("colspan", "1") or "1"),
                "is_header": tag == "th",
                "style": {
                    key: value
                    for key, value in attributes.items()
                    if key in {"style", "class"}
                },
            }

    def handle_data(self, data: str) -> None:
        if self._current_cell is not None:
            self._current_cell["text"] += data

    def handle_endtag(self, tag: str) -> None:
        if tag in {"td", "th"} and self._current_cell is not None and self._current_row is not None:
            self._current_cell["text"] = " ".join(self._current_cell["text"].split())
            self._current_row.append(self._current_cell)
            self._current_cell = None
        elif tag == "tr" and self._current_row is not None and self._table_stack:
            self._table_stack[-1].append(self._current_row)
            self._current_row = None
        elif tag == "table" and self._table_stack:
            table = self._table_stack.pop()
            self.tables.append(table)


def infer_scalar(value: str) -> tuple[Any, str]:
    """Infer a scalar value and a normalized data type from text."""
    cleaned = " ".join(value.replace("\xa0", " ").split())
    if not cleaned:
        return None, "empty"

    lowered = cleaned.casefold()
    if lowered in BOOL_TRUE:
        return True, "bool"
    if lowered in BOOL_FALSE:
        return False, "bool"

    if PERCENT_RE.match(cleaned):
        return float(cleaned.rstrip("%").replace(",", "")) / 100.0, "float"
    if CURRENCY_RE.match(cleaned):
        return float(cleaned[1:].strip().replace(",", "")), "float"
    if INT_RE.match(cleaned):
        return int(cleaned), "int"
    if FLOAT_RE.match(cleaned):
        return float(cleaned.replace(",", "")), "float"

    for fmt in (
        "%Y-%m-%d",
        "%Y/%m/%d",
        "%d/%m/%Y",
        "%m/%d/%Y",
        "%d-%m-%Y",
        "%m-%d-%Y",
        "%Y-%m-%d %H:%M:%S",
        "%Y/%m/%d %H:%M:%S",
    ):
        try:
            return datetime.strptime(cleaned, fmt), "date"
        except ValueError:
            continue
    return cleaned, "text"


def _escape_markdown(value: str) -> str:
    return str(value).replace("|", "\\|").replace("\n", "<br>")


def _style_to_css(style: dict[str, Any]) -> str:
    css_parts: list[str] = []
    if background := style.get("background_color"):
        css_parts.append(f"background-color: {background}")
    if alignment := style.get("alignment") or style.get("horizontal"):
        css_parts.append(f"text-align: {alignment}")
    if style.get("bold"):
        css_parts.append("font-weight: bold")
    return "; ".join(css_parts)


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


__all__ = ["Cell", "PANDAS_AVAILABLE", "Table", "TableExtractor", "infer_scalar"]
