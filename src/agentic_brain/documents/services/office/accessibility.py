"""Accessibility support for Office documents.

This module integrates the Phase 3 accessibility stack with Office documents,
primarily OOXML formats:

- Word (.docx)
- Excel (.xlsx)
- PowerPoint (.pptx)

It reuses the shared accessibility models and engines from
``agentic_brain.documents.services.accessibility`` while adding Office-aware
heuristics for:

- alternative text management
- reading order analysis
- heading and list validation
- table accessibility checks
- color contrast analysis
- aggregated WCAG-style reporting

The implementation is intentionally best-effort. Some Office accessibility
features are not fully exposed through the upstream Python libraries, so this
processor combines high-level document APIs with direct OOXML package/XML
inspection when necessary.
"""

from __future__ import annotations

import io
import logging
import mimetypes
import posixpath
import re
import shutil
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable, Optional
from xml.etree import ElementTree as ET

from .exceptions import DocumentValidationError, UnsupportedOfficeFormatError
from .models import Image
from ..accessibility.alttext import AltTextGenerator
from ..accessibility.contrast import (
    Color,
    ContrastChecker,
    TextSize,
    WCAGLevel as ContrastWCAGLevel,
    suggest_accessible_color,
)
from ..accessibility.models import (
    AccessibilityIssue,
    AccessibilityReport,
    AccessibilityStandard,
    ContrastIssue,
    HeadingStructure,
    ImageDescription,
    IssueSeverity,
    IssueType,
    ReadingOrderItem,
    TableSummary,
)
from ..accessibility.tables import CellType, TableAccessibilityEngine

logger = logging.getLogger(__name__)

LIST_MARKER_RE = re.compile(
    r"^\s*(?:[\u2022\u2023\u25E6\u2043\u2219\-*+]|(?:\d+|[A-Za-z]|[ivxlcdmIVXLCDM]+)[\.\)])\s+"
)
HEX_COLOR_RE = re.compile(r"^[0-9A-Fa-f]{6,8}$")

OOXML_EXTENSIONS = {"docx", "xlsx", "pptx"}

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
    "dc": "http://purl.org/dc/elements/1.1/",
    "dcterms": "http://purl.org/dc/terms/",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
    "pr": "http://schemas.openxmlformats.org/package/2006/relationships",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
    "x": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "xdr": "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing",
    "xsi": "http://www.w3.org/2001/XMLSchema-instance",
}

SEVERITY_WEIGHTS = {
    IssueSeverity.CRITICAL: 20.0,
    IssueSeverity.MAJOR: 10.0,
    IssueSeverity.MINOR: 4.0,
    IssueSeverity.ADVISORY: 1.5,
}


try:  # pragma: no cover - optional dependency
    from PIL import Image as PILImage
except ImportError:  # pragma: no cover - optional dependency
    PILImage = None

try:  # pragma: no cover - optional dependency
    from docx import Document as DocxDocument
    from docx.table import Table as DocxTableElement
    from docx.text.paragraph import Paragraph as DocxParagraphElement

    DOCX_AVAILABLE = True
except ImportError:  # pragma: no cover - optional dependency
    DocxDocument = None
    DocxParagraphElement = None
    DocxTableElement = None
    DOCX_AVAILABLE = False

try:  # pragma: no cover - optional dependency
    from openpyxl import load_workbook
    from openpyxl.styles import PatternFill
    from openpyxl.utils.cell import range_boundaries

    OPENPYXL_AVAILABLE = True
except ImportError:  # pragma: no cover - optional dependency
    load_workbook = None
    PatternFill = None
    range_boundaries = None
    OPENPYXL_AVAILABLE = False

try:  # pragma: no cover - optional dependency
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None
    MSO_SHAPE_TYPE = None
    PP_PLACEHOLDER = None
    PPTX_AVAILABLE = False


@dataclass(slots=True)
class ReadingOrderReport:
    """Summarized Office reading-order analysis."""

    document_path: Path
    file_type: str
    logical_order: bool
    total_elements: int
    items: list[ReadingOrderItem] = field(default_factory=list)
    issues: list[str] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(slots=True)
class ContrastFix:
    """Suggested change for a contrast failure."""

    element_id: str
    page: int
    original_color: str
    background_color: str
    suggested_color: str
    original_ratio: float
    suggested_ratio: float
    required_ratio: float
    description: str = ""


@dataclass(slots=True)
class WCAGReport:
    """Aggregated WCAG-style report for an Office document."""

    document_path: Path
    level: str
    is_compliant: bool
    score: float
    issues: list[AccessibilityIssue] = field(default_factory=list)
    category_counts: dict[str, int] = field(default_factory=dict)
    recommendations: list[str] = field(default_factory=list)


@dataclass(slots=True)
class Suggestion:
    """Remediation recommendation for an Office accessibility issue."""

    category: str
    priority: str
    message: str
    target_id: Optional[str] = None
    auto_fixable: bool = False


@dataclass(slots=True)
class _OfficeImageRecord:
    """Internal representation of an OOXML image and its accessibility state."""

    image_id: str
    xml_id: str
    display_name: str
    part_name: str
    media_target: str
    location: str
    page: int
    alt_text: str
    title: str
    data: bytes


@dataclass(slots=True)
class _ContrastCandidate:
    """Internal color-contrast candidate."""

    element_id: str
    page: int
    foreground: str
    background: str
    text_size: TextSize
    description: str


class OfficeAccessibilityProcessor:
    """Phase 3 accessibility processor for Office documents."""

    def __init__(self, alt_text_backend: str = "auto") -> None:
        self.alt_text_backend = alt_text_backend
        self.table_engine = TableAccessibilityEngine(
            auto_detect_headers=True,
            generate_summaries=True,
        )
        self.contrast_checker = ContrastChecker()
        self._alt_text_generator: AltTextGenerator | None = None

    @property
    def alt_text_generator(self) -> AltTextGenerator:
        """Lazily initialize the shared alt-text generator."""
        if self._alt_text_generator is None:
            self._alt_text_generator = AltTextGenerator(backend=self.alt_text_backend)
        return self._alt_text_generator

    def check_alt_text(self, path: str | Path) -> list[AccessibilityIssue]:
        """Return accessibility issues for images missing alternative text."""
        source = self._validate_path(path)
        issues: list[AccessibilityIssue] = []
        for record in self._iter_images(source):
            if record.alt_text.strip() or record.title.strip():
                continue
            issues.append(
                self._issue(
                    issue_type=IssueType.MISSING_ALT_TEXT,
                    severity=IssueSeverity.CRITICAL,
                    page=record.page,
                    description=(
                        f"Image '{record.display_name}' at {record.location} "
                        "is missing alternative text."
                    ),
                    wcag="1.1.1",
                    hint="Add descriptive alt text or mark the image as decorative.",
                    auto_fixable=True,
                )
            )
        return issues

    def get_images_without_alt_text(self, path: str | Path) -> list[Image]:
        """Return image objects for Office images without alternative text."""
        source = self._validate_path(path)
        images: list[Image] = []
        for record in self._iter_images(source):
            if record.alt_text.strip() or record.title.strip():
                continue
            images.append(
                Image(
                    data=record.data,
                    mime_type=self._guess_mime(record.media_target),
                    title=record.title or record.display_name,
                    alternate_text=record.alt_text or None,
                    properties={
                        "image_id": record.image_id,
                        "location": record.location,
                        "resource": record.media_target,
                    },
                )
            )
        return images

    def add_alt_text(
        self,
        path: str | Path,
        image_id: str,
        alt_text: str,
        output_path: str | Path,
    ) -> Path:
        """Write alt text for a specific Office image."""
        source = self._validate_path(path)
        record = next(
            (candidate for candidate in self._iter_images(source) if candidate.image_id == image_id),
            None,
        )
        if record is None:
            raise DocumentValidationError(
                "Image not found in Office document",
                details=image_id,
            )
        return self._apply_alt_text_updates(
            source,
            Path(output_path),
            {record.part_name: {record.xml_id: alt_text.strip()}},
        )

    def generate_alt_text(self, path: str | Path, output_path: str | Path) -> Path:
        """Generate and write alt text for Office images that do not have any."""
        source = self._validate_path(path)
        updates: dict[str, dict[str, str]] = {}
        for record in self._iter_images(source):
            if record.alt_text.strip() or record.title.strip():
                continue
            generated = self._generate_image_alt_text(record)
            updates.setdefault(record.part_name, {})[record.xml_id] = generated
        if not updates:
            return self._copy_document(source, output_path)
        return self._apply_alt_text_updates(source, Path(output_path), updates)

    def analyze_reading_order(self, path: str | Path) -> ReadingOrderReport:
        """Analyze Office reading order using format-specific heuristics."""
        source = self._validate_path(path)
        file_type = self._detect_file_type(source)
        if file_type == "docx":
            return self._analyze_docx_reading_order(source)
        if file_type == "pptx":
            return self._analyze_pptx_reading_order(source)
        if file_type == "xlsx":
            return self._analyze_xlsx_reading_order(source)
        raise UnsupportedOfficeFormatError(file_type)

    def check_logical_order(self, path: str | Path) -> bool:
        """Return whether the Office document appears to use a logical order."""
        return self.analyze_reading_order(path).logical_order

    def fix_reading_order(self, path: str | Path, output_path: str | Path) -> Path:
        """Best-effort reading-order remediation for Office documents."""
        source = self._validate_path(path)
        file_type = self._detect_file_type(source)
        if file_type == "pptx":
            return self._fix_pptx_reading_order(source, Path(output_path))
        return self._copy_document(source, output_path)

    def check_heading_hierarchy(self, path: str | Path) -> list[AccessibilityIssue]:
        """Check heading hierarchy and missing title/title-slide structure."""
        source = self._validate_path(path)
        file_type = self._detect_file_type(source)
        if file_type == "docx":
            return self._check_docx_headings(source)
        if file_type == "pptx":
            return self._check_pptx_headings(source)
        return []

    def check_list_structure(self, path: str | Path) -> list[AccessibilityIssue]:
        """Check for manually-authored or malformed lists."""
        source = self._validate_path(path)
        file_type = self._detect_file_type(source)
        if file_type == "docx":
            return self._check_docx_lists(source)
        if file_type == "pptx":
            return self._check_pptx_lists(source)
        return []

    def add_document_title(
        self,
        path: str | Path,
        title: str,
        output_path: str | Path,
    ) -> Path:
        """Set the Office document core title."""
        source = self._validate_path(path)
        file_type = self._detect_file_type(source)
        target = Path(output_path)
        if file_type == "docx":
            self._require_dependency(DOCX_AVAILABLE, "python-docx", file_type)
            document = DocxDocument(source)
            document.core_properties.title = title
            target.parent.mkdir(parents=True, exist_ok=True)
            document.save(target)
            return target
        if file_type == "pptx":
            self._require_dependency(PPTX_AVAILABLE, "python-pptx", file_type)
            presentation = Presentation(source)
            presentation.core_properties.title = title
            target.parent.mkdir(parents=True, exist_ok=True)
            presentation.save(target)
            return target
        if file_type == "xlsx":
            self._require_dependency(OPENPYXL_AVAILABLE, "openpyxl", file_type)
            workbook = load_workbook(source)
            workbook.properties.title = title
            target.parent.mkdir(parents=True, exist_ok=True)
            workbook.save(target)
            return target
        raise UnsupportedOfficeFormatError(file_type)

    def check_table_headers(self, path: str | Path) -> list[AccessibilityIssue]:
        """Check for tables missing discoverable header rows/columns."""
        source = self._validate_path(path)
        issues: list[AccessibilityIssue] = []
        for page, table_id, matrix in self._iter_table_matrices(source):
            analyzed = self.table_engine.analyze_table(matrix, page=page)
            if not analyzed.has_header_row and not analyzed.has_header_col:
                issues.append(
                    self._issue(
                        issue_type=IssueType.MISSING_TABLE_HEADERS,
                        severity=IssueSeverity.MAJOR,
                        page=page,
                        description=f"Table '{table_id}' is missing header rows or columns.",
                        wcag="1.3.1",
                        hint="Identify table headers so screen readers can associate data cells.",
                        auto_fixable=True,
                    )
                )
        return issues

    def add_table_headers(
        self,
        path: str | Path,
        table_id: str,
        headers: list[str],
        output_path: str | Path,
    ) -> Path:
        """Apply a header row to a specific Office table."""
        source = self._validate_path(path)
        file_type = self._detect_file_type(source)
        target = Path(output_path)
        if file_type == "docx":
            return self._add_docx_table_headers(source, table_id, headers, target)
        if file_type == "pptx":
            return self._add_pptx_table_headers(source, table_id, headers, target)
        if file_type == "xlsx":
            return self._add_xlsx_table_headers(source, table_id, headers, target)
        raise UnsupportedOfficeFormatError(file_type)

    def check_table_scope(self, path: str | Path) -> list[AccessibilityIssue]:
        """Check that table headers meaningfully scope data cells."""
        source = self._validate_path(path)
        issues: list[AccessibilityIssue] = []
        for page, table_id, matrix in self._iter_table_matrices(source):
            analyzed = self.table_engine.analyze_table(matrix, page=page)
            if not analyzed.cells:
                continue
            unresolved = 0
            for cell in analyzed.cells:
                if cell.cell_type is not CellType.DATA or not cell.content.strip():
                    continue
                if not analyzed.get_headers_for_cell(cell.row, cell.col):
                    unresolved += 1
            if unresolved:
                issues.append(
                    self._issue(
                        issue_type=IssueType.MISSING_TABLE_HEADERS,
                        severity=IssueSeverity.MINOR,
                        page=page,
                        description=(
                            f"Table '{table_id}' has {unresolved} data cell(s) without "
                            "clear header scope."
                        ),
                        wcag="1.3.1",
                        hint="Add or normalize row and column headers for ambiguous table cells.",
                        auto_fixable=False,
                    )
                )
        return issues

    def check_color_contrast(self, path: str | Path) -> list[ContrastIssue]:
        """Analyze text/background color contrast in Office documents."""
        source = self._validate_path(path)
        issues: list[ContrastIssue] = []
        for candidate in self._collect_contrast_candidates(source):
            result = self.contrast_checker.check_contrast(
                candidate.foreground,
                candidate.background,
                candidate.text_size,
            )
            required_ratio = 3.0 if candidate.text_size is TextSize.LARGE else 4.5
            if result.ratio >= required_ratio:
                continue
            issues.append(
                ContrastIssue(
                    element_id=candidate.element_id,
                    page=candidate.page,
                    ratio=round(result.ratio, 2),
                    required_ratio=required_ratio,
                    passes=False,
                    description=(
                        f"{candidate.description} uses {candidate.foreground} on "
                        f"{candidate.background}, which is below the required "
                        f"{required_ratio}:1 contrast ratio."
                    ),
                )
            )
        return issues

    def suggest_contrast_fixes(self, path: str | Path) -> list[ContrastFix]:
        """Suggest accessible color replacements for Office contrast failures."""
        source = self._validate_path(path)
        fixes: list[ContrastFix] = []
        for candidate in self._collect_contrast_candidates(source):
            text_level = TextSize.LARGE if candidate.text_size is TextSize.LARGE else TextSize.NORMAL
            current = self.contrast_checker.check_contrast(
                candidate.foreground,
                candidate.background,
                text_level,
            )
            required_ratio = 3.0 if text_level is TextSize.LARGE else 4.5
            if current.ratio >= required_ratio:
                continue
            suggestion = suggest_accessible_color(
                candidate.foreground,
                candidate.background,
                level=ContrastWCAGLevel.AA,
                text_size=text_level,
            )
            suggested = suggestion["suggested_fg"].to_hex()
            fixes.append(
                ContrastFix(
                    element_id=candidate.element_id,
                    page=candidate.page,
                    original_color=candidate.foreground,
                    background_color=candidate.background,
                    suggested_color=suggested,
                    original_ratio=round(float(suggestion["original_ratio"]), 2),
                    suggested_ratio=round(float(suggestion["suggested_ratio"]), 2),
                    required_ratio=required_ratio,
                    description=candidate.description,
                )
            )
        return fixes

    def check_wcag_compliance(self, path: str | Path, level: str = "AA") -> WCAGReport:
        """Run an aggregated WCAG-style accessibility analysis for Office content."""
        source = self._validate_path(path)
        report = self.generate_accessibility_report(source)
        requested_level = level.upper().strip()
        category_counts: dict[str, int] = {}
        recommendations: list[str] = []
        for issue in report.issues:
            key = issue.issue_type.value
            category_counts[key] = category_counts.get(key, 0) + 1
            if issue.remediation_hint and issue.remediation_hint not in recommendations:
                recommendations.append(issue.remediation_hint)
        is_compliant = report.score >= (92.0 if requested_level == "AAA" else 85.0) and not any(
            issue.severity in {IssueSeverity.CRITICAL, IssueSeverity.MAJOR}
            for issue in report.issues
        )
        return WCAGReport(
            document_path=source,
            level=requested_level,
            is_compliant=is_compliant,
            score=report.score,
            issues=report.issues,
            category_counts=category_counts,
            recommendations=recommendations[:10],
        )

    def generate_accessibility_report(self, path: str | Path) -> AccessibilityReport:
        """Generate a Phase 3 accessibility report for an Office document."""
        source = self._validate_path(path)
        issues = self._collect_all_issues(source)
        image_descriptions = self._build_image_descriptions(source)
        headings = self._build_heading_structures(source)
        tables = self._build_table_summaries(source)
        reading_order = self.analyze_reading_order(source).items
        score = max(
            0.0,
            100.0 - sum(SEVERITY_WEIGHTS.get(issue.severity, 2.0) for issue in issues),
        )
        title = self._get_document_title(source) or source.stem
        standard = AccessibilityStandard.WCAG_2_1_AA
        is_compliant = score >= 85.0 and not any(
            issue.severity in {IssueSeverity.CRITICAL, IssueSeverity.MAJOR}
            for issue in issues
        )
        return AccessibilityReport(
            document_path=source,
            standard=standard,
            is_compliant=is_compliant,
            score=round(score, 2),
            issues=issues,
            images=image_descriptions,
            headings=headings,
            tables=tables,
            reading_order=reading_order,
            language="en",
            title=title,
            metadata={
                "file_type": self._detect_file_type(source),
                "issue_count": len(issues),
                "image_count": len(image_descriptions),
                "table_count": len(tables),
                "heading_count": len(headings),
            },
        )

    def get_remediation_suggestions(self, path: str | Path) -> list[Suggestion]:
        """Return actionable remediation suggestions ordered by severity."""
        source = self._validate_path(path)
        suggestions: list[Suggestion] = []
        for issue in self._collect_all_issues(source):
            priority = issue.severity.value.upper()
            message = issue.remediation_hint or issue.description or issue.issue_type.value
            suggestions.append(
                Suggestion(
                    category=issue.issue_type.value,
                    priority=priority,
                    message=message,
                    auto_fixable=issue.auto_fixable,
                )
            )
        return suggestions

    def _collect_all_issues(self, path: Path) -> list[AccessibilityIssue]:
        """Aggregate all supported Office accessibility findings."""
        issues = []
        issues.extend(self.check_alt_text(path))
        issues.extend(self.check_heading_hierarchy(path))
        issues.extend(self.check_list_structure(path))
        issues.extend(self.check_table_headers(path))
        issues.extend(self.check_table_scope(path))
        issues.extend(self._check_document_title(path))
        reading_order = self.analyze_reading_order(path)
        if not reading_order.logical_order:
            issues.append(
                self._issue(
                    issue_type=IssueType.INCORRECT_READING_ORDER,
                    severity=IssueSeverity.MAJOR,
                    page=1,
                    description="Document reading order is not fully logical for assistive technology.",
                    wcag="1.3.2",
                    hint="Reorder body elements or slide shapes to match the intended reading sequence.",
                    auto_fixable=self._detect_file_type(path) == "pptx",
                )
            )
        for contrast_issue in self.check_color_contrast(path):
            issues.append(
                self._issue(
                    issue_type=IssueType.LOW_CONTRAST,
                    severity=IssueSeverity.MAJOR,
                    page=contrast_issue.page,
                    description=contrast_issue.description,
                    wcag="1.4.3",
                    hint="Increase contrast by darkening text or lightening the background.",
                    auto_fixable=False,
                )
            )
        return issues

    def _check_document_title(self, path: Path) -> list[AccessibilityIssue]:
        """Validate document metadata title presence."""
        title = self._get_document_title(path)
        if title:
            return []
        return [
            self._issue(
                issue_type=IssueType.MISSING_TITLE,
                severity=IssueSeverity.MINOR,
                page=1,
                description="Document title metadata is missing.",
                wcag="2.4.2",
                hint="Set the document title in Office core properties.",
                auto_fixable=True,
            )
        ]

    def _build_image_descriptions(self, path: Path) -> list[ImageDescription]:
        """Build Phase 3 image description models from Office images."""
        images: list[ImageDescription] = []
        for record in self._iter_images(path):
            images.append(
                ImageDescription(
                    image_id=record.image_id,
                    page=record.page,
                    bbox=(0.0, 0.0, 0.0, 0.0),
                    alt_text=record.alt_text or record.title or "",
                    confidence=1.0 if (record.alt_text or record.title) else 0.0,
                    description_type="extracted" if (record.alt_text or record.title) else "missing",
                )
            )
        return images

    def _build_heading_structures(self, path: Path) -> list[HeadingStructure]:
        """Build Phase 3 heading structures from Office documents."""
        file_type = self._detect_file_type(path)
        headings: list[HeadingStructure] = []
        if file_type == "docx" and DOCX_AVAILABLE:
            document = DocxDocument(path)
            for idx, paragraph in enumerate(document.paragraphs, start=1):
                style_name = (paragraph.style.name if paragraph.style else "").strip()
                if not style_name.lower().startswith("heading"):
                    continue
                level = self._parse_heading_level(style_name) or 1
                headings.append(
                    HeadingStructure(
                        level=level,
                        text=paragraph.text.strip(),
                        page=1,
                        bbox=(0.0, float(idx), 0.0, float(idx + 1)),
                    )
                )
        elif file_type == "pptx" and PPTX_AVAILABLE:
            presentation = Presentation(path)
            for slide_index, slide in enumerate(presentation.slides, start=1):
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    if not self._is_ppt_title(shape):
                        continue
                    text = shape.text.strip()
                    if not text:
                        continue
                    headings.append(
                        HeadingStructure(
                            level=1,
                            text=text,
                            page=slide_index,
                            bbox=self._shape_bbox(shape),
                        )
                    )
        return headings

    def _build_table_summaries(self, path: Path) -> list[TableSummary]:
        """Build Phase 3 table summaries from Office tables."""
        summaries: list[TableSummary] = []
        for page, table_id, matrix in self._iter_table_matrices(path):
            analyzed = self.table_engine.analyze_table(matrix, page=page)
            summaries.append(
                TableSummary(
                    table_id=table_id,
                    page=page,
                    rows=analyzed.rows,
                    cols=analyzed.cols,
                    has_headers=analyzed.has_header_row or analyzed.has_header_col,
                    header_row=1 if analyzed.has_header_row else None,
                    header_col=1 if analyzed.has_header_col else None,
                    summary=analyzed.summary,
                    caption=analyzed.caption,
                )
            )
        return summaries

    def _analyze_docx_reading_order(self, path: Path) -> ReadingOrderReport:
        """Analyze Word reading order from XML body sequence."""
        self._require_dependency(DOCX_AVAILABLE, "python-docx", "docx")
        document = DocxDocument(path)
        items: list[ReadingOrderItem] = []
        issues: list[str] = []
        order = 0
        saw_body_text = False
        for child in document.element.body.iterchildren():
            tag = self._local_name(child.tag)
            if tag == "p":
                paragraph = DocxParagraphElement(child, document)
                text = paragraph.text.strip()
                if not text:
                    continue
                style_name = (paragraph.style.name if paragraph.style else "").lower()
                element_type = "heading" if style_name.startswith("heading") else "text"
                if element_type == "heading" and saw_body_text:
                    issues.append("Heading appears after body content; verify document outline flow.")
                if element_type == "text":
                    saw_body_text = True
                items.append(
                    ReadingOrderItem(
                        element_id=f"docx-p-{order + 1}",
                        element_type=element_type,
                        page=1,
                        order=order,
                        bbox=(0.0, float(order), 0.0, float(order + 1)),
                        content_preview=text[:80],
                    )
                )
                order += 1
            elif tag == "tbl":
                table = DocxTableElement(child, document)
                preview = " | ".join(cell.text.strip() for cell in table.rows[0].cells[:4]) if table.rows else ""
                items.append(
                    ReadingOrderItem(
                        element_id=f"docx-tbl-{order + 1}",
                        element_type="table",
                        page=1,
                        order=order,
                        bbox=(0.0, float(order), 0.0, float(order + 1)),
                        content_preview=preview[:80],
                    )
                )
                order += 1
        return ReadingOrderReport(
            document_path=path,
            file_type="docx",
            logical_order=not issues,
            total_elements=len(items),
            items=items,
            issues=issues,
            metadata={"sections": len(document.sections)},
        )

    def _analyze_pptx_reading_order(self, path: Path) -> ReadingOrderReport:
        """Analyze PowerPoint reading order based on slide XML order and geometry."""
        self._require_dependency(PPTX_AVAILABLE, "python-pptx", "pptx")
        presentation = Presentation(path)
        items: list[ReadingOrderItem] = []
        issues: list[str] = []
        order = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            readable_shapes = [shape for shape in slide.shapes if self._shape_is_readable(shape)]
            expected_ids = [shape.shape_id for shape in sorted(readable_shapes, key=self._ppt_shape_sort_key)]
            current_ids = [shape.shape_id for shape in readable_shapes]
            if current_ids != expected_ids:
                issues.append(f"Slide {slide_index} shape order does not match expected reading order.")
            for shape in readable_shapes:
                items.append(
                    ReadingOrderItem(
                        element_id=f"slide-{slide_index}-shape-{shape.shape_id}",
                        element_type=self._ppt_shape_element_type(shape),
                        page=slide_index,
                        order=order,
                        bbox=self._shape_bbox(shape),
                        content_preview=self._shape_text_preview(shape),
                    )
                )
                order += 1
        return ReadingOrderReport(
            document_path=path,
            file_type="pptx",
            logical_order=not issues,
            total_elements=len(items),
            items=items,
            issues=issues,
            metadata={"slides": len(presentation.slides)},
        )

    def _analyze_xlsx_reading_order(self, path: Path) -> ReadingOrderReport:
        """Analyze workbook/sheet order for spreadsheet accessibility."""
        self._require_dependency(OPENPYXL_AVAILABLE, "openpyxl", "xlsx")
        workbook = load_workbook(path, read_only=True, data_only=True)
        items: list[ReadingOrderItem] = []
        issues: list[str] = []
        order = 0
        visible_sheets = [ws for ws in workbook.worksheets if ws.sheet_state == "visible"]
        if not visible_sheets:
            issues.append("Workbook contains no visible worksheets.")
        for index, sheet in enumerate(workbook.worksheets, start=1):
            if sheet.sheet_state != "visible" and index == 1:
                issues.append("The first worksheet is hidden; assistive navigation may be confusing.")
            items.append(
                ReadingOrderItem(
                    element_id=f"sheet-{index}",
                    element_type="worksheet",
                    page=index,
                    order=order,
                    bbox=(0.0, float(index), 0.0, float(index + 1)),
                    content_preview=sheet.title,
                )
            )
            order += 1
        return ReadingOrderReport(
            document_path=path,
            file_type="xlsx",
            logical_order=not issues,
            total_elements=len(items),
            items=items,
            issues=issues,
            metadata={"worksheets": len(workbook.worksheets)},
        )

    def _check_docx_headings(self, path: Path) -> list[AccessibilityIssue]:
        """Check Word heading hierarchy using paragraph styles."""
        self._require_dependency(DOCX_AVAILABLE, "python-docx", "docx")
        issues: list[AccessibilityIssue] = []
        document = DocxDocument(path)
        headings: list[tuple[int, str]] = []
        for paragraph in document.paragraphs:
            style_name = (paragraph.style.name if paragraph.style else "").strip()
            level = self._parse_heading_level(style_name)
            if level is None:
                continue
            headings.append((level, paragraph.text.strip()))
        previous_level = 0
        if headings and headings[0][0] != 1:
            issues.append(
                self._issue(
                    issue_type=IssueType.MISSING_HEADINGS,
                    severity=IssueSeverity.MINOR,
                    page=1,
                    description="The first heading is not Heading 1.",
                    wcag="2.4.6",
                    hint="Start the heading hierarchy with Heading 1.",
                    auto_fixable=False,
                )
            )
        for level, text in headings:
            if not text:
                issues.append(
                    self._issue(
                        issue_type=IssueType.MISSING_HEADINGS,
                        severity=IssueSeverity.MAJOR,
                        page=1,
                        description=f"Heading {level} is empty.",
                        wcag="1.3.1",
                        hint="Provide meaningful text for every heading.",
                        auto_fixable=False,
                    )
                )
            if previous_level and level > previous_level + 1:
                issues.append(
                    self._issue(
                        issue_type=IssueType.MISSING_HEADINGS,
                        severity=IssueSeverity.MINOR,
                        page=1,
                        description=f"Heading level skips from H{previous_level} to H{level}.",
                        wcag="1.3.1",
                        hint="Avoid skipping heading levels.",
                        auto_fixable=False,
                    )
                )
            previous_level = level
        return issues

    def _check_pptx_headings(self, path: Path) -> list[AccessibilityIssue]:
        """Check slide title structure as a presentation heading hierarchy."""
        self._require_dependency(PPTX_AVAILABLE, "python-pptx", "pptx")
        issues: list[AccessibilityIssue] = []
        presentation = Presentation(path)
        for slide_index, slide in enumerate(presentation.slides, start=1):
            titles = [
                shape.text.strip()
                for shape in slide.shapes
                if self._is_ppt_title(shape) and getattr(shape, "text", "").strip()
            ]
            if not titles:
                issues.append(
                    self._issue(
                        issue_type=IssueType.MISSING_HEADINGS,
                        severity=IssueSeverity.MAJOR,
                        page=slide_index,
                        description=f"Slide {slide_index} is missing a title placeholder.",
                        wcag="2.4.6",
                        hint="Add a descriptive slide title so screen readers can announce context.",
                        auto_fixable=False,
                    )
                )
        return issues

    def _check_docx_lists(self, path: Path) -> list[AccessibilityIssue]:
        """Detect Word paragraphs that visually look like lists but are not semantic lists."""
        self._require_dependency(DOCX_AVAILABLE, "python-docx", "docx")
        issues: list[AccessibilityIssue] = []
        document = DocxDocument(path)
        for index, paragraph in enumerate(document.paragraphs, start=1):
            text = paragraph.text.strip()
            if not text or not LIST_MARKER_RE.match(text):
                continue
            style_name = (paragraph.style.name if paragraph.style else "").lower()
            if "list" in style_name:
                continue
            num_pr = paragraph._p.pPr.numPr if paragraph._p.pPr is not None else None
            if num_pr is not None:
                continue
            issues.append(
                self._issue(
                    issue_type=IssueType.UNTAGGED_CONTENT,
                    severity=IssueSeverity.MINOR,
                    page=1,
                    description=f"Paragraph {index} appears to be a manual list item.",
                    wcag="1.3.1",
                    hint="Use Word list formatting instead of typing bullets or numbers manually.",
                    auto_fixable=False,
                )
            )
        return issues

    def _check_pptx_lists(self, path: Path) -> list[AccessibilityIssue]:
        """Detect PowerPoint text that looks like a manual list."""
        self._require_dependency(PPTX_AVAILABLE, "python-pptx", "pptx")
        issues: list[AccessibilityIssue] = []
        presentation = Presentation(path)
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if LIST_MARKER_RE.match(text) and paragraph.level == 0:
                        issues.append(
                            self._issue(
                                issue_type=IssueType.UNTAGGED_CONTENT,
                                severity=IssueSeverity.MINOR,
                                page=slide_index,
                                description=(
                                    f"Slide {slide_index} contains text that appears to be a "
                                    "manual bullet or numbered list."
                                ),
                                wcag="1.3.1",
                                hint="Use real slide bullets and levels instead of manual list markers.",
                                auto_fixable=False,
                            )
                        )
        return issues

    def _iter_table_matrices(self, path: Path) -> Iterable[tuple[int, str, list[list[str]]]]:
        """Yield table matrices across supported Office formats."""
        file_type = self._detect_file_type(path)
        if file_type == "docx":
            self._require_dependency(DOCX_AVAILABLE, "python-docx", "docx")
            document = DocxDocument(path)
            for index, table in enumerate(document.tables, start=1):
                yield 1, f"docx-table-{index}", [
                    [cell.text.strip() for cell in row.cells] for row in table.rows
                ]
            return
        if file_type == "pptx":
            self._require_dependency(PPTX_AVAILABLE, "python-pptx", "pptx")
            presentation = Presentation(path)
            for slide_index, slide in enumerate(presentation.slides, start=1):
                table_index = 0
                for shape in slide.shapes:
                    if not getattr(shape, "has_table", False):
                        continue
                    table_index += 1
                    yield slide_index, f"pptx-slide-{slide_index}-table-{table_index}", [
                        [cell.text.strip() for cell in row.cells] for row in shape.table.rows
                    ]
            return
        if file_type == "xlsx":
            self._require_dependency(OPENPYXL_AVAILABLE, "openpyxl", "xlsx")
            workbook = load_workbook(path, read_only=False, data_only=True)
            for sheet in workbook.worksheets:
                if not getattr(sheet, "tables", None):
                    continue
                for display_name, table in sheet.tables.items():
                    min_col, min_row, max_col, max_row = range_boundaries(table.ref)
                    matrix = []
                    for row in sheet.iter_rows(
                        min_row=min_row,
                        max_row=max_row,
                        min_col=min_col,
                        max_col=max_col,
                    ):
                        matrix.append([self._cell_to_text(cell.value) for cell in row])
                    yield workbook.worksheets.index(sheet) + 1, f"xlsx:{sheet.title}:{display_name}", matrix

    def _add_docx_table_headers(
        self,
        path: Path,
        table_id: str,
        headers: list[str],
        output_path: Path,
    ) -> Path:
        """Write first-row headers into a Word table."""
        self._require_dependency(DOCX_AVAILABLE, "python-docx", "docx")
        match = re.fullmatch(r"docx-table-(\d+)", table_id)
        if not match:
            raise DocumentValidationError("Invalid DOCX table identifier", details=table_id)
        table_index = int(match.group(1)) - 1
        document = DocxDocument(path)
        if table_index < 0 or table_index >= len(document.tables):
            raise DocumentValidationError("DOCX table not found", details=table_id)
        table = document.tables[table_index]
        if not table.rows:
            raise DocumentValidationError("DOCX table is empty", details=table_id)
        header_cells = table.rows[0].cells
        for index, text in enumerate(headers[: len(header_cells)]):
            header_cells[index].text = text
        output_path.parent.mkdir(parents=True, exist_ok=True)
        document.save(output_path)
        return output_path

    def _add_pptx_table_headers(
        self,
        path: Path,
        table_id: str,
        headers: list[str],
        output_path: Path,
    ) -> Path:
        """Write first-row headers into a PowerPoint table."""
        self._require_dependency(PPTX_AVAILABLE, "python-pptx", "pptx")
        match = re.fullmatch(r"pptx-slide-(\d+)-table-(\d+)", table_id)
        if not match:
            raise DocumentValidationError("Invalid PPTX table identifier", details=table_id)
        slide_index = int(match.group(1)) - 1
        wanted_table = int(match.group(2))
        presentation = Presentation(path)
        if slide_index < 0 or slide_index >= len(presentation.slides):
            raise DocumentValidationError("PPTX slide not found", details=table_id)
        slide = presentation.slides[slide_index]
        current = 0
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            current += 1
            if current != wanted_table:
                continue
            for index, text in enumerate(headers[: len(shape.table.rows[0].cells)]):
                shape.table.cell(0, index).text = text
            output_path.parent.mkdir(parents=True, exist_ok=True)
            presentation.save(output_path)
            return output_path
        raise DocumentValidationError("PPTX table not found", details=table_id)

    def _add_xlsx_table_headers(
        self,
        path: Path,
        table_id: str,
        headers: list[str],
        output_path: Path,
    ) -> Path:
        """Write first-row headers into an Excel table region."""
        self._require_dependency(OPENPYXL_AVAILABLE, "openpyxl", "xlsx")
        match = re.fullmatch(r"xlsx:([^:]+):(.+)", table_id)
        if not match:
            raise DocumentValidationError("Invalid XLSX table identifier", details=table_id)
        sheet_name, display_name = match.group(1), match.group(2)
        workbook = load_workbook(path)
        if sheet_name not in workbook.sheetnames:
            raise DocumentValidationError("Worksheet not found", details=sheet_name)
        sheet = workbook[sheet_name]
        table = sheet.tables.get(display_name)
        if table is None:
            raise DocumentValidationError("Excel table not found", details=table_id)
        min_col, min_row, max_col, _ = range_boundaries(table.ref)
        for index, text in enumerate(headers, start=min_col):
            if index > max_col:
                break
            sheet.cell(row=min_row, column=index, value=text)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        workbook.save(output_path)
        return output_path

    def _collect_contrast_candidates(self, path: Path) -> list[_ContrastCandidate]:
        """Collect text/background pairs for Office contrast analysis."""
        file_type = self._detect_file_type(path)
        if file_type == "pptx":
            return self._collect_pptx_contrast_candidates(path)
        if file_type == "docx":
            return self._collect_docx_contrast_candidates(path)
        if file_type == "xlsx":
            return self._collect_xlsx_contrast_candidates(path)
        return []

    def _collect_pptx_contrast_candidates(self, path: Path) -> list[_ContrastCandidate]:
        """Collect contrast candidates from PowerPoint text shapes."""
        self._require_dependency(PPTX_AVAILABLE, "python-pptx", "pptx")
        candidates: list[_ContrastCandidate] = []
        presentation = Presentation(path)
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                background = self._ppt_shape_background(shape, slide) or "#FFFFFF"
                for paragraph in shape.text_frame.paragraphs:
                    size = TextSize.NORMAL
                    for run in paragraph.runs:
                        foreground = self._ppt_font_color(run.font)
                        if not foreground or not run.text.strip():
                            continue
                        if run.font.size and getattr(run.font.size, "pt", 0) >= 18:
                            size = TextSize.LARGE
                        candidates.append(
                            _ContrastCandidate(
                                element_id=f"slide-{slide_index}-shape-{shape.shape_id}",
                                page=slide_index,
                                foreground=foreground,
                                background=background,
                                text_size=size,
                                description="Slide text",
                            )
                        )
        return candidates

    def _collect_docx_contrast_candidates(self, path: Path) -> list[_ContrastCandidate]:
        """Collect contrast candidates from Word runs with explicit formatting."""
        self._require_dependency(DOCX_AVAILABLE, "python-docx", "docx")
        candidates: list[_ContrastCandidate] = []
        document = DocxDocument(path)
        for paragraph_index, paragraph in enumerate(document.paragraphs, start=1):
            for run_index, run in enumerate(paragraph.runs, start=1):
                foreground = self._docx_run_color(run)
                if not foreground or not run.text.strip():
                    continue
                size = TextSize.LARGE if run.font.size and run.font.size.pt >= 18 else TextSize.NORMAL
                candidates.append(
                    _ContrastCandidate(
                        element_id=f"paragraph-{paragraph_index}-run-{run_index}",
                        page=1,
                        foreground=foreground,
                        background="#FFFFFF",
                        text_size=size,
                        description="Word text",
                    )
                )
        return candidates

    def _collect_xlsx_contrast_candidates(self, path: Path) -> list[_ContrastCandidate]:
        """Collect contrast candidates from Excel cells with explicit colors."""
        self._require_dependency(OPENPYXL_AVAILABLE, "openpyxl", "xlsx")
        candidates: list[_ContrastCandidate] = []
        workbook = load_workbook(path, data_only=True)
        for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
            for row in sheet.iter_rows():
                for cell in row:
                    if not isinstance(cell.value, str) or not cell.value.strip():
                        continue
                    foreground = self._openpyxl_color_to_hex(getattr(cell.font, "color", None))
                    if not foreground:
                        continue
                    background = self._openpyxl_fill_to_hex(getattr(cell, "fill", None)) or "#FFFFFF"
                    candidates.append(
                        _ContrastCandidate(
                            element_id=f"{sheet.title}!{cell.coordinate}",
                            page=sheet_index,
                            foreground=foreground,
                            background=background,
                            text_size=TextSize.NORMAL,
                            description="Worksheet cell text",
                        )
                    )
        return candidates

    def _fix_pptx_reading_order(self, path: Path, output_path: Path) -> Path:
        """Reorder slide shapes in XML order to better match reading order."""
        self._require_dependency(PPTX_AVAILABLE, "python-pptx", "pptx")
        presentation = Presentation(path)
        for slide in presentation.slides:
            sp_tree = slide.shapes._spTree
            shape_elements = list(sp_tree)[2:]
            if not shape_elements:
                continue
            ordered_shapes = sorted(
                slide.shapes,
                key=self._ppt_shape_sort_key,
            )
            desired = [shape.element for shape in ordered_shapes if shape.element in shape_elements]
            if not desired:
                continue
            for element in shape_elements:
                sp_tree.remove(element)
            for element in desired:
                sp_tree.append(element)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        return output_path

    def _iter_images(self, path: Path) -> list[_OfficeImageRecord]:
        """Extract OOXML images and their alt-text metadata."""
        file_type = self._detect_file_type(path)
        if file_type not in OOXML_EXTENSIONS:
            raise UnsupportedOfficeFormatError(file_type)
        with zipfile.ZipFile(path, "r") as archive:
            if file_type == "docx":
                return self._extract_docx_images(archive)
            if file_type == "pptx":
                return self._extract_pptx_images(archive)
            return self._extract_xlsx_images(archive)

    def _extract_docx_images(self, archive: zipfile.ZipFile) -> list[_OfficeImageRecord]:
        """Extract image metadata from WordprocessingML parts."""
        records: list[_OfficeImageRecord] = []
        part_names = [
            name
            for name in archive.namelist()
            if name.startswith("word/") and name.endswith(".xml") and "/_rels/" not in name
        ]
        for part_name in part_names:
            rels = self._load_relationships(archive, part_name)
            root = ET.fromstring(archive.read(part_name))
            for docpr in root.findall(".//wp:docPr", NS):
                rel_id = None
                parent = self._find_parent_with_blip(root, docpr)
                if parent is not None:
                    blip = parent.find(".//a:blip", NS)
                    rel_id = blip.get(f"{{{NS['r']}}}embed") if blip is not None else None
                if not rel_id or rel_id not in rels:
                    continue
                media_target = rels[rel_id]
                if media_target not in archive.namelist():
                    continue
                xml_id = docpr.get("id", "")
                display_name = docpr.get("name", media_target.rsplit("/", 1)[-1])
                page = 1
                records.append(
                    _OfficeImageRecord(
                        image_id=f"{part_name}:{xml_id}",
                        xml_id=xml_id,
                        display_name=display_name,
                        part_name=part_name,
                        media_target=media_target,
                        location=part_name,
                        page=page,
                        alt_text=docpr.get("descr", ""),
                        title=docpr.get("title", ""),
                        data=archive.read(media_target),
                    )
                )
        return records

    def _extract_pptx_images(self, archive: zipfile.ZipFile) -> list[_OfficeImageRecord]:
        """Extract image metadata from PresentationML slides."""
        records: list[_OfficeImageRecord] = []
        slide_parts = sorted(
            name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        for index, part_name in enumerate(slide_parts, start=1):
            rels = self._load_relationships(archive, part_name)
            root = ET.fromstring(archive.read(part_name))
            for pic in root.findall(".//p:pic", NS):
                c_nv_pr = pic.find(".//p:cNvPr", NS)
                blip = pic.find(".//a:blip", NS)
                if c_nv_pr is None or blip is None:
                    continue
                rel_id = blip.get(f"{{{NS['r']}}}embed")
                if not rel_id or rel_id not in rels:
                    continue
                media_target = rels[rel_id]
                if media_target not in archive.namelist():
                    continue
                xml_id = c_nv_pr.get("id", "")
                display_name = c_nv_pr.get("name", media_target.rsplit("/", 1)[-1])
                records.append(
                    _OfficeImageRecord(
                        image_id=f"{part_name}:{xml_id}",
                        xml_id=xml_id,
                        display_name=display_name,
                        part_name=part_name,
                        media_target=media_target,
                        location=f"slide {index}",
                        page=index,
                        alt_text=c_nv_pr.get("descr", ""),
                        title=c_nv_pr.get("title", ""),
                        data=archive.read(media_target),
                    )
                )
        return records

    def _extract_xlsx_images(self, archive: zipfile.ZipFile) -> list[_OfficeImageRecord]:
        """Extract image metadata from SpreadsheetML drawings."""
        records: list[_OfficeImageRecord] = []
        drawing_to_sheet = self._xlsx_drawing_sheet_map(archive)
        drawing_parts = sorted(
            name for name in archive.namelist() if name.startswith("xl/drawings/drawing") and name.endswith(".xml")
        )
        for part_name in drawing_parts:
            rels = self._load_relationships(archive, part_name)
            root = ET.fromstring(archive.read(part_name))
            for picture in root.findall(".//xdr:pic", NS):
                c_nv_pr = picture.find(".//xdr:cNvPr", NS)
                blip = picture.find(".//a:blip", NS)
                if c_nv_pr is None or blip is None:
                    continue
                rel_id = blip.get(f"{{{NS['r']}}}embed")
                if not rel_id or rel_id not in rels:
                    continue
                media_target = rels[rel_id]
                if media_target not in archive.namelist():
                    continue
                xml_id = c_nv_pr.get("id", "")
                sheet_name = drawing_to_sheet.get(part_name, part_name)
                records.append(
                    _OfficeImageRecord(
                        image_id=f"{part_name}:{xml_id}",
                        xml_id=xml_id,
                        display_name=c_nv_pr.get("name", media_target.rsplit("/", 1)[-1]),
                        part_name=part_name,
                        media_target=media_target,
                        location=sheet_name,
                        page=1,
                        alt_text=c_nv_pr.get("descr", ""),
                        title=c_nv_pr.get("title", ""),
                        data=archive.read(media_target),
                    )
                )
        return records

    def _apply_alt_text_updates(
        self,
        source: Path,
        output_path: Path,
        updates: dict[str, dict[str, str]],
    ) -> Path:
        """Write alt-text updates back into an OOXML archive."""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(source, "r") as src_zip, zipfile.ZipFile(
            output_path,
            "w",
            compression=zipfile.ZIP_DEFLATED,
        ) as dst_zip:
            for info in src_zip.infolist():
                data = src_zip.read(info.filename)
                if info.filename in updates:
                    root = ET.fromstring(data)
                    changed = False
                    for element in root.iter():
                        local = self._local_name(element.tag)
                        if local not in {"docPr", "cNvPr"}:
                            continue
                        xml_id = element.get("id", "")
                        if xml_id in updates[info.filename]:
                            value = updates[info.filename][xml_id]
                            element.set("descr", value)
                            element.set("title", value[:255])
                            changed = True
                    if changed:
                        data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
                dst_zip.writestr(info, data)
        return output_path

    def _generate_image_alt_text(self, record: _OfficeImageRecord) -> str:
        """Generate alt text using the shared Phase 3 alt-text service."""
        if PILImage is None:
            return f"Image: {record.display_name.replace('_', ' ').replace('-', ' ')}".strip()
        try:
            with PILImage.open(io.BytesIO(record.data)) as image:
                result = self.alt_text_generator.generate(image=image.convert("RGB"))
                return result.text.strip() or f"Image: {record.display_name}"
        except Exception as exc:  # pragma: no cover - defensive fallback
            logger.warning("Alt-text generation failed for %s: %s", record.image_id, exc)
            return f"Image: {record.display_name.replace('_', ' ').replace('-', ' ')}".strip()

    def _get_document_title(self, path: Path) -> str:
        """Extract Office core title metadata using the package core properties."""
        try:
            with zipfile.ZipFile(path, "r") as archive:
                if "docProps/core.xml" not in archive.namelist():
                    return ""
                root = ET.fromstring(archive.read("docProps/core.xml"))
                title = root.findtext("dc:title", default="", namespaces=NS)
                return (title or "").strip()
        except Exception:  # pragma: no cover - malformed package fallback
            return ""

    def _xlsx_drawing_sheet_map(self, archive: zipfile.ZipFile) -> dict[str, str]:
        """Map spreadsheet drawing parts to worksheet titles."""
        mapping: dict[str, str] = {}
        if "xl/workbook.xml" not in archive.namelist():
            return mapping
        workbook_root = ET.fromstring(archive.read("xl/workbook.xml"))
        workbook_rels = self._load_relationships(archive, "xl/workbook.xml")
        sheet_parts: dict[str, str] = {}
        for sheet in workbook_root.findall(".//x:sheets/x:sheet", NS):
            rel_id = sheet.get(f"{{{NS['r']}}}id")
            name = sheet.get("name", "")
            if rel_id and rel_id in workbook_rels:
                sheet_parts[workbook_rels[rel_id]] = name
        for sheet_part, sheet_name in sheet_parts.items():
            rels = self._load_relationships(archive, sheet_part)
            root = ET.fromstring(archive.read(sheet_part))
            drawing = root.find(".//x:drawing", NS)
            if drawing is None:
                continue
            rel_id = drawing.get(f"{{{NS['r']}}}id")
            if rel_id and rel_id in rels:
                mapping[rels[rel_id]] = sheet_name
        return mapping

    def _load_relationships(self, archive: zipfile.ZipFile, part_name: str) -> dict[str, str]:
        """Load OOXML relationships for a part."""
        rels_path = self._relationships_path(part_name)
        if rels_path not in archive.namelist():
            return {}
        root = ET.fromstring(archive.read(rels_path))
        rels: dict[str, str] = {}
        for rel in root.findall(".//pr:Relationship", NS):
            rel_id = rel.get("Id")
            target = rel.get("Target")
            if not rel_id or not target:
                continue
            rels[rel_id] = self._resolve_part_target(part_name, target)
        return rels

    def _relationships_path(self, part_name: str) -> str:
        """Return the OOXML relationships path for a part."""
        directory, filename = posixpath.split(part_name)
        return posixpath.join(directory, "_rels", f"{filename}.rels")

    def _resolve_part_target(self, part_name: str, target: str) -> str:
        """Resolve an OOXML relationship target relative to its source part."""
        if target.startswith("/"):
            return target.lstrip("/")
        base_dir = posixpath.dirname(part_name)
        return posixpath.normpath(posixpath.join(base_dir, target))

    def _guess_mime(self, resource: str) -> str:
        """Guess MIME type for an embedded image resource."""
        mime_type, _ = mimetypes.guess_type(resource)
        return mime_type or "application/octet-stream"

    def _copy_document(self, source: Path, output_path: str | Path) -> Path:
        """Copy Office document to an output path."""
        target = Path(output_path)
        target.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source, target)
        return target

    def _validate_path(self, path: str | Path) -> Path:
        """Validate the source Office document path."""
        source = Path(path)
        if not source.exists():
            raise FileNotFoundError(f"Office document not found: {source}")
        if source.is_dir():
            raise DocumentValidationError("Expected Office document file, got directory")
        extension = source.suffix.lower().lstrip(".")
        if extension not in OOXML_EXTENSIONS:
            raise UnsupportedOfficeFormatError(extension)
        return source

    def _detect_file_type(self, path: Path) -> str:
        """Return lowercase Office file extension without the dot."""
        return path.suffix.lower().lstrip(".")

    def _issue(
        self,
        *,
        issue_type: IssueType,
        severity: IssueSeverity,
        page: int,
        description: str,
        wcag: str,
        hint: str,
        auto_fixable: bool,
    ) -> AccessibilityIssue:
        """Create a normalized accessibility issue."""
        return AccessibilityIssue(
            issue_type=issue_type,
            severity=severity,
            page=page,
            description=description,
            wcag_criterion=wcag,
            remediation_hint=hint,
            auto_fixable=auto_fixable,
        )

    def _parse_heading_level(self, style_name: str) -> int | None:
        """Parse a heading level from an Office style name."""
        match = re.search(r"heading\s*(\d+)", style_name, re.IGNORECASE)
        return int(match.group(1)) if match else None

    def _require_dependency(self, available: bool, package: str, file_type: str) -> None:
        """Raise a readable error for missing optional Office dependencies."""
        if not available:
            raise DocumentValidationError(
                f"{package} is required for {file_type} accessibility processing"
            )

    def _cell_to_text(self, value: Any) -> str:
        """Normalize spreadsheet cell values into strings."""
        return "" if value is None else str(value).strip()

    def _local_name(self, tag: str) -> str:
        """Extract local name from an XML tag."""
        return tag.rsplit("}", 1)[-1]

    def _find_parent_with_blip(self, root: ET.Element, target: ET.Element) -> Optional[ET.Element]:
        """Find a nearby drawing container for a given docPr element."""
        for candidate in root.iter():
            if target not in list(candidate):
                continue
            if candidate.find(".//a:blip", NS) is not None:
                return candidate
        for candidate in root.iter():
            if target in list(candidate.iter()) and candidate.find(".//a:blip", NS) is not None:
                return candidate
        return None

    def _shape_is_readable(self, shape: Any) -> bool:
        """Return whether a PowerPoint shape participates in reading order."""
        return (
            getattr(shape, "has_text_frame", False)
            or getattr(shape, "has_table", False)
            or getattr(shape, "shape_type", None) == getattr(MSO_SHAPE_TYPE, "PICTURE", None)
        )

    def _ppt_shape_sort_key(self, shape: Any) -> tuple[int, int, int, int]:
        """Sort PowerPoint shapes by semantic role and geometry."""
        priority = 0 if self._is_ppt_title(shape) else 1
        top = int(getattr(shape, "top", 0) or 0)
        left = int(getattr(shape, "left", 0) or 0)
        shape_id = int(getattr(shape, "shape_id", 0) or 0)
        return (priority, top, left, shape_id)

    def _is_ppt_title(self, shape: Any) -> bool:
        """Return whether a PowerPoint shape is a title placeholder."""
        try:
            if not getattr(shape, "is_placeholder", False):
                return False
            placeholder = shape.placeholder_format
            return placeholder.type in {
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            }
        except Exception:
            return False

    def _ppt_shape_element_type(self, shape: Any) -> str:
        """Infer Phase 3 element type for a PowerPoint shape."""
        if self._is_ppt_title(shape):
            return "heading"
        if getattr(shape, "has_table", False):
            return "table"
        if getattr(shape, "shape_type", None) == getattr(MSO_SHAPE_TYPE, "PICTURE", None):
            return "image"
        return "text"

    def _shape_text_preview(self, shape: Any) -> str:
        """Generate a short PowerPoint shape preview string."""
        if getattr(shape, "has_text_frame", False):
            return getattr(shape, "text", "").strip()[:80]
        if getattr(shape, "has_table", False):
            first_row = shape.table.rows[0].cells if shape.table.rows else []
            return " | ".join(cell.text.strip() for cell in first_row[:4])[:80]
        return getattr(shape, "name", "")[:80]

    def _shape_bbox(self, shape: Any) -> tuple[float, float, float, float]:
        """Convert a PowerPoint shape geometry into a bbox tuple."""
        left = float(getattr(shape, "left", 0) or 0)
        top = float(getattr(shape, "top", 0) or 0)
        width = float(getattr(shape, "width", 0) or 0)
        height = float(getattr(shape, "height", 0) or 0)
        return (left, top, left + width, top + height)

    def _ppt_font_color(self, font: Any) -> str:
        """Extract PowerPoint font RGB as hex if available."""
        try:
            rgb = getattr(font.color, "rgb", None)
            if rgb is None:
                return ""
            return f"#{str(rgb)[-6:]}"
        except Exception:
            return ""

    def _ppt_shape_background(self, shape: Any, slide: Any) -> str:
        """Extract a PowerPoint shape background color when available."""
        try:
            fill = shape.fill
            fore_color = getattr(fill, "fore_color", None)
            rgb = getattr(fore_color, "rgb", None)
            if rgb is not None:
                return f"#{str(rgb)[-6:]}"
        except Exception:
            pass
        try:
            bg = slide.background.fill
            rgb = getattr(getattr(bg, "fore_color", None), "rgb", None)
            if rgb is not None:
                return f"#{str(rgb)[-6:]}"
        except Exception:
            pass
        return "#FFFFFF"

    def _docx_run_color(self, run: Any) -> str:
        """Extract Word run color as hex."""
        try:
            rgb = getattr(run.font.color, "rgb", None)
            if rgb is None:
                return ""
            return f"#{str(rgb)[-6:]}"
        except Exception:
            return ""

    def _openpyxl_color_to_hex(self, color: Any) -> str:
        """Convert an openpyxl Color object into an RGB hex string."""
        rgb = getattr(color, "rgb", None)
        if not rgb:
            return ""
        rgb = str(rgb)
        if len(rgb) == 8:
            rgb = rgb[2:]
        return f"#{rgb[-6:]}" if HEX_COLOR_RE.match(rgb) else ""

    def _openpyxl_fill_to_hex(self, fill: Any) -> str:
        """Convert spreadsheet fill to hex where possible."""
        if fill is None:
            return ""
        if not isinstance(fill, PatternFill):
            return ""
        return self._openpyxl_color_to_hex(fill.fgColor)


__all__ = [
    "ContrastFix",
    "OfficeAccessibilityProcessor",
    "ReadingOrderReport",
    "Suggestion",
    "WCAGReport",
]
