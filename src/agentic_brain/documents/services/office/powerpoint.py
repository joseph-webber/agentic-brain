# SPDX-License-Identifier: Apache-2.0

"""PowerPoint processing built on top of python-pptx.

This module provides a high-level, structured PowerPoint processor capable of
reading, writing, and inspecting presentation content while mapping extracted
content into the office document models used across agentic-brain.

The implementation is intentionally defensive:

* python-pptx is treated as an optional dependency and only required at runtime
  when PowerPoint operations are requested.
* XML-backed features not exposed directly by python-pptx, such as transition
  and animation discovery, are implemented using Open XML inspection.
* Thumbnail generation is best-effort. python-pptx does not render slides, so
  textual preview thumbnails are generated with Pillow when available.
"""

from __future__ import annotations

import xml.etree.ElementTree as ET
from collections.abc import Iterable, Mapping, Sequence
from io import BytesIO
from pathlib import Path
from typing import Any, BinaryIO, cast
from zipfile import ZipFile

try:
    from PIL import Image as PILImage
    from PIL import ImageDraw, ImageFont

    PIL_AVAILABLE = True
except ImportError:  # pragma: no cover - optional dependency
    PIL_AVAILABLE = False
    PILImage = None
    ImageDraw = None
    ImageFont = None

try:
    from pptx import Presentation as PptxPresentationFactory
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu, Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover - optional dependency
    PptxPresentationFactory = None
    ChartData = None
    XL_CHART_TYPE = None
    XL_LEGEND_POSITION = None
    MSO_SHAPE_TYPE = None
    Emu = None
    Inches = None
    Pt = None
    PPTX_AVAILABLE = False

from .exceptions import DocumentValidationError, UnsupportedOfficeFormatError
from .models import (
    Chart,
    DocumentContent,
    DocumentStyle,
    Image,
    Metadata,
    OfficeFormat,
    Paragraph,
    Shape,
    Slide,
    Table,
    TableCell,
    TextRun,
)

EMU_PER_POINT = 12_700
EMU_PER_INCH = 914_400
DEFAULT_THUMBNAIL_SIZE = (400, 225)

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
    "dc": "http://purl.org/dc/elements/1.1/",
    "ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "vt": "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes",
}

MEDIA_REL_TYPES = {
    "video": {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video",
        "http://schemas.microsoft.com/office/2007/relationships/media",
    },
    "audio": {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio",
        "http://schemas.microsoft.com/office/2007/relationships/media",
    },
    "media": {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/media",
        "http://schemas.microsoft.com/office/2007/relationships/media",
    },
    "ole": {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package",
    },
}


def _local_name(tag: str) -> str:
    """Return XML local-name for a tag like ``{namespace}name``."""

    if "}" in tag:
        return tag.rsplit("}", 1)[-1]
    return tag


def _pt(value: Any) -> float | None:
    """Convert a PPTX length value expressed in EMU to points."""

    if value is None:
        return None
    try:
        return float(value) / EMU_PER_POINT
    except (TypeError, ValueError):
        return None


def _inches_to_emu(value: float) -> int:
    """Convert inches to EMU, independent of pptx.util availability."""

    return int(round(value * EMU_PER_INCH))


def _is_truthy(value: Any) -> bool:
    """Normalize truth-like XML attribute values."""

    if isinstance(value, bool):
        return value
    if value is None:
        return False
    return str(value).strip().lower() in {"1", "true", "yes", "on"}


class PowerPointProcessor:
    """Comprehensive processor for PowerPoint presentations.

    The processor supports both reading and writing:

    * Parse an existing ``.pptx`` file into :class:`DocumentContent`
    * Extract slides, notes, images, tables, charts, and metadata
    * Discover hyperlinks, media, transitions, and animation markers
    * Create or modify presentations using python-pptx
    """

    def __init__(self, path: str | Path | None = None) -> None:
        self._source_path: Path | None = None
        self._presentation: Any | None = None
        self._parsed_content: DocumentContent | None = None
        self._hyperlinks_cache: list[dict[str, Any]] | None = None
        self._animations_cache: list[dict[str, Any]] | None = None
        self._media_cache: list[dict[str, Any]] | None = None
        self._thumbnails_cache: list[Image] | None = None

        if path is not None:
            self._source_path = Path(path)

    # ------------------------------------------------------------------
    # Public read API
    # ------------------------------------------------------------------
    def parse(self, path: str | Path | None = None) -> DocumentContent:
        """Parse a PowerPoint file and return normalized document content."""

        if path is not None:
            self._source_path = Path(path)

        self._ensure_pptx_available()
        resolved_path = self._require_source_path()
        self._validate_path(resolved_path)

        self._presentation = PptxPresentationFactory(str(resolved_path))
        self._invalidate_caches()

        metadata = self.extract_metadata()
        slides = self.extract_all_slides()
        content = DocumentContent(
            format=OfficeFormat.PPTX,
            slides=slides,
            images=[image for slide in slides for image in slide.images],
            shapes=[shape for slide in slides for shape in slide.shapes],
            tables=[table for slide in slides for table in slide.tables],
            charts=[chart for slide in slides for chart in slide.charts],
            metadata=metadata,
            document_properties={
                "slide_count": len(slides),
                "slide_width_pt": _pt(getattr(self._presentation, "slide_width", None)),
                "slide_height_pt": _pt(
                    getattr(self._presentation, "slide_height", None)
                ),
                "hyperlink_count": len(self.extract_hyperlinks()),
                "animation_count": len(self.detect_animations()),
                "media_count": len(self.detect_embedded_media()),
            },
        )
        self._parsed_content = content
        return content

    def get_slide_count(self) -> int:
        """Return the number of slides in the loaded presentation."""

        presentation = self._require_presentation()
        return len(presentation.slides)

    def extract_slide(self, index: int) -> Slide:
        """Extract a single slide by zero-based index."""

        slides = self.extract_all_slides()
        if index < 0 or index >= len(slides):
            raise IndexError(f"Slide index out of range: {index}")
        return slides[index]

    def extract_all_slides(self) -> list[Slide]:
        """Extract all slides as normalized slide models."""

        if self._parsed_content is not None and self._parsed_content.slides:
            return list(self._parsed_content.slides)

        presentation = self._require_presentation()
        transitions = self.detect_transitions()
        slides: list[Slide] = []

        for index, pptx_slide in enumerate(presentation.slides):
            slide_model = self._extract_slide_model(pptx_slide, index)
            slide_model.transition = transitions.get(index)
            slides.append(slide_model)

        if self._parsed_content is not None:
            self._parsed_content.slides = slides
        return slides

    def extract_text(self) -> str:
        """Return concatenated text from all slides."""

        text_chunks: list[str] = []
        for slide in self.extract_all_slides():
            if slide.title:
                title_text = self._paragraph_to_plain_text(slide.title).strip()
                if title_text:
                    text_chunks.append(title_text)

            for paragraph in slide.body:
                text = self._paragraph_to_plain_text(paragraph).strip()
                if text:
                    text_chunks.append(text)

            for table in slide.tables:
                for row in table.rows:
                    for cell in row:
                        for paragraph in cell.paragraphs:
                            text = self._paragraph_to_plain_text(paragraph).strip()
                            if text:
                                text_chunks.append(text)

        return "\n".join(text_chunks)

    def extract_speaker_notes(self) -> list[str]:
        """Return speaker notes text for each slide."""

        notes: list[str] = []
        for slide in self.extract_all_slides():
            note_text = "\n".join(
                self._paragraph_to_plain_text(paragraph).strip()
                for paragraph in slide.notes
                if self._paragraph_to_plain_text(paragraph).strip()
            )
            notes.append(note_text)
        return notes

    def extract_shapes(self) -> list[Shape]:
        """Return all non-table, non-chart, non-image shapes."""

        return [shape for slide in self.extract_all_slides() for shape in slide.shapes]

    def extract_tables(self) -> list[Table]:
        """Return all tables from all slides."""

        return [table for slide in self.extract_all_slides() for table in slide.tables]

    def extract_images(self) -> list[Image]:
        """Return all embedded images from all slides."""

        return [image for slide in self.extract_all_slides() for image in slide.images]

    def extract_charts(self) -> list[Chart]:
        """Return all charts from all slides."""

        return [chart for slide in self.extract_all_slides() for chart in slide.charts]

    def extract_metadata(self) -> Metadata:
        """Extract core, app, and custom document metadata."""

        presentation = self._require_presentation()
        core = presentation.core_properties

        metadata = Metadata(
            title=self._none_if_blank(getattr(core, "title", None)),
            subject=self._none_if_blank(getattr(core, "subject", None)),
            author=self._none_if_blank(getattr(core, "author", None)),
            category=self._none_if_blank(getattr(core, "category", None)),
            keywords=self._split_keywords(getattr(core, "keywords", None)),
            created_at=getattr(core, "created", None),
            modified_at=getattr(core, "modified", None),
            last_printed_at=getattr(core, "last_printed", None),
            revision=self._normalize_revision(getattr(core, "revision", None)),
            custom_properties={},
        )

        if self._source_path is not None and self._source_path.exists():
            app_props = self._read_app_properties(self._source_path)
            custom_props = self._read_custom_properties(self._source_path)
            metadata.company = cast(str | None, app_props.get("company"))
            metadata.custom_properties.update(custom_props)

        return metadata

    def get_slide_layouts(self) -> list[dict[str, Any]]:
        """Return metadata for all slide layouts in the presentation."""

        presentation = self._require_presentation()
        layouts: list[dict[str, Any]] = []
        for master_index, master in enumerate(presentation.slide_masters):
            for layout_index, layout in enumerate(master.slide_layouts):
                layout_name = self._none_if_blank(getattr(layout, "name", None))
                layouts.append(
                    {
                        "master_index": master_index,
                        "layout_index": layout_index,
                        "name": layout_name or f"Layout {layout_index}",
                        "shape_count": len(layout.shapes),
                        "placeholder_count": len(layout.placeholders),
                        "used_by_count": len(
                            getattr(layout, "used_by_slides", ()) or ()
                        ),
                    }
                )
        return layouts

    def get_master_slides(self) -> list[dict[str, Any]]:
        """Return metadata describing slide masters and their layouts."""

        presentation = self._require_presentation()
        masters: list[dict[str, Any]] = []
        for index, master in enumerate(presentation.slide_masters):
            masters.append(
                {
                    "index": index,
                    "name": self._master_name(master, index),
                    "shape_count": len(master.shapes),
                    "placeholder_count": len(master.placeholders),
                    "layout_names": [
                        self._none_if_blank(getattr(layout, "name", None))
                        or f"Layout {layout_index}"
                        for layout_index, layout in enumerate(master.slide_layouts)
                    ],
                }
            )
        return masters

    # ------------------------------------------------------------------
    # Public advanced read API
    # ------------------------------------------------------------------
    def detect_animations(self) -> list[dict[str, Any]]:
        """Return animation markers discovered from slide XML timing data."""

        if self._animations_cache is not None:
            return list(self._animations_cache)

        presentation = self._require_presentation()
        animations: list[dict[str, Any]] = []

        for slide_index, slide in enumerate(presentation.slides):
            timing = slide.element.find(f"{{{NSMAP['p']}}}timing")
            if timing is None:
                continue

            for node in timing.iter():
                name = _local_name(node.tag)
                if name not in {
                    "anim",
                    "animClr",
                    "animEffect",
                    "animMotion",
                    "animRot",
                    "animScale",
                    "audio",
                    "cmd",
                    "par",
                    "seq",
                    "set",
                    "video",
                }:
                    continue

                attributes = {
                    _local_name(key): value for key, value in dict(node.attrib).items()
                }
                animation = {
                    "slide_index": slide_index,
                    "type": name,
                    "attributes": attributes,
                }

                target_element = node.find(".//p:spTgt", NSMAP)
                if target_element is not None:
                    animation["shape_id"] = target_element.get("spid")

                animations.append(animation)

        self._animations_cache = animations
        return list(animations)

    def detect_transitions(self) -> dict[int, dict[str, Any] | None]:
        """Return transition metadata keyed by slide index."""

        presentation = self._require_presentation()
        transitions: dict[int, dict[str, Any] | None] = {}

        for slide_index, slide in enumerate(presentation.slides):
            transition = slide.element.find(f"{{{NSMAP['p']}}}transition")
            if transition is None:
                transitions[slide_index] = None
                continue

            data: dict[str, Any] = {
                "type": "transition",
                "speed": transition.get("spd"),
                "advance_on_click": not _is_truthy(transition.get("advTm")),
                "advance_after_ms": transition.get("advTm"),
            }

            child = next(iter(transition), None)
            if child is not None:
                data["type"] = _local_name(child.tag)
                if child.attrib:
                    data["options"] = {
                        _local_name(key): value
                        for key, value in dict(child.attrib).items()
                    }

            transitions[slide_index] = data

        return transitions

    def extract_hyperlinks(self) -> list[dict[str, Any]]:
        """Return hyperlinks found in text runs and clickable shapes."""

        if self._hyperlinks_cache is not None:
            return list(self._hyperlinks_cache)

        presentation = self._require_presentation()
        hyperlinks: list[dict[str, Any]] = []

        for slide_index, slide in enumerate(presentation.slides):
            for shape_index, shape in enumerate(self._iter_shapes(slide.shapes)):
                click_action = getattr(shape, "click_action", None)
                if click_action is not None:
                    hyperlink = getattr(click_action, "hyperlink", None)
                    address = self._none_if_blank(getattr(hyperlink, "address", None))
                    target_slide = getattr(click_action, "target_slide", None)
                    if address or target_slide is not None:
                        hyperlinks.append(
                            {
                                "slide_index": slide_index,
                                "shape_index": shape_index,
                                "shape_name": self._none_if_blank(
                                    getattr(shape, "name", None)
                                ),
                                "kind": "shape",
                                "address": address,
                                "target_slide_id": getattr(
                                    target_slide, "slide_id", None
                                ),
                            }
                        )

                if not getattr(shape, "has_text_frame", False):
                    continue

                for paragraph_index, paragraph in enumerate(
                    shape.text_frame.paragraphs
                ):
                    for run_index, run in enumerate(paragraph.runs):
                        hyperlink = getattr(run, "hyperlink", None)
                        address = self._none_if_blank(
                            getattr(hyperlink, "address", None)
                        )
                        if not address:
                            continue
                        hyperlinks.append(
                            {
                                "slide_index": slide_index,
                                "shape_index": shape_index,
                                "paragraph_index": paragraph_index,
                                "run_index": run_index,
                                "kind": "text",
                                "address": address,
                                "text": run.text,
                            }
                        )

        self._hyperlinks_cache = hyperlinks
        return list(hyperlinks)

    def detect_embedded_media(self) -> list[dict[str, Any]]:
        """Detect embedded audio, video, and OLE payloads."""

        if self._media_cache is not None:
            return list(self._media_cache)

        presentation = self._require_presentation()
        media_items: list[dict[str, Any]] = []

        for slide_index, slide in enumerate(presentation.slides):
            rels = getattr(slide.part, "rels", {})
            for rel_id, rel in rels.items():
                rel_type = getattr(rel, "reltype", "")
                target_ref = getattr(rel, "target_ref", None)
                is_external = bool(getattr(rel, "is_external", False))

                media_kind: str | None = None
                for kind, rel_types in MEDIA_REL_TYPES.items():
                    if rel_type in rel_types:
                        media_kind = kind
                        break

                if media_kind is None and target_ref:
                    target_lower = str(target_ref).lower()
                    if any(
                        target_lower.endswith(ext) for ext in (".mp4", ".mov", ".avi")
                    ):
                        media_kind = "video"
                    elif any(
                        target_lower.endswith(ext) for ext in (".mp3", ".wav", ".m4a")
                    ):
                        media_kind = "audio"

                if media_kind is None:
                    continue

                media_items.append(
                    {
                        "slide_index": slide_index,
                        "relationship_id": rel_id,
                        "type": media_kind,
                        "relationship_type": rel_type,
                        "target_ref": target_ref,
                        "is_external": is_external,
                    }
                )

        self._media_cache = media_items
        return list(media_items)

    def generate_slide_thumbnails(
        self,
        size: tuple[int, int] = DEFAULT_THUMBNAIL_SIZE,
    ) -> list[Image]:
        """Generate best-effort slide thumbnail previews.

        python-pptx does not render slides to bitmap output. This method creates
        compact textual preview cards containing the slide number, title, and a
        short excerpt of body content when Pillow is available.
        """

        if self._thumbnails_cache is not None and size == DEFAULT_THUMBNAIL_SIZE:
            return list(self._thumbnails_cache)

        if not PIL_AVAILABLE:
            return []

        thumbnails: list[Image] = []
        slides = self.extract_all_slides()

        for slide_index, slide in enumerate(slides):
            image_bytes = self._build_thumbnail_bytes(slide, slide_index, size)
            thumbnails.append(
                Image(
                    data=image_bytes,
                    mime_type="image/png",
                    title=f"Slide {slide_index + 1} thumbnail",
                    description="Generated preview thumbnail",
                    alternate_text=f"Preview for slide {slide_index + 1}",
                    width=float(size[0]),
                    height=float(size[1]),
                    properties={
                        "generated": True,
                        "slide_index": slide_index,
                        "kind": "thumbnail",
                    },
                )
            )

        if size == DEFAULT_THUMBNAIL_SIZE:
            self._thumbnails_cache = thumbnails
        return thumbnails

    # ------------------------------------------------------------------
    # Public write API
    # ------------------------------------------------------------------
    def create_presentation(self, template_path: str | Path | None = None) -> Any:
        """Create a new presentation or load a template for editing."""

        self._ensure_pptx_available()
        if template_path is None:
            self._presentation = PptxPresentationFactory()
            self._source_path = None
        else:
            template = Path(template_path)
            self._validate_path(template)
            self._presentation = PptxPresentationFactory(str(template))
            self._source_path = template

        self._invalidate_caches()
        return self._presentation

    def add_slide(self, layout: int | str | Any = 0) -> Any:
        """Add a slide using a layout index, layout name, or layout object."""

        presentation = self._require_writable_presentation()
        layout_obj = self._resolve_layout(layout)
        slide = presentation.slides.add_slide(layout_obj)
        self._invalidate_caches()
        return slide

    def add_text_box(
        self,
        slide: Any,
        text: str,
        position: Mapping[str, float] | Sequence[float],
    ) -> Any:
        """Add a text box to ``slide`` at the supplied position."""

        left, top, width, height = self._normalize_position(position)
        shape = slide.shapes.add_textbox(left, top, width, height)
        shape.text_frame.text = text
        self._invalidate_caches()
        return shape

    def add_image(
        self,
        slide: Any,
        path: str | Path | BinaryIO,
        position: Mapping[str, float] | Sequence[float],
    ) -> Any:
        """Add an image to ``slide`` at the supplied position."""

        left, top, width, height = self._normalize_position(position)
        image_source: str | BinaryIO = (
            str(Path(path)) if isinstance(path, (str, Path)) else path
        )
        picture = slide.shapes.add_picture(image_source, left, top, width, height)
        self._invalidate_caches()
        return picture

    def add_table(
        self,
        slide: Any,
        data: Sequence[Sequence[Any]],
        position: Mapping[str, float] | Sequence[float],
    ) -> Any:
        """Add a table populated from a row/column sequence."""

        rows = len(data)
        cols = max((len(row) for row in data), default=0)
        if rows == 0 or cols == 0:
            raise DocumentValidationError("Table data must contain at least one cell")

        left, top, width, height = self._normalize_position(position)
        graphic_frame = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = graphic_frame.table

        for row_index, row in enumerate(data):
            for col_index, value in enumerate(row):
                table.cell(row_index, col_index).text = (
                    "" if value is None else str(value)
                )

        self._invalidate_caches()
        return graphic_frame

    def add_chart(
        self,
        slide: Any,
        chart_type: str | Any,
        data: Mapping[str, Any],
        position: Mapping[str, float] | Sequence[float] | None = None,
    ) -> Any:
        """Add a chart to a slide using normalized chart input."""

        if ChartData is None or XL_CHART_TYPE is None:
            raise DocumentValidationError(
                "python-pptx chart support is unavailable. Install agentic-brain[documents]."
            )

        chart_enum = self._resolve_chart_type(chart_type)
        chart_data = ChartData()
        categories = list(data.get("categories", []))
        chart_data.categories = categories

        for series in data.get("series", []):
            if isinstance(series, Mapping):
                name = str(series.get("name", "Series"))
                values = list(series.get("values", []))
            elif isinstance(series, Sequence) and len(series) >= 2:
                name = str(series[0])
                values = list(series[1])
            else:
                raise DocumentValidationError("Invalid chart series payload")
            chart_data.add_series(name, values)

        left, top, width, height = self._normalize_position(
            position or {"x": 1.0, "y": 1.5, "width": 8.0, "height": 4.5}
        )
        graphic_frame = slide.shapes.add_chart(
            chart_enum,
            left,
            top,
            width,
            height,
            chart_data,
        )
        self._invalidate_caches()
        return graphic_frame

    def add_speaker_notes(self, slide: Any, notes: str | Sequence[str]) -> Any:
        """Attach speaker notes to a slide."""

        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        if isinstance(notes, str):
            lines = notes.splitlines() or [notes]
        else:
            lines = [str(line) for line in notes]

        text_frame.clear()
        for index, line in enumerate(lines):
            paragraph = (
                text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            )
            paragraph.text = line

        self._invalidate_caches()
        return notes_slide

    def save(self, path: str | Path) -> None:
        """Save the current presentation to disk."""

        presentation = self._require_writable_presentation()
        output_path = Path(path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))
        self._source_path = output_path
        self._invalidate_caches()

    # ------------------------------------------------------------------
    # Internal extraction helpers
    # ------------------------------------------------------------------
    def _extract_slide_model(self, slide: Any, slide_index: int) -> Slide:
        """Convert a python-pptx slide to the normalized slide model."""

        title = self._extract_slide_title(slide)
        body: list[Paragraph] = []
        notes = self._extract_notes_paragraphs(slide)
        images: list[Image] = []
        shapes: list[Shape] = []
        tables: list[Table] = []
        charts: list[Chart] = []

        for shape in self._iter_shapes(slide.shapes):
            if self._is_image_shape(shape):
                image = self._extract_image(shape)
                if image is not None:
                    images.append(image)
                continue

            if getattr(shape, "has_table", False):
                tables.append(self._extract_table(shape.table))
                continue

            if getattr(shape, "has_chart", False):
                charts.append(self._extract_chart(shape.chart, shape))
                continue

            normalized_shape = self._extract_shape(shape)
            if normalized_shape is not None:
                shapes.append(normalized_shape)

            if getattr(shape, "has_text_frame", False):
                for paragraph in self._extract_text_frame(shape.text_frame):
                    if self._is_title_shape(shape, slide):
                        continue
                    if self._paragraph_to_plain_text(paragraph).strip():
                        body.append(paragraph)

        return Slide(
            title=title,
            body=body,
            notes=notes,
            images=images,
            shapes=shapes,
            tables=tables,
            charts=charts,
            background=self._extract_background(slide),
            layout_name=self._none_if_blank(getattr(slide.slide_layout, "name", None)),
        )

    def _extract_slide_title(self, slide: Any) -> Paragraph | None:
        """Return the title placeholder paragraph for a slide, if present."""

        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is None or not getattr(title_shape, "has_text_frame", False):
            return None

        paragraphs = self._extract_text_frame(title_shape.text_frame)
        if not paragraphs:
            return None

        combined_runs: list[TextRun] = []
        for paragraph in paragraphs:
            combined_runs.extend(paragraph.runs)

        if not combined_runs:
            return None

        title_style = paragraphs[0].style if paragraphs else DocumentStyle()
        return Paragraph(
            runs=combined_runs, style=title_style, is_heading=True, heading_level=1
        )

    def _extract_notes_paragraphs(self, slide: Any) -> list[Paragraph]:
        """Extract speaker notes from a slide without forcing empty note creation."""

        if not getattr(slide, "has_notes_slide", False):
            return []

        notes_slide = slide.notes_slide
        text_frame = getattr(notes_slide, "notes_text_frame", None)
        if text_frame is None:
            return []
        return self._extract_text_frame(text_frame)

    def _extract_text_frame(self, text_frame: Any) -> list[Paragraph]:
        """Convert a pptx text frame into paragraph models."""

        paragraphs: list[Paragraph] = []
        for paragraph in getattr(text_frame, "paragraphs", []):
            runs: list[TextRun] = []
            for run in getattr(paragraph, "runs", []):
                style = self._extract_font_style(getattr(run, "font", None))
                hyperlink = getattr(getattr(run, "hyperlink", None), "address", None)
                runs.append(
                    TextRun(
                        text=run.text or "",
                        style=style,
                        hyperlink=self._none_if_blank(hyperlink),
                    )
                )

            if not runs and getattr(paragraph, "text", ""):
                runs = [TextRun(text=paragraph.text, style=DocumentStyle())]

            paragraph_style = self._extract_paragraph_style(paragraph)
            paragraphs.append(
                Paragraph(
                    runs=runs,
                    style=paragraph_style,
                    numbering_level=getattr(paragraph, "level", None),
                )
            )

        return paragraphs

    def _extract_font_style(self, font: Any) -> DocumentStyle:
        """Build a style model from a run font."""

        style = DocumentStyle()
        if font is None:
            return style

        style.font_family = (
            self._none_if_blank(getattr(font, "name", None)) or style.font_family
        )
        font_size = getattr(font, "size", None)
        if font_size is not None:
            try:
                style.font_size = float(font_size.pt)
            except AttributeError:
                style.font_size = float(font_size)
        style.bold = bool(getattr(font, "bold", False))
        style.italic = bool(getattr(font, "italic", False))
        style.underline = bool(getattr(font, "underline", False))
        style.text_color = (
            self._extract_color(getattr(font, "color", None)) or style.text_color
        )
        return style

    def _extract_paragraph_style(self, paragraph: Any) -> DocumentStyle:
        """Extract paragraph-level styling details."""

        style = DocumentStyle()
        alignment = getattr(paragraph, "alignment", None)
        if alignment is not None:
            style.alignment = getattr(alignment, "name", str(alignment)).lower()
        return style

    def _extract_shape(self, shape: Any) -> Shape | None:
        """Convert a non-image/table/chart shape into a normalized shape model."""

        shape_type = self._shape_type_name(shape)
        if shape_type == "GROUP":
            return None

        text_paragraph = None
        if getattr(shape, "has_text_frame", False):
            paragraphs = self._extract_text_frame(shape.text_frame)
            if paragraphs:
                text_paragraph = self._merge_paragraphs(paragraphs)

        return Shape(
            shape_type=shape_type,
            path=[],
            fill_color=self._extract_fill_color(shape),
            stroke_color=self._extract_line_color(shape),
            stroke_width=self._extract_line_width(shape),
            text=text_paragraph,
            properties={
                "name": self._none_if_blank(getattr(shape, "name", None)) or "",
                "rotation": float(getattr(shape, "rotation", 0) or 0),
                "left_pt": _pt(getattr(shape, "left", None)) or 0.0,
                "top_pt": _pt(getattr(shape, "top", None)) or 0.0,
                "width_pt": _pt(getattr(shape, "width", None)) or 0.0,
                "height_pt": _pt(getattr(shape, "height", None)) or 0.0,
                "placeholder": bool(getattr(shape, "is_placeholder", False)),
            },
        )

    def _extract_table(self, table: Any) -> Table:
        """Convert a pptx table into the normalized table model."""

        rows: list[list[TableCell]] = []
        for row in table.rows:
            normalized_row: list[TableCell] = []
            for cell in row.cells:
                normalized_row.append(
                    TableCell(
                        paragraphs=self._extract_text_frame(cell.text_frame),
                        width=_pt(getattr(cell, "width", None)),
                        height=None,
                    )
                )
            rows.append(normalized_row)

        return Table(
            rows=rows,
            width=None,
            has_header_row=False,
        )

    def _extract_chart(self, chart: Any, shape: Any | None = None) -> Chart:
        """Convert a pptx chart into the normalized chart model."""

        title: str | None = None
        if getattr(chart, "has_title", False):
            chart_title = getattr(chart, "chart_title", None)
            text_frame = getattr(chart_title, "text_frame", None)
            if text_frame is not None:
                title = "\n".join(
                    filter(
                        None,
                        (
                            self._paragraph_to_plain_text(paragraph).strip()
                            for paragraph in self._extract_text_frame(text_frame)
                        ),
                    )
                )
                title = self._none_if_blank(title)

        categories = self._extract_chart_categories(chart)
        series_payload: list[dict[str, Any]] = []
        for series in getattr(chart, "series", []):
            values = self._extract_series_values(series)
            series_payload.append(
                {
                    "name": self._none_if_blank(getattr(series, "name", None))
                    or "Series",
                    "values": values,
                }
            )

        legend: dict[str, Any] = {}
        if getattr(chart, "has_legend", False):
            chart_legend = getattr(chart, "legend", None)
            legend = {
                "position": getattr(
                    getattr(chart_legend, "position", None),
                    "name",
                    getattr(chart_legend, "position", None),
                ),
                "include_in_layout": getattr(chart_legend, "include_in_layout", None),
            }

        position: dict[str, float] = {}
        if shape is not None:
            position = {
                "left_pt": _pt(getattr(shape, "left", None)) or 0.0,
                "top_pt": _pt(getattr(shape, "top", None)) or 0.0,
                "width_pt": _pt(getattr(shape, "width", None)) or 0.0,
                "height_pt": _pt(getattr(shape, "height", None)) or 0.0,
            }

        chart_type = getattr(getattr(chart, "chart_type", None), "name", None) or str(
            getattr(chart, "chart_type", "UNKNOWN")
        )
        return Chart(
            chart_type=str(chart_type),
            title=title,
            series=series_payload,
            categories=categories,
            legend=legend,
            position=position,
        )

    def _extract_chart_categories(self, chart: Any) -> list[str]:
        """Extract chart categories using public API when possible."""

        categories: list[str] = []
        for plot in getattr(chart, "plots", []):
            plot_categories = getattr(plot, "categories", None)
            if plot_categories is None:
                continue
            try:
                for category in plot_categories:
                    label = getattr(category, "label", category)
                    categories.append(str(label))
            except TypeError:
                continue
            if categories:
                break
        return categories

    def _extract_series_values(self, series: Any) -> list[float]:
        """Extract numeric values from a chart series."""

        values: list[float] = []
        series_values = getattr(series, "values", None)
        if series_values is None:
            return values
        try:
            for value in series_values:
                if value is None:
                    values.append(0.0)
                else:
                    values.append(float(value))
        except TypeError:
            return values
        return values

    def _extract_image(self, shape: Any) -> Image | None:
        """Convert a picture shape to the normalized image model."""

        image_obj = getattr(shape, "image", None)
        if image_obj is None:
            return None

        mime_type = (
            getattr(image_obj, "content_type", None) or "application/octet-stream"
        )
        description = self._extract_alt_text(shape)
        return Image(
            data=getattr(image_obj, "blob", b""),
            mime_type=mime_type,
            description=description,
            title=self._none_if_blank(getattr(shape, "name", None)),
            alternate_text=description,
            width=_pt(getattr(shape, "width", None)),
            height=_pt(getattr(shape, "height", None)),
            position={
                "left_pt": _pt(getattr(shape, "left", None)) or 0.0,
                "top_pt": _pt(getattr(shape, "top", None)) or 0.0,
                "width_pt": _pt(getattr(shape, "width", None)) or 0.0,
                "height_pt": _pt(getattr(shape, "height", None)) or 0.0,
            },
            properties={
                "filename": getattr(image_obj, "filename", ""),
                "ext": getattr(image_obj, "ext", ""),
                "dpi": self._extract_image_dpi(image_obj),
            },
        )

    def _extract_background(self, slide: Any) -> str | None:
        """Extract slide background color if explicitly available."""

        background = getattr(slide, "background", None)
        if background is None:
            return None

        fill = getattr(background, "fill", None)
        if fill is None:
            return None
        return self._extract_color(getattr(fill, "fore_color", None))

    # ------------------------------------------------------------------
    # Internal writer helpers
    # ------------------------------------------------------------------
    def _resolve_layout(self, layout: int | str | Any) -> Any:
        """Resolve a layout specifier into a python-pptx slide layout object."""

        presentation = self._require_writable_presentation()
        if isinstance(layout, int):
            return presentation.slide_layouts[layout]

        if isinstance(layout, str):
            direct = presentation.slide_layouts.get_by_name(layout)
            if direct is not None:
                return direct
            for master in presentation.slide_masters:
                found = master.slide_layouts.get_by_name(layout)
                if found is not None:
                    return found
            raise DocumentValidationError(f"Slide layout not found: {layout}")

        return layout

    def _resolve_chart_type(self, chart_type: str | Any) -> Any:
        """Normalize a chart type into an XL_CHART_TYPE enum value."""

        if XL_CHART_TYPE is None:
            raise DocumentValidationError("Chart support is unavailable")

        if not isinstance(chart_type, str):
            return chart_type

        normalized = chart_type.strip().upper()
        aliases = {
            "BAR": "BAR_CLUSTERED",
            "COLUMN": "COLUMN_CLUSTERED",
            "LINE": "LINE",
            "PIE": "PIE",
            "AREA": "AREA",
            "DOUGHNUT": "DOUGHNUT",
            "SCATTER": "XY_SCATTER",
        }
        normalized = aliases.get(normalized, normalized)
        try:
            return getattr(XL_CHART_TYPE, normalized)
        except AttributeError as exc:
            raise DocumentValidationError(
                f"Unsupported chart type: {chart_type}"
            ) from exc

    def _normalize_position(
        self,
        position: Mapping[str, float] | Sequence[float],
    ) -> tuple[Any, Any, Any, Any]:
        """Normalize user-friendly positional input into EMU lengths."""

        if Inches is None and Emu is None:
            raise DocumentValidationError(
                "python-pptx length utilities are unavailable"
            )

        unit = "in"
        if isinstance(position, Mapping):
            unit = str(position.get("unit", "in")).lower()
            x = float(position.get("x", position.get("left", 1.0)))
            y = float(position.get("y", position.get("top", 1.0)))
            width = float(position.get("width", 4.0))
            height = float(position.get("height", 3.0))
        else:
            if len(position) != 4:
                raise DocumentValidationError(
                    "Position sequences must contain exactly four values"
                )
            x, y, width, height = (float(value) for value in position)

        if unit in {"pt", "point", "points"}:
            return Pt(x), Pt(y), Pt(width), Pt(height)
        if unit in {"emu"}:
            return Emu(int(x)), Emu(int(y)), Emu(int(width)), Emu(int(height))
        return Inches(x), Inches(y), Inches(width), Inches(height)

    # ------------------------------------------------------------------
    # Shared utility helpers
    # ------------------------------------------------------------------
    def _ensure_pptx_available(self) -> None:
        """Raise a helpful error when python-pptx is unavailable."""

        if not PPTX_AVAILABLE:
            raise UnsupportedOfficeFormatError(OfficeFormat.PPTX)

    def _validate_path(self, path: Path) -> None:
        """Validate that the path exists and points to a PowerPoint file."""

        if not path.exists():
            raise DocumentValidationError("PowerPoint file not found", str(path))
        if not path.is_file():
            raise DocumentValidationError("PowerPoint path is not a file", str(path))
        if path.suffix.lower() != ".pptx":
            raise DocumentValidationError(
                "Unsupported PowerPoint file extension",
                f"Expected .pptx, got {path.suffix or '<none>'}",
            )

    def _require_source_path(self) -> Path:
        """Return the currently configured source path."""

        if self._source_path is None:
            raise DocumentValidationError("No PowerPoint source path configured")
        return self._source_path

    def _require_presentation(self) -> Any:
        """Return a loaded presentation, loading from disk if needed."""

        if self._presentation is not None:
            return self._presentation

        path = self._require_source_path()
        self._validate_path(path)
        self._ensure_pptx_available()
        self._presentation = PptxPresentationFactory(str(path))
        return self._presentation

    def _require_writable_presentation(self) -> Any:
        """Return a presentation for writing operations."""

        if self._presentation is None:
            return self.create_presentation()
        return self._presentation

    def _invalidate_caches(self) -> None:
        """Clear cached extracted content after load or write changes."""

        self._parsed_content = None
        self._hyperlinks_cache = None
        self._animations_cache = None
        self._media_cache = None
        self._thumbnails_cache = None

    def _iter_shapes(self, shapes: Iterable[Any]) -> Iterable[Any]:
        """Yield shapes recursively, flattening groups."""

        for shape in shapes:
            yield shape
            if self._shape_type_name(shape) == "GROUP":
                yield from self._iter_shapes(shape.shapes)

    def _shape_type_name(self, shape: Any) -> str:
        """Return a readable shape type name."""

        shape_type = getattr(shape, "shape_type", None)
        return str(getattr(shape_type, "name", shape_type or "UNKNOWN"))

    def _is_image_shape(self, shape: Any) -> bool:
        """Return True when the shape is a picture-like element."""

        shape_type_name = self._shape_type_name(shape)
        if shape_type_name in {"PICTURE", "LINKED_PICTURE"}:
            return True
        return getattr(shape, "image", None) is not None

    def _is_title_shape(self, shape: Any, slide: Any) -> bool:
        """Return True when a shape is the title placeholder for a slide."""

        title_shape = getattr(slide.shapes, "title", None)
        return title_shape is not None and shape == title_shape

    def _extract_fill_color(self, shape: Any) -> str | None:
        """Extract the foreground fill color for a shape."""

        fill = getattr(getattr(shape, "fill", None), "fore_color", None)
        return self._extract_color(fill)

    def _extract_line_color(self, shape: Any) -> str | None:
        """Extract line color for a shape."""

        line = getattr(shape, "line", None)
        return self._extract_color(getattr(line, "color", None))

    def _extract_line_width(self, shape: Any) -> float:
        """Extract line width in points for a shape."""

        line = getattr(shape, "line", None)
        width = getattr(line, "width", None)
        return _pt(width) or 1.0

    def _extract_color(self, color_obj: Any) -> str | None:
        """Convert a python-pptx color object into a CSS-like hex string."""

        if color_obj is None:
            return None

        rgb = getattr(color_obj, "rgb", None)
        if rgb is not None:
            value = str(rgb)
            value = value.replace("0x", "").replace("0X", "")
            if value:
                return f"#{value}"

        theme_color = getattr(color_obj, "theme_color", None)
        if theme_color is not None:
            theme_name = getattr(theme_color, "name", str(theme_color))
            return f"theme:{theme_name}"

        return None

    def _extract_alt_text(self, shape: Any) -> str | None:
        """Extract alternate text / description from the underlying XML."""

        nv_props = shape.element.find(".//p:cNvPr", NSMAP)
        if nv_props is None:
            return None

        description = self._none_if_blank(nv_props.get("descr"))
        if description:
            return description
        return self._none_if_blank(nv_props.get("title"))

    def _extract_image_dpi(self, image_obj: Any) -> str:
        """Extract human-readable image dpi details when available."""

        dpi = getattr(image_obj, "dpi", None)
        if not dpi:
            return ""
        if isinstance(dpi, Sequence) and len(dpi) >= 2:
            return f"{dpi[0]}x{dpi[1]}"
        return str(dpi)

    def _merge_paragraphs(self, paragraphs: Sequence[Paragraph]) -> Paragraph:
        """Combine multiple paragraphs into a single paragraph-like summary."""

        merged_runs: list[TextRun] = []
        for index, paragraph in enumerate(paragraphs):
            if index > 0:
                merged_runs.append(TextRun(text="\n"))
            merged_runs.extend(paragraph.runs)
        style = paragraphs[0].style if paragraphs else DocumentStyle()
        return Paragraph(runs=merged_runs, style=style)

    def _paragraph_to_plain_text(self, paragraph: Paragraph) -> str:
        """Flatten a paragraph into plain text."""

        return "".join(run.text for run in paragraph.runs)

    def _none_if_blank(self, value: Any) -> str | None:
        """Convert blank strings to None."""

        if value is None:
            return None
        text = str(value).strip()
        return text or None

    def _split_keywords(self, keywords: Any) -> list[str]:
        """Split keyword strings into individual values."""

        if keywords is None:
            return []
        if isinstance(keywords, Sequence) and not isinstance(keywords, str):
            return [str(item).strip() for item in keywords if str(item).strip()]

        text = str(keywords).replace(";", ",")
        return [part.strip() for part in text.split(",") if part.strip()]

    def _normalize_revision(self, revision: Any) -> str | None:
        """Normalize revision values to a string."""

        if revision is None:
            return None
        text = str(revision).strip()
        return text or None

    def _master_name(self, master: Any, index: int) -> str:
        """Derive a readable name for a slide master."""

        name = self._none_if_blank(getattr(master, "name", None))
        if name:
            return name

        c_sld = master.element.find(".//p:cSld", NSMAP)
        if c_sld is not None:
            derived = self._none_if_blank(c_sld.get("name"))
            if derived:
                return derived
        return f"Master {index}"

    def _read_app_properties(self, path: Path) -> dict[str, Any]:
        """Read selected extended properties from ``docProps/app.xml``."""

        try:
            with ZipFile(path) as archive:
                with archive.open("docProps/app.xml") as handle:
                    root = ET.parse(handle).getroot()
        except (FileNotFoundError, KeyError, ET.ParseError):
            return {}

        company = root.findtext("ep:Company", default=None, namespaces=NSMAP)
        manager = root.findtext("ep:Manager", default=None, namespaces=NSMAP)
        application = root.findtext("ep:Application", default=None, namespaces=NSMAP)
        return {
            "company": self._none_if_blank(company),
            "manager": self._none_if_blank(manager),
            "application": self._none_if_blank(application),
        }

    def _read_custom_properties(self, path: Path) -> dict[str, Any]:
        """Read custom document properties from ``docProps/custom.xml``."""

        try:
            with ZipFile(path) as archive:
                with archive.open("docProps/custom.xml") as handle:
                    root = ET.parse(handle).getroot()
        except (FileNotFoundError, KeyError, ET.ParseError):
            return {}

        properties: dict[str, Any] = {}
        for prop in root.findall(
            ".//{http://schemas.openxmlformats.org/officeDocument/2006/custom-properties}property"
        ):
            name = prop.get("name")
            if not name:
                continue

            value: Any = None
            for child in prop:
                tag = _local_name(child.tag)
                text = child.text
                if tag in {"lpwstr", "bstr", "str"}:
                    value = text
                elif tag in {"i1", "i2", "i4", "i8", "int", "uint"}:
                    value = int(text) if text is not None else None
                elif tag in {"r4", "r8", "decimal"}:
                    value = float(text) if text is not None else None
                elif tag == "bool":
                    value = _is_truthy(text)
                elif tag in {"filetime", "date"}:
                    value = text
                else:
                    value = text
                break
            properties[name] = value
        return properties

    def _build_thumbnail_bytes(
        self,
        slide: Slide,
        slide_index: int,
        size: tuple[int, int],
    ) -> bytes:
        """Create a thumbnail preview image for a slide."""

        if not PIL_AVAILABLE or PILImage is None or ImageDraw is None:
            return b""

        image = PILImage.new("RGB", size, color=(245, 247, 250))
        draw = ImageDraw.Draw(image)
        font = ImageFont.load_default() if ImageFont is not None else None
        accent = (54, 96, 146)
        text_color = (30, 41, 59)
        subtle = (100, 116, 139)

        draw.rectangle((0, 0, size[0], 30), fill=accent)
        draw.text((12, 8), f"Slide {slide_index + 1}", fill=(255, 255, 255), font=font)

        title = (
            self._paragraph_to_plain_text(slide.title).strip() if slide.title else ""
        )
        body_lines = [
            self._paragraph_to_plain_text(paragraph).strip()
            for paragraph in slide.body
            if self._paragraph_to_plain_text(paragraph).strip()
        ]
        preview = "\n".join(body_lines[:4])

        y = 42
        if title:
            draw.text((12, y), title[:70], fill=text_color, font=font)
            y += 24

        if preview:
            wrapped = preview[:220]
            for line in self._wrap_text(wrapped, width=48):
                draw.text((12, y), line, fill=subtle, font=font)
                y += 16

        footer = (
            f"Images: {len(slide.images)}  "
            f"Tables: {len(slide.tables)}  "
            f"Charts: {len(slide.charts)}"
        )
        draw.text((12, size[1] - 20), footer, fill=subtle, font=font)

        buffer = BytesIO()
        image.save(buffer, format="PNG")
        return buffer.getvalue()

    def _wrap_text(self, text: str, width: int) -> list[str]:
        """Wrap text into lines with a simple character-based heuristic."""

        words = text.split()
        if not words:
            return []

        lines: list[str] = []
        current = words[0]
        for word in words[1:]:
            candidate = f"{current} {word}"
            if len(candidate) <= width:
                current = candidate
            else:
                lines.append(current)
                current = word
        lines.append(current)
        return lines


__all__ = ["PowerPointProcessor"]
