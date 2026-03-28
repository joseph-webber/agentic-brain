# SPDX-License-Identifier: Apache-2.0

"""Office document format converter for agentic-brain."""

from __future__ import annotations

import builtins
import csv
import logging
import os
import platform
import re
import shutil
import subprocess
import tempfile
from enum import StrEnum
from pathlib import Path
from typing import Callable, Dict, List, Optional, Sequence, Tuple, Union

from .models import DocumentContent, OfficeFormat as ModelOfficeFormat, Paragraph, TextRun

logger = logging.getLogger(__name__)

try:  # html2text for HTML -> Markdown fallback
    import html2text

    HAS_HTML2TEXT = True
except ImportError:  # pragma: no cover - optional dependency
    html2text = None
    HAS_HTML2TEXT = False

try:  # python-pptx for presentation parsing
    from pptx import Presentation as PptxPresentation

    HAS_PPTX = True
except ImportError:  # pragma: no cover - optional dependency
    PptxPresentation = None
    HAS_PPTX = False


class ConversionError(RuntimeError):
    """Raised when a conversion workflow fails."""


class OfficeFormat(StrEnum):
    """Converter-facing office format identifiers."""

    DOCX = "docx"
    DOC = "doc"
    XLSX = "xlsx"
    XLS = "xls"
    PPTX = "pptx"
    PPT = "ppt"
    PAGES = "pages"
    NUMBERS = "numbers"
    KEYNOTE = "key"
    ODT = "odt"
    ODS = "ods"
    ODP = "odp"
    TXT = "txt"


class OfficeConverter:
    """Cross-format Office document converter."""

    WORD_LIKE_EXTENSIONS = {"docx", "doc", "rtf", "odt"}
    EXCEL_LIKE_EXTENSIONS = {"xlsx", "xls", "ods", "csv"}
    POWERPOINT_LIKE_EXTENSIONS = {"pptx", "ppt", "odp"}
    IWORK_EXTENSIONS = {"pages", "numbers", "key", "keynote"}

    FORMAT_ALIAS: Dict[str, ModelOfficeFormat] = {
        "docx": ModelOfficeFormat.DOCX,
        "doc": ModelOfficeFormat.DOCX,
        "rtf": ModelOfficeFormat.RTF,
        "odt": ModelOfficeFormat.ODT,
        "xlsx": ModelOfficeFormat.XLSX,
        "xls": ModelOfficeFormat.XLSX,
        "ods": ModelOfficeFormat.ODS,
        "csv": ModelOfficeFormat.ODS,  # treat CSV as spreadsheet content
        "pptx": ModelOfficeFormat.PPTX,
        "ppt": ModelOfficeFormat.PPTX,
        "odp": ModelOfficeFormat.ODP,
        "pages": ModelOfficeFormat.PAGES,
        "numbers": ModelOfficeFormat.NUMBERS,
        "key": ModelOfficeFormat.KEYNOTE,
        "keynote": ModelOfficeFormat.KEYNOTE,
    }

    def __init__(
        self, libreoffice_path: Optional[str | Path] = None, timeout: int = 60
    ) -> None:
        self.platform = platform.system()
        self.timeout = timeout
        self.libreoffice_path = (
            str(Path(libreoffice_path).expanduser())
            if libreoffice_path
            else self._find_libreoffice()
        )
        self.textutil_path = (
            Path(shutil.which("textutil")).resolve()  # type: ignore[arg-type]
            if self.platform == "Darwin" and shutil.which("textutil")
            else None
        )
        self._has_docx2pdf = False
        self._has_mammoth = False
        self._has_python_docx = False
        self._has_openpyxl = False
        self._check_dependencies()

        if not self.libreoffice_path:
            logger.debug("LibreOffice not detected. Some conversions will be limited.")
        if self.platform == "Darwin" and not self.textutil_path:
            logger.debug(
                "textutil not available. iWork conversions will rely on LibreOffice."
            )

    # ------------------------------------------------------------------ Public API

    def to_pdf(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert a document into PDF."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".pdf")
        extension = source.suffix.lstrip(".").lower()

        logger.debug("Converting %s -> %s (pdf)", source, target)

        if extension in self.WORD_LIKE_EXTENSIONS:
            if self._can_use_docx2pdf():
                try:
                    return self._convert_with_docx2pdf(source, target)
                except ConversionError as exc:
                    logger.debug(
                        "docx2pdf failed (%s); falling back to LibreOffice", exc
                    )
            return self._convert_with_libreoffice(source, target, "pdf")

        if extension in self.IWORK_EXTENSIONS and self.textutil_path:
            try:
                return self._convert_with_textutil(source, target, "pdf")
            except ConversionError as exc:
                logger.debug("textutil failed (%s); falling back to LibreOffice", exc)

        # Excel / PowerPoint / OpenDocument or any other supported format
        return self._convert_with_libreoffice(source, target, "pdf")

    def to_text(
        self,
        input_path: str | Path,
        output_path: Optional[str | Path] = None,
        *,
        encoding: str = "utf-8",
    ) -> Path:
        """Extract structured plain text and persist it to ``*.txt``."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".txt")

        text_content, _ = self._extract_textual_content(source)
        target.write_text(text_content, encoding=encoding)

        logger.info("Extracted text from %s -> %s", source.name, target.name)
        self._maybe_emit_debug_preview(text_content)

        return target

    def docx_to_pdf(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert DOCX (or legacy DOC) into PDF."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".pdf")

        if self._has_docx2pdf and self.platform in {"Windows", "Darwin"}:
            try:
                docx2pdf = builtins.__import__("docx2pdf")
                docx2pdf.convert(str(source), str(target))
                return target
            except Exception as exc:  # pragma: no cover - external dependency
                logger.debug("docx2pdf failed (%s); falling back to LibreOffice", exc)

        converted = self._run_libreoffice_conversion(source, "pdf", target.parent)
        return self._rename_if_needed(converted, target)

    def docx_to_txt(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert DOCX into plain text using python-docx when available."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".txt")

        if self._has_python_docx:
            try:
                docx = builtins.__import__("docx")
                document = docx.Document(str(source))
                lines = [para.text for para in document.paragraphs if para.text]
                target.write_text("\n".join(lines), encoding="utf-8")
                return target
            except Exception as exc:
                logger.debug("python-docx failed (%s); falling back to LibreOffice", exc)

        converted = self._run_libreoffice_conversion(source, "txt", target.parent)
        return self._rename_if_needed(converted, target)

    def docx_to_html(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert DOCX into HTML using mammoth when available."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".html")

        if self._has_mammoth:
            try:
                mammoth = builtins.__import__("mammoth")
                with source.open("rb") as handle:
                    result = mammoth.convert_to_html(handle)
                for message in getattr(result, "messages", []):
                    logger.debug("mammoth notice: %s", message)
                target.write_text(result.value, encoding="utf-8")
                return target
            except Exception as exc:
                logger.debug("mammoth failed (%s); falling back to LibreOffice", exc)

        converted = self._run_libreoffice_conversion(source, "html", target.parent)
        return self._rename_if_needed(converted, target)

    def docx_to_markdown(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert DOCX into Markdown using mammoth with HTML fallback."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".md")

        if self._has_mammoth:
            try:
                mammoth = builtins.__import__("mammoth")
                with source.open("rb") as handle:
                    result = mammoth.convert_to_markdown(handle)
                target.write_text(result.value, encoding="utf-8")
                return target
            except Exception as exc:
                logger.debug("mammoth markdown failed (%s); falling back to HTML", exc)

        html_path = target.with_suffix(".html")
        try:
            self.docx_to_html(source, html_path)
            html_text = html_path.read_text(encoding="utf-8", errors="ignore")
            markdown = self._html_to_markdown(html_text)
            target.write_text(markdown, encoding="utf-8")
        finally:
            html_path.unlink(missing_ok=True)

        return target

    def xlsx_to_pdf(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert XLSX to PDF using LibreOffice."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".pdf")
        converted = self._run_libreoffice_conversion(source, "pdf", target.parent)
        return self._rename_if_needed(converted, target)

    def xlsx_to_csv(
        self,
        input_path: str | Path,
        output_path: Optional[str | Path] = None,
        *,
        sheet: Optional[str] = None,
    ) -> Path:
        """Convert XLSX to CSV using openpyxl when available."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".csv")

        if self._has_openpyxl:
            try:
                openpyxl = builtins.__import__("openpyxl")
                workbook = openpyxl.load_workbook(str(source), data_only=True)
                worksheet = (
                    workbook[sheet] if sheet is not None else workbook.active
                )
                with target.open("w", encoding="utf-8", newline="") as handle:
                    writer = csv.writer(handle)
                    for row in worksheet.iter_rows(values_only=True):
                        writer.writerow(
                            ["" if cell is None else str(cell) for cell in row]
                        )
                if hasattr(workbook, "close"):
                    workbook.close()
                return target
            except KeyError:
                logger.debug(
                    "Requested sheet %s missing; falling back to LibreOffice", sheet
                )
            except Exception as exc:
                logger.debug("openpyxl failed (%s); falling back to LibreOffice", exc)

        converted = self._run_libreoffice_conversion(source, "csv", target.parent)
        return self._rename_if_needed(converted, target)

    def xlsx_to_ods(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert XLSX to ODS using LibreOffice."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".ods")
        converted = self._run_libreoffice_conversion(source, "ods", target.parent)
        return self._rename_if_needed(converted, target)

    def pptx_to_pdf(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert PPTX to PDF using LibreOffice."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".pdf")
        converted = self._run_libreoffice_conversion(source, "pdf", target.parent)
        return self._rename_if_needed(converted, target)

    def pages_to_pdf(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert Pages to PDF using macOS automation."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".pdf")
        return self._run_macos_automation("Pages", source, target)

    def numbers_to_pdf(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert Numbers to PDF using macOS automation."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".pdf")
        return self._run_macos_automation("Numbers", source, target)

    def keynote_to_pdf(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert Keynote to PDF using macOS automation."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".pdf")
        return self._run_macos_automation("Keynote", source, target)

    def odt_to_pdf(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert ODT to PDF using LibreOffice."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".pdf")
        converted = self._run_libreoffice_conversion(source, "pdf", target.parent)
        return self._rename_if_needed(converted, target)

    def docx_to_odt(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert DOCX to ODT using LibreOffice."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".odt")
        converted = self._run_libreoffice_conversion(source, "odt", target.parent)
        return self._rename_if_needed(converted, target)

    def odt_to_docx(
        self, input_path: str | Path, output_path: Optional[str | Path] = None
    ) -> Path:
        """Convert ODT to DOCX using LibreOffice."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".docx")
        converted = self._run_libreoffice_conversion(source, "docx", target.parent)
        return self._rename_if_needed(converted, target)

    def detect_format(self, path: str | Path) -> Optional[OfficeFormat]:
        """Detect supported office format for conversion purposes."""

        extension = Path(path).suffix.lower().lstrip(".")
        if not extension:
            return None
        if extension == "rtf":
            return None
        mapping = {
            "docx": OfficeFormat.DOCX,
            "doc": OfficeFormat.DOC,
            "xlsx": OfficeFormat.XLSX,
            "xls": OfficeFormat.XLS,
            "pptx": OfficeFormat.PPTX,
            "ppt": OfficeFormat.PPT,
            "pages": OfficeFormat.PAGES,
            "numbers": OfficeFormat.NUMBERS,
            "key": OfficeFormat.KEYNOTE,
            "keynote": OfficeFormat.KEYNOTE,
            "odt": OfficeFormat.ODT,
            "ods": OfficeFormat.ODS,
            "odp": OfficeFormat.ODP,
            "txt": OfficeFormat.TXT,
        }
        return mapping.get(extension)

    def get_converter(
        self, source_format: str, target_format: str
    ) -> Optional[Callable[[Path, Optional[Path]], Path]]:
        """Return a conversion callable for the requested pair."""

        source = source_format.lower().lstrip(".")
        target = target_format.lower().lstrip(".")
        converters: Dict[Tuple[str, str], Callable[[Path, Optional[Path]], Path]] = {
            ("docx", "pdf"): self.docx_to_pdf,
            ("docx", "txt"): self.docx_to_txt,
            ("docx", "html"): self.docx_to_html,
            ("docx", "md"): self.docx_to_markdown,
            ("docx", "markdown"): self.docx_to_markdown,
            ("docx", "odt"): self.docx_to_odt,
            ("xlsx", "pdf"): self.xlsx_to_pdf,
            ("xlsx", "csv"): self.xlsx_to_csv,
            ("xlsx", "ods"): self.xlsx_to_ods,
            ("pptx", "pdf"): self.pptx_to_pdf,
            ("pages", "pdf"): self.pages_to_pdf,
            ("numbers", "pdf"): self.numbers_to_pdf,
            ("key", "pdf"): self.keynote_to_pdf,
            ("keynote", "pdf"): self.keynote_to_pdf,
            ("odt", "pdf"): self.odt_to_pdf,
            ("odt", "docx"): self.odt_to_docx,
        }
        return converters.get((source, target))

    def is_conversion_supported(self, source_format: str, target_format: str) -> bool:
        """Return True if a conversion workflow exists."""

        return self.get_converter(source_format, target_format) is not None

    def convert_batch(
        self,
        input_files: Sequence[str | Path],
        target_format: str,
        output_dir: Optional[str | Path] = None,
    ) -> List[Path]:
        """Convert multiple files to the same target format."""

        fmt = target_format.lower().lstrip(".")
        outputs: List[Path] = []
        output_dir_path = Path(output_dir).expanduser() if output_dir else None
        if output_dir_path:
            output_dir_path.mkdir(parents=True, exist_ok=True)

        for file_path in input_files:
            detected = self.detect_format(file_path)
            if detected is None or detected == OfficeFormat.TXT:
                continue

            converter = self.get_converter(detected.value, fmt)
            if converter is None:
                continue

            source = Path(file_path)
            destination = (
                output_dir_path / f"{source.stem}.{fmt}" if output_dir_path else None
            )
            try:
                outputs.append(converter(source, destination))
            except ConversionError:
                continue

        return outputs

    def convert_folder(
        self,
        input_dir: str | Path,
        target_format: str,
        output_dir: Optional[str | Path] = None,
        *,
        recursive: bool = False,
    ) -> List[Path]:
        """Convert all supported files in a directory."""

        directory = Path(input_dir).expanduser()
        if not directory.exists():
            raise FileNotFoundError(f"Input directory not found: {directory}")

        iterator = directory.rglob("*") if recursive else directory.iterdir()
        candidates = [
            path
            for path in iterator
            if path.is_file()
            and self.detect_format(path) is not None
            and self.detect_format(path) != OfficeFormat.TXT
        ]
        candidates.sort(key=lambda path: path.name)
        return self.convert_batch(candidates, target_format, output_dir)

    def to_html(
        self,
        input_path: str | Path,
        output_path: Optional[str | Path] = None,
        *,
        encoding: str = "utf-8",
    ) -> Path:
        """Convert document to HTML."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".html")
        extension = source.suffix.lstrip(".").lower()

        if extension == "docx" and self._has_mammoth:
            html_content = self._convert_docx_with_mammoth(source, to_markdown=False)
            target.write_text(html_content, encoding=encoding)
            return target

        if extension in self.IWORK_EXTENSIONS and self.textutil_path:
            html_content = self._convert_textutil_to_string(source, "html")
            target.write_text(html_content, encoding=encoding)
            return target

        # Fallback to LibreOffice export
        return self._convert_with_libreoffice(source, target, "html")

    def to_markdown(
        self,
        input_path: str | Path,
        output_path: Optional[str | Path] = None,
        *,
        encoding: str = "utf-8",
    ) -> Path:
        """Convert document to Markdown using mammoth or HTML fallbacks."""

        source = self._normalize_input(input_path)
        target = self._prepare_output(source, output_path, suffix=".md")
        extension = source.suffix.lstrip(".").lower()

        if extension == "docx" and self._has_mammoth:
            markdown_content = self._convert_docx_with_mammoth(source, to_markdown=True)
            target.write_text(markdown_content, encoding=encoding)
            return target

        html_source: Optional[Path] = None
        html_text: Optional[str] = None

        if extension in self.IWORK_EXTENSIONS and self.textutil_path:
            html_text = self._convert_textutil_to_string(source, "html")
        else:
            html_source = self._temporary_conversion(source, "html")
            html_text = html_source.read_text(encoding="utf-8", errors="ignore")

        markdown_content = self._html_to_markdown(html_text)
        target.write_text(markdown_content, encoding=encoding)

        if html_source:
            html_source.unlink(missing_ok=True)

        return target

    def batch_convert(
        self,
        input_files: Sequence[str | Path],
        target_format: str,
        output_dir: Optional[str | Path] = None,
    ) -> List[Path]:
        """Convert multiple files to the same target format."""

        fmt = target_format.lower()
        method_map: Dict[str, Tuple[Callable[..., Path], str]] = {
            "pdf": (self.to_pdf, ".pdf"),
            "html": (self.to_html, ".html"),
            "markdown": (self.to_markdown, ".md"),
            "md": (self.to_markdown, ".md"),
            "text": (self.to_text, ".txt"),
            "txt": (self.to_text, ".txt"),
        }

        if fmt not in method_map:
            raise ValueError(f"Unsupported batch target format: {target_format}")

        method, suffix = method_map[fmt]
        outputs: List[Path] = []
        output_dir_path = Path(output_dir).expanduser() if output_dir else None
        if output_dir_path:
            output_dir_path.mkdir(parents=True, exist_ok=True)

        for file_path in input_files:
            try:
                source = self._normalize_input(file_path)
            except FileNotFoundError:
                logger.warning("Skipping missing file: %s", file_path)
                continue

            destination = (
                output_dir_path / f"{source.stem}{suffix}" if output_dir_path else None
            )

            try:
                result_path = method(source, destination)
                outputs.append(result_path)
            except Exception as exc:  # pragma: no cover - defensive logging
                logger.error("Failed to convert %s: %s", source.name, exc)

        logger.info(
            "Batch conversion complete: %d/%d succeeded",
            len(outputs),
            len(input_files),
        )
        return outputs

    # ------------------------------------------------------------------ Core helpers

    def _extract_textual_content(self, source: Path) -> Tuple[str, DocumentContent]:
        """Return text content alongside a lightweight DocumentContent."""

        extension = source.suffix.lstrip(".").lower()

        if extension in self.WORD_LIKE_EXTENSIONS:
            text = self._extract_word_text(source)
        elif extension in self.EXCEL_LIKE_EXTENSIONS:
            text = self._extract_spreadsheet_text(source)
        elif extension in self.POWERPOINT_LIKE_EXTENSIONS:
            text = self._extract_presentation_text(source)
        elif extension in self.IWORK_EXTENSIONS:
            text = self._extract_iwork_text(source)
        else:
            text = self._extract_via_libreoffice_text(source)

        document = DocumentContent(
            format=self.FORMAT_ALIAS.get(extension, ModelOfficeFormat.DOCX),
            paragraphs=[
                Paragraph(runs=[TextRun(text=line)])
                for line in text.splitlines()
                if line.strip()
            ],
        )
        return text, document

    def _extract_word_text(self, source: Path) -> str:
        """Extract text using python-docx with LibreOffice fallback."""

        if not self._has_python_docx:
            logger.debug(
                "python-docx unavailable; using LibreOffice for text extraction"
            )
            return self._extract_via_libreoffice_text(source)

        docx = builtins.__import__("docx")
        doc = docx.Document(str(source))
        lines: List[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)

        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    lines.append(row_text)

        return "\n".join(lines)

    def _extract_spreadsheet_text(self, source: Path) -> str:
        """Extract text from spreadsheet formats using openpyxl."""

        if not self._has_openpyxl:
            logger.debug("openpyxl unavailable; using LibreOffice for text extraction")
            return self._extract_via_libreoffice_text(source)

        openpyxl = builtins.__import__("openpyxl")
        workbook = openpyxl.load_workbook(str(source), data_only=True, read_only=True)
        lines: List[str] = []

        for sheet in workbook.worksheets:
            lines.append(f"[{sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                if all(cell is None or str(cell).strip() == "" for cell in row):
                    continue
                row_text = "\t".join("" if cell is None else str(cell) for cell in row)
                lines.append(row_text)

        workbook.close()
        return "\n".join(lines)

    def _extract_presentation_text(self, source: Path) -> str:
        """Extract text from PPT/PPTX using python-pptx."""

        if not HAS_PPTX:
            logger.debug(
                "python-pptx unavailable; using LibreOffice for text extraction"
            )
            return self._extract_via_libreoffice_text(source)

        presentation = PptxPresentation(str(source))
        lines: List[str] = []

        for index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"Slide {index}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = getattr(shape, "text", "").strip()
                    if text:
                        lines.append(text)

        return "\n".join(lines)

    def _extract_iwork_text(self, source: Path) -> str:
        """Extract text from iWork formats via textutil on macOS."""

        if self.textutil_path:
            return self._convert_textutil_to_string(source, "txt")

        logger.debug(
            "textutil unavailable; using LibreOffice for iWork text extraction"
        )
        return self._extract_via_libreoffice_text(source)

    def _extract_via_libreoffice_text(self, source: Path) -> str:
        """Generic LibreOffice TXT extraction used as a fallback."""

        temp_txt = self._temporary_conversion(source, "txt")
        try:
            return temp_txt.read_text(encoding="utf-8", errors="ignore")
        finally:
            temp_txt.unlink(missing_ok=True)

    def _convert_with_docx2pdf(self, source: Path, target: Path) -> Path:
        """Use docx2pdf for Word conversions on macOS/Windows."""

        if not self._has_docx2pdf:
            raise ConversionError("docx2pdf is not installed.")
        if self.platform not in {"Windows", "Darwin"}:
            raise ConversionError("docx2pdf only supports Windows/macOS.")

        target.parent.mkdir(parents=True, exist_ok=True)
        try:
            docx2pdf = builtins.__import__("docx2pdf")
            docx2pdf.convert(str(source), str(target))
        except Exception as exc:  # pragma: no cover - external dependency
            raise ConversionError(f"docx2pdf failed: {exc}") from exc
        return target

    def _convert_with_libreoffice(
        self, source: Path, target: Path, target_format: str
    ) -> Path:
        """Invoke LibreOffice headless mode for conversion."""

        if not self.libreoffice_path:
            raise ConversionError("LibreOffice binary not found.")

        target.parent.mkdir(parents=True, exist_ok=True)
        cmd = [
            str(self.libreoffice_path),
            "--headless",
            "--convert-to",
            target_format,
            "--outdir",
            str(target.parent),
            str(source),
        ]

        logger.debug("Running LibreOffice command: %s", " ".join(cmd))

        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=self.timeout,
            check=False,
        )

        if result.returncode != 0:
            raise ConversionError(
                result.stderr or result.stdout or "LibreOffice conversion failed."
            )

        produced = target.parent / f"{source.stem}.{target_format.split(':', 1)[0]}"
        if produced != target:
            if produced.exists():
                produced.replace(target)
            else:
                raise ConversionError(f"Expected output not produced: {produced}")

        return target

    def _convert_with_textutil(
        self, source: Path, target: Path, target_format: str
    ) -> Path:
        """Use macOS textutil for Pages/Numbers/Keynote exports."""

        if not self.textutil_path:
            raise ConversionError("textutil not available on this platform.")

        target.parent.mkdir(parents=True, exist_ok=True)
        cmd = [
            str(self.textutil_path),
            "-convert",
            target_format,
            str(source),
            "-output",
            str(target),
        ]

        logger.debug("Running textutil command: %s", " ".join(cmd))

        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=self.timeout,
            check=False,
        )

        if result.returncode != 0:
            raise ConversionError(
                result.stderr or result.stdout or "textutil conversion failed."
            )

        if not target.exists():
            raise ConversionError(f"textutil reported success but {target} is missing.")

        return target

    def _convert_textutil_to_string(self, source: Path, target_format: str) -> str:
        """Convert using textutil and return the produced text without persisting."""

        with tempfile.TemporaryDirectory() as tmp_dir:
            temp_path = Path(tmp_dir) / f"{source.stem}.{target_format}"
            self._convert_with_textutil(source, temp_path, target_format)
            return temp_path.read_text(encoding="utf-8", errors="ignore")

    def _convert_docx_with_mammoth(self, source: Path, *, to_markdown: bool) -> str:
        """Use mammoth for DOCX -> HTML/Markdown conversion."""

        if not self._has_mammoth:
            raise ConversionError("mammoth is not installed.")

        mammoth = builtins.__import__("mammoth")

        with source.open("rb") as handle:
            if to_markdown:
                result = mammoth.convert_to_markdown(handle)
            else:
                result = mammoth.convert_to_html(handle)

        for message in result.messages:  # pragma: no cover - best effort logging
            logger.debug("mammoth notice: %s", message)

        return result.value

    def _find_libreoffice(self) -> Optional[str]:
        """Return best-effort LibreOffice binary location."""

        candidates: List[str] = []
        if self.platform == "Darwin":
            candidates = [
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "/usr/local/bin/soffice",
                "/opt/homebrew/bin/soffice",
            ]
        elif self.platform == "Linux":
            candidates = [
                "/usr/bin/soffice",
                "/usr/bin/libreoffice",
                "/usr/local/bin/soffice",
            ]
        elif self.platform == "Windows":
            candidates = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            ]

        for candidate in candidates:
            if os.path.exists(candidate):
                return candidate

        return shutil.which("soffice")

    def _check_dependencies(self) -> None:
        """Detect optional conversion dependencies."""

        def _import_available(name: str) -> bool:
            try:
                builtins.__import__(name)
                return True
            except ImportError:
                return False

        self._has_docx2pdf = _import_available("docx2pdf")
        self._has_mammoth = _import_available("mammoth")
        try:
            docx = builtins.__import__("docx")
            self._has_python_docx = hasattr(docx, "Document")
        except ImportError:
            self._has_python_docx = False
        self._has_openpyxl = _import_available("openpyxl")

    def _run_libreoffice_conversion(
        self, input_path: Path, output_format: str, output_dir: Path
    ) -> Path:
        """Run LibreOffice headless conversion and return the output path."""

        if not self.libreoffice_path:
            raise ConversionError("LibreOffice not found")
        source = self._normalize_input(input_path)
        output_dir = Path(output_dir).expanduser()
        output_dir.mkdir(parents=True, exist_ok=True)
        output_format = output_format.lower().lstrip(".")

        cmd = [
            self.libreoffice_path,
            "--headless",
            "--convert-to",
            output_format,
            "--outdir",
            str(output_dir),
            str(source.resolve()),
        ]
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=self.timeout,
                check=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise ConversionError("LibreOffice conversion timed out") from exc

        if result.returncode != 0:
            message = result.stderr or result.stdout or "LibreOffice conversion failed."
            raise ConversionError(f"Conversion failed: LibreOffice conversion failed: {message}")

        output_path = output_dir / f"{source.stem}.{output_format}"
        if not output_path.exists():
            raise ConversionError("Conversion failed: LibreOffice output file not found")
        return output_path

    def _run_macos_automation(
        self, app_name: str, input_path: Path, output_path: Path
    ) -> Path:
        """Run macOS automation for iWork PDF exports."""

        if self.platform != "Darwin":
            raise ConversionError(f"macOS automation not available on {self.platform}")

        source = self._normalize_input(input_path)
        output = Path(output_path).expanduser()
        output.parent.mkdir(parents=True, exist_ok=True)

        script = (
            f'tell application "{app_name}"\n'
            f'  set inputFile to POSIX file "{source}"\n'
            f'  set outputFile to POSIX file "{output}"\n'
            f"  open inputFile\n"
            f"  export to outputFile as PDF\n"
            f"  close front document\n"
            f"end tell"
        )
        try:
            result = subprocess.run(
                ["osascript", "-e", script],
                capture_output=True,
                text=True,
                timeout=30,
                check=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise ConversionError(f"{app_name} automation timed out") from exc

        if result.returncode != 0:
            message = result.stderr or result.stdout or "automation failure"
            raise ConversionError(
                f"macOS automation failed: {app_name} automation failed: {message}"
            )

        if not output.exists():
            raise ConversionError(f"{app_name} automation did not produce output")

        return output

    @staticmethod
    def _rename_if_needed(produced: Path, target: Path) -> Path:
        if produced.resolve() == target.resolve():
            return target
        target.parent.mkdir(parents=True, exist_ok=True)
        if produced.exists():
            produced.replace(target)
            return target
        raise ConversionError(f"Expected output not produced: {produced}")

    def _temporary_conversion(self, source: Path, target_format: str) -> Path:
        """Perform a temporary LibreOffice conversion and return the temp file path."""

        with tempfile.TemporaryDirectory() as tmp_dir:
            temp_target = Path(tmp_dir) / f"{source.stem}.{target_format}"
            self._convert_with_libreoffice(source, temp_target, target_format)
            # Copy to a persistent temp path because TemporaryDirectory cleans up immediately
            fd, temp_name = tempfile.mkstemp(suffix=f".{target_format}")
            os.close(fd)
            persistent_temp = Path(temp_name)
            shutil.copy(temp_target, persistent_temp)
        return persistent_temp

    # ------------------------------------------------------------------ Utilities

    def _normalize_input(self, input_path: str | Path) -> Path:
        """Validate and normalise the incoming path."""

        path = Path(input_path).expanduser()
        if not path.exists():
            raise FileNotFoundError(f"Input file not found: {path}")
        return path

    def _prepare_output(
        self,
        source: Path,
        output_path: Optional[str | Path],
        *,
        suffix: str,
    ) -> Path:
        """Return a writable output path, creating parent directories as needed."""

        target = (
            Path(output_path).expanduser()
            if output_path
            else source.with_suffix(suffix)
        )
        target.parent.mkdir(parents=True, exist_ok=True)
        return target

    def _auto_detect_libreoffice(self) -> Optional[Path]:
        """Best-effort detection of the LibreOffice binary."""

        candidates: List[Path] = []
        if self.platform == "Darwin":
            candidates = [
                Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"),
                Path("/usr/local/bin/soffice"),
                Path("/opt/homebrew/bin/soffice"),
            ]
        elif self.platform == "Linux":
            candidates = [
                Path("/usr/bin/soffice"),
                Path("/usr/bin/libreoffice"),
                Path("/usr/local/bin/soffice"),
            ]
        elif self.platform == "Windows":
            candidates = [
                Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
                Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
            ]

        for candidate in candidates:
            if candidate.exists():
                return candidate

        env_path = shutil.which("soffice")
        return Path(env_path).resolve() if env_path else None

    def _can_use_docx2pdf(self) -> bool:
        """Return True when docx2pdf is usable on the current platform."""

        return self._has_docx2pdf and self.platform in {"Windows", "Darwin"}

    def _html_to_markdown(self, html: str) -> str:
        """Convert HTML to Markdown using html2text if available."""

        if HAS_HTML2TEXT:
            parser = html2text.HTML2Text()
            parser.ignore_links = False
            parser.ignore_images = False
            parser.body_width = 0
            return parser.handle(html)

        # Lightweight fallback: strip tags and normalise whitespace
        text = re.sub(
            r"<(script|style)[^>]*>.*?</\1>", "", html, flags=re.IGNORECASE | re.DOTALL
        )
        text = re.sub(r"<[^>]+>", "", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def _maybe_emit_debug_preview(self, text: str, limit: int = 400) -> None:
        """Log a short preview of extracted text for debugging."""

        preview = text.strip().splitlines()[:10]
        if preview:
            logger.debug("Text preview:\n%s", "\n".join(preview)[:limit])


__all__ = ["OfficeConverter", "OfficeFormat", "ConversionError"]
